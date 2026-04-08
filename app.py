import io
import itertools
import math
import textwrap
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import altair as alt
import numpy as np
import pandas as pd
import streamlit as st


st.set_page_config(page_title="Supplier and Sourcing Analysis Tool", layout="wide")
st.markdown(
    """
    <style>
    div[role="tablist"] > button[role="tab"],
    [data-baseweb="tab-list"] button[role="tab"] {
        color: #ffffff !important;
        font-size: 1.05rem !important;
        font-weight: 800 !important;
        padding: 0.9rem 1.35rem !important;
        min-height: 3.15rem !important;
        border-radius: 0.8rem 0.8rem 0 0 !important;
        border-width: 2px !important;
        border-style: solid !important;
        opacity: 0.92 !important;
        transition: transform 0.12s ease, opacity 0.12s ease, filter 0.12s ease !important;
    }
    div[role="tablist"] > button[role="tab"] p,
    div[role="tablist"] > button[role="tab"] div,
    div[role="tablist"] > button[role="tab"] span,
    [data-baseweb="tab-list"] button[role="tab"] p,
    [data-baseweb="tab-list"] button[role="tab"] div,
    [data-baseweb="tab-list"] button[role="tab"] span {
        color: #ffffff !important;
        font-size: 1.05rem !important;
        font-weight: 800 !important;
    }
    div[role="tablist"] > button[role="tab"]:hover,
    [data-baseweb="tab-list"] button[role="tab"]:hover {
        transform: translateY(-1px);
        opacity: 1 !important;
        filter: brightness(1.03);
    }
    div[role="tablist"] > button[role="tab"][aria-selected="true"],
    [data-baseweb="tab-list"] button[role="tab"][aria-selected="true"] {
        opacity: 1 !important;
        filter: saturate(1.08) brightness(0.95);
        box-shadow: inset 0 -4px 0 rgba(255, 255, 255, 0.24) !important;
    }
    div[role="tablist"] > button[role="tab"]:nth-of-type(1),
    [data-baseweb="tab-list"] button[role="tab"]:nth-of-type(1) {
        background: #1f4e79 !important;
        border-color: #1f4e79 !important;
    }
    div[role="tablist"] > button[role="tab"]:nth-of-type(2),
    [data-baseweb="tab-list"] button[role="tab"]:nth-of-type(2) {
        background: #0f766e !important;
        border-color: #0f766e !important;
    }
    div[role="tablist"] > button[role="tab"]:nth-of-type(3),
    [data-baseweb="tab-list"] button[role="tab"]:nth-of-type(3) {
        background: #b45309 !important;
        border-color: #b45309 !important;
    }
    div[role="tablist"] > button[role="tab"]:nth-of-type(4),
    [data-baseweb="tab-list"] button[role="tab"]:nth-of-type(4) {
        background: #7c3aed !important;
        border-color: #7c3aed !important;
    }
    div[role="tablist"] > button[role="tab"]:nth-of-type(5),
    [data-baseweb="tab-list"] button[role="tab"]:nth-of-type(5) {
        background: #b91c1c !important;
        border-color: #b91c1c !important;
    }
    div[role="tablist"] > button[role="tab"]:nth-of-type(6),
    [data-baseweb="tab-list"] button[role="tab"]:nth-of-type(6) {
        background: #475569 !important;
        border-color: #475569 !important;
    }
    div[role="tablist"] > button[role="tab"]:nth-of-type(7),
    [data-baseweb="tab-list"] button[role="tab"]:nth-of-type(7) {
        background: #065f46 !important;
        border-color: #065f46 !important;
    }
    div[role="tablist"] > button[role="tab"]:nth-of-type(8),
    [data-baseweb="tab-list"] button[role="tab"]:nth-of-type(8) {
        background: #1d4ed8 !important;
        border-color: #1d4ed8 !important;
    }
    div[role="tablist"] > button[role="tab"]:nth-of-type(9),
    [data-baseweb="tab-list"] button[role="tab"]:nth-of-type(9) {
        background: #92400e !important;
        border-color: #92400e !important;
    }
    div.stButton > button {
        border-radius: 0.8rem !important;
        font-weight: 700 !important;
        min-height: 2.9rem !important;
        border-width: 2px !important;
        box-shadow: 0 6px 14px rgba(31, 78, 121, 0.12) !important;
        transition: transform 0.15s ease, box-shadow 0.15s ease, background-color 0.15s ease !important;
    }
    div.stButton > button:hover {
        transform: translateY(-1px);
        box-shadow: 0 10px 18px rgba(31, 78, 121, 0.18) !important;
    }
    div.stButton > button[kind="primary"] {
        background: linear-gradient(135deg, #0f766e, #115e59) !important;
        border-color: #115e59 !important;
        color: #ffffff !important;
    }
    div.stButton > button[kind="primary"]:hover {
        background: linear-gradient(135deg, #0d6760, #0f4f4b) !important;
        border-color: #0f4f4b !important;
        color: #ffffff !important;
    }
    div.stButton > button[kind="secondary"] {
        background: #fff8e8 !important;
        border-color: #d97706 !important;
        color: #8a4b00 !important;
    }
    div.stButton > button[kind="secondary"]:hover {
        background: #ffefc7 !important;
        border-color: #b45309 !important;
        color: #7c3f00 !important;
    }
    div.stButton > button[kind="secondary"] p,
    div.stButton > button[kind="primary"] p {
        font-weight: 700 !important;
    }
    </style>
    """,
    unsafe_allow_html=True,
)

DEFAULT_DATASET_CANDIDATES = [
    Path(__file__).with_name("sample_data.xlsx"),
    Path(__file__).with_name("sample_data.xls"),
]
RISK_LEVEL_DOMAIN = ["High", "Medium", "Low"]
RISK_LEVEL_RANGE = ["#d73027", "#f39c12", "#2e8b57"]
ABC_CLASS_DOMAIN = ["A", "B", "C"]
ABC_CLASS_RANGE = ["#1f4e79", "#7a5c99", "#8a6f42"]


COLUMN_ALIASES = {
    "supplier": ["supplier", "vendor", "supplier_name", "vendor_name", "source"],
    "component": ["component", "part", "item", "sku", "material", "part_name", "component_name", "product"],
    "spend": ["spend", "amount", "total_spend", "purchase_amount", "value", "cost_total", "extended_cost", "total_price", "total_cost"],
    "units": ["units", "qty", "quantity", "volume", "ordered_qty", "received_qty"],
    "unit_cost": ["unit_cost", "price", "cost_per_unit", "unit_price", "rate"],
    "defects": ["defects", "defect_qty", "defect_count", "returns", "rejections", "nonconforming_qty"],
    "defect_rate": ["defect_rate", "ppm", "failure_rate", "quality_rate"],
    "lead_time": ["lead_time", "lead_time_days", "cycle_time", "delivery_days", "days_to_deliver"],
    "order_date": ["order_date", "po_date", "purchase_date", "request_date", "created_date"],
    "receipt_date": ["receipt_date", "delivery_date", "received_date", "arrival_date", "ship_date", "fulfilled_date"],
    "risk_score": ["risk_score", "supplier_risk", "risk", "external_risk", "risk_rating"],
    "criticality": ["criticality", "criticality_score", "business_criticality", "importance", "priority"],
}

CANONICAL_COLUMNS = [
    "supplier",
    "component",
    "spend",
    "units",
    "unit_cost",
    "defects",
    "defect_rate",
    "lead_time",
    "order_date",
    "receipt_date",
    "risk_score",
    "criticality",
]


def cleanup_excel_frame(df: pd.DataFrame) -> pd.DataFrame:
    cleaned = df.loc[:, ~df.columns.astype(str).str.startswith("Unnamed:")].copy()
    cleaned = cleaned.dropna(how="all").dropna(axis=1, how="all")
    return cleaned


def safe_divide(numerator, denominator, fill_value=0.0):
    denominator_series = pd.Series(denominator)
    numerator_series = pd.Series(numerator)
    result = np.where(denominator_series.fillna(0).astype(float) == 0, fill_value, numerator_series.astype(float) / denominator_series.astype(float))
    return pd.Series(result, index=numerator_series.index)


def min_max_scale(series: pd.Series, invert: bool = False) -> pd.Series:
    numeric = pd.to_numeric(series, errors="coerce").fillna(0.0).astype(float)
    minimum = float(numeric.min()) if len(numeric) else 0.0
    maximum = float(numeric.max()) if len(numeric) else 0.0
    if np.isclose(maximum, minimum):
        scaled = pd.Series(np.full(len(numeric), 50.0), index=numeric.index)
    else:
        scaled = (numeric - minimum) / (maximum - minimum) * 100.0
    if invert:
        scaled = 100.0 - scaled
    return scaled.clip(0, 100)


def as_percentage(series: pd.Series) -> pd.Series:
    numeric = pd.to_numeric(series, errors="coerce").fillna(0.0).astype(float)
    if numeric.max() > 1.0:
        return numeric / 100.0
    return numeric


def identify_columns(df: pd.DataFrame) -> Dict[str, Optional[str]]:
    normalized_map = {str(col).strip().lower().replace(" ", "_"): col for col in df.columns}
    matches: Dict[str, Optional[str]] = {}
    for canonical, aliases in COLUMN_ALIASES.items():
        matches[canonical] = None
        for alias in aliases:
            alias_key = alias.strip().lower().replace(" ", "_")
            if alias_key in normalized_map:
                matches[canonical] = normalized_map[alias_key]
                break
    return matches


def normalize_input_data(df: pd.DataFrame) -> pd.DataFrame:
    working = df.copy()
    column_map = identify_columns(working)
    normalized = pd.DataFrame(index=working.index)

    for canonical in CANONICAL_COLUMNS:
        source_column = column_map.get(canonical)
        if source_column is not None:
            normalized[canonical] = working[source_column]
        else:
            normalized[canonical] = np.nan

    normalized["supplier"] = normalized["supplier"].fillna("Unknown Supplier").astype(str).str.strip()
    normalized["component"] = normalized["component"].fillna("Unknown Component").astype(str).str.strip()

    numeric_cols = ["spend", "units", "unit_cost", "defects", "defect_rate", "lead_time", "risk_score", "criticality"]
    for col in numeric_cols:
        normalized[col] = pd.to_numeric(normalized[col], errors="coerce")

    normalized["order_date"] = pd.to_datetime(normalized["order_date"], errors="coerce")
    normalized["receipt_date"] = pd.to_datetime(normalized["receipt_date"], errors="coerce")

    normalized["units"] = normalized["units"].fillna(0.0)
    normalized["defects"] = normalized["defects"].fillna(0.0)
    normalized["risk_score"] = normalized["risk_score"].fillna(normalized["risk_score"].median())
    normalized["criticality"] = normalized["criticality"].fillna(normalized["criticality"].median())
    normalized["risk_score"] = normalized["risk_score"].fillna(50.0)
    normalized["criticality"] = normalized["criticality"].fillna(50.0)

    calc_unit_cost = safe_divide(normalized["spend"], normalized["units"], fill_value=np.nan)
    normalized["unit_cost"] = normalized["unit_cost"].fillna(calc_unit_cost)

    calc_spend = normalized["units"].fillna(0.0) * normalized["unit_cost"].fillna(0.0)
    normalized["spend"] = normalized["spend"].fillna(calc_spend)

    calculated_lead_time = (normalized["receipt_date"] - normalized["order_date"]).dt.days
    normalized["lead_time"] = normalized["lead_time"].fillna(calculated_lead_time)
    normalized["lead_time"] = normalized["lead_time"].fillna(normalized["lead_time"].median())
    normalized["lead_time"] = normalized["lead_time"].fillna(0.0).clip(lower=0)

    defect_rate_calc = safe_divide(normalized["defects"], normalized["units"], fill_value=np.nan)
    normalized["defect_rate"] = normalized["defect_rate"].fillna(defect_rate_calc)
    normalized["defect_rate"] = as_percentage(normalized["defect_rate"]).fillna(0.0).clip(lower=0.0)

    normalized["spend"] = normalized["spend"].fillna(0.0).clip(lower=0.0)
    normalized["unit_cost"] = normalized["unit_cost"].fillna(0.0).clip(lower=0.0)
    normalized["defects"] = normalized["defects"].fillna(0.0).clip(lower=0.0)
    normalized["record_id"] = np.arange(1, len(normalized) + 1)
    return normalized


def build_input_diagnostics(df: pd.DataFrame) -> List[str]:
    cleaned = cleanup_excel_frame(df)
    column_map = identify_columns(cleaned)
    diagnostics: List[str] = []

    for field_name, fallback_value in [("risk_score", 50.0), ("criticality", 50.0)]:
        source_column = column_map.get(field_name)
        if source_column is None:
            diagnostics.append(
                f"`{field_name}` was not found in the data source, so the normalized input defaults it to {fallback_value:.0f}."
            )
            continue

        numeric_values = pd.to_numeric(cleaned[source_column], errors="coerce")
        if int(numeric_values.notna().sum()) == 0:
            diagnostics.append(
                f"`{field_name}` was matched to source column `{source_column}`, but no numeric values were usable, so the normalized input defaults it to {fallback_value:.0f}."
            )

    return diagnostics


def build_input_field_status(df: pd.DataFrame) -> pd.DataFrame:
    cleaned = cleanup_excel_frame(df)
    column_map = identify_columns(cleaned)
    field_rules = {
        "supplier": "Required identifier. Missing values normalize to `Unknown Supplier`.",
        "component": "Required identifier. Missing values normalize to `Unknown Component`.",
        "spend": "Used directly when present; otherwise derived from units x unit cost where possible.",
        "units": "Missing values default to 0 for normalization.",
        "defects": "Missing values default to 0 and feed the defect-rate calculation.",
        "lead_time": "Used directly when present; otherwise derived from receipt date minus order date where possible.",
        "risk_score": "Used as an input to the modeled risk score; defaults to 50 when not provided or unusable.",
        "criticality": "Used as an input to strategic priority; defaults to 50 when not provided or unusable.",
    }
    numeric_fields = {"spend", "units", "defects", "lead_time", "risk_score", "criticality"}
    rows: List[Dict[str, str]] = []
    for field_name, handling in field_rules.items():
        source_column = column_map.get(field_name)
        if source_column is None:
            status = "Missing"
            source_text = "Not found"
        elif field_name in numeric_fields:
            numeric_values = pd.to_numeric(cleaned[source_column], errors="coerce")
            if int(numeric_values.notna().sum()) == 0:
                status = "Found, unusable"
            else:
                status = "Provided"
            source_text = str(source_column)
        else:
            status = "Provided"
            source_text = str(source_column)
        rows.append(
            {
                "Field": field_name,
                "Status": status,
                "Source Column": source_text,
                "Handling": handling,
            }
        )
    return pd.DataFrame(rows)


def build_data_quality_summary(field_status: pd.DataFrame) -> str:
    if field_status.empty:
        return "Data quality checks summarize which core fields were provided versus inferred during normalization."
    provided_count = int(field_status["Status"].eq("Provided").sum())
    missing_count = int(field_status["Status"].eq("Missing").sum())
    unusable_count = int(field_status["Status"].eq("Found, unusable").sum())
    inferred_fields = field_status.loc[field_status["Status"].ne("Provided"), "Field"].tolist()
    inferred_text = ", ".join(inferred_fields) if inferred_fields else "none"
    return (
        f"Data quality snapshot: {provided_count} core fields were provided as usable source data, "
        f"{missing_count} were missing, and {unusable_count} were found but unusable. "
        f"Fields currently relying on inference or fallback logic are {inferred_text}."
    )


def render_glossary_drawer() -> None:
    with st.expander("Glossary & Help"):
        st.write("`Single-source`: a component with only one effective supplier in the current view.")
        st.write("`Mitigation supplier`: a planned backup supplier assigned to reduce single-source or uncovered exposure.")
        st.write("`Effective supplier count`: selected suppliers plus any explicit mitigation suppliers assigned to that component.")
        st.write("`Risk-adjusted spend`: spend weighted by quality burden, sourcing concentration, and modeled supply risk.")
        st.write("`Kraljic quadrant`: component positioning across business impact and supply risk to guide sourcing strategy.")
        st.write("`Covered spend`: spend attached to components that still have at least one effective supplier in the evaluated scenario.")
        st.write("`Net savings`: modeled gross consolidation savings minus modeled mitigation qualification cost.")


def build_sample_data() -> pd.DataFrame:
    sample_rows = [
        ["Alpha Metals", "Bearing Assembly", 160000, 2000, 80, 18, 21, 62, 88],
        ["Beta Industrial", "Bearing Assembly", 90000, 1100, 82, 10, 18, 45, 88],
        ["Alpha Metals", "Housing Unit", 120000, 1500, 80, 9, 26, 58, 78],
        ["Delta Components", "Housing Unit", 60000, 800, 75, 22, 34, 71, 78],
        ["Nova Circuits", "Controller Board", 230000, 900, 255, 32, 42, 84, 95],
        ["Prime Source", "Controller Board", 40000, 150, 267, 7, 37, 64, 95],
        ["Omega Plastics", "Seal Kit", 50000, 5000, 10, 14, 10, 32, 52],
        ["Vertex Supply", "Seal Kit", 35000, 3300, 11, 8, 12, 28, 52],
        ["Vertex Supply", "Fastener Pack", 42000, 12000, 3.5, 30, 8, 25, 40],
        ["Sigma Parts", "Fastener Pack", 28000, 9000, 3.1, 18, 11, 38, 40],
        ["Nova Circuits", "Sensor Module", 145000, 500, 290, 21, 45, 87, 91],
        ["Core Micro", "Sensor Module", 60000, 220, 273, 12, 39, 74, 91],
        ["Gamma Forging", "Valve Body", 98000, 950, 103, 11, 29, 55, 73],
        ["Titan Works", "Valve Body", 66000, 620, 106, 15, 31, 68, 73],
        ["Solo Precision", "Custom Bracket", 76000, 400, 190, 16, 47, 82, 86],
        ["Alpha Metals", "Rotor Shaft", 110000, 700, 157, 6, 24, 51, 84],
        ["Titan Works", "Rotor Shaft", 45000, 300, 150, 9, 27, 66, 84],
        ["Prime Source", "Power Relay", 72000, 1200, 60, 11, 20, 53, 65],
        ["Electra Hub", "Power Relay", 68000, 1050, 64.8, 5, 16, 41, 65],
        ["Solo Precision", "Safety Latch", 54000, 240, 225, 9, 51, 79, 93],
    ]
    sample_df = pd.DataFrame(
        sample_rows,
        columns=["supplier", "component", "spend", "units", "unit_cost", "defects", "lead_time", "risk_score", "criticality"],
    )
    return normalize_input_data(sample_df)


@st.cache_data(show_spinner=False)
def get_default_data() -> Tuple[pd.DataFrame, str, List[str], pd.DataFrame]:
    for dataset_path in DEFAULT_DATASET_CANDIDATES:
        if dataset_path.exists():
            return load_uploaded_data(dataset_path.name, dataset_path.read_bytes())
    sample_df = build_sample_data()
    return sample_df, "Built-in sample data", [], build_input_field_status(sample_df)


@st.cache_data(show_spinner=False)
def load_uploaded_data(file_name: str, file_bytes: bytes) -> Tuple[pd.DataFrame, str, List[str], pd.DataFrame]:
    buffer = io.BytesIO(file_bytes)
    lowered_name = file_name.lower()
    if lowered_name.endswith(".csv"):
        raw = pd.read_csv(buffer)
        source_label = f"Loaded file: {file_name} (CSV)"
    elif lowered_name.endswith(".xlsx") or lowered_name.endswith(".xls"):
        workbook = pd.ExcelFile(buffer)
        best_sheet = None
        best_score = -1
        best_rows = -1

        for sheet_name in workbook.sheet_names:
            candidate = cleanup_excel_frame(pd.read_excel(workbook, sheet_name=sheet_name))
            if candidate.empty:
                continue

            column_hits = sum(1 for matched in identify_columns(candidate).values() if matched is not None)
            row_count = len(candidate)
            if column_hits > best_score or (column_hits == best_score and row_count > best_rows):
                best_sheet = sheet_name
                best_score = column_hits
                best_rows = row_count

        if best_sheet is None or best_score <= 0:
            raise ValueError("No usable worksheet was found in the uploaded Excel file.")

        raw = pd.read_excel(workbook, sheet_name=best_sheet)
        source_label = f"Loaded file: {file_name} | Worksheet: {best_sheet}"
    else:
        raise ValueError("Please upload a CSV or Excel file.")
    cleaned = cleanup_excel_frame(raw)
    return normalize_input_data(cleaned), source_label, build_input_diagnostics(cleaned), build_input_field_status(cleaned)


def assign_abc_categories(sorted_df: pd.DataFrame, value_col: str, prefix: str) -> pd.DataFrame:
    result = sorted_df.copy()
    total = float(result[value_col].sum())
    if total <= 0:
        result[f"{prefix}_cum_share"] = 0.0
        result[f"{prefix}_abc"] = "C"
        return result
    result[f"{prefix}_cum_share"] = result[value_col].cumsum() / total
    result[f"{prefix}_abc"] = np.select(
        [result[f"{prefix}_cum_share"] <= 0.8, result[f"{prefix}_cum_share"] <= 0.95],
        ["A", "B"],
        default="C",
    )
    return result


def classify_sourcing_risk_level(supplier_count: pd.Series, risk_score: pd.Series) -> pd.Series:
    levels = np.select(
        [supplier_count.fillna(0).astype(float) <= 1, supplier_count.fillna(0).astype(float) <= 3, risk_score.fillna(0).astype(float) >= 60],
        ["High", "Medium", "Medium"],
        default="Low",
    )
    return pd.Series(levels, index=risk_score.index)


def classify_kraljic_quadrant(profit_impact_score: pd.Series, supply_risk_score: pd.Series) -> pd.Series:
    profit_threshold = float(profit_impact_score.median()) if len(profit_impact_score) else 0.0
    risk_threshold = float(supply_risk_score.median()) if len(supply_risk_score) else 0.0

    quadrants = np.select(
        [
            (profit_impact_score >= profit_threshold) & (supply_risk_score >= risk_threshold),
            (profit_impact_score >= profit_threshold) & (supply_risk_score < risk_threshold),
            (profit_impact_score < profit_threshold) & (supply_risk_score >= risk_threshold),
        ],
        ["Strategic", "Leverage", "Bottleneck"],
        default="Non-Critical",
    )
    return pd.Series(quadrants, index=profit_impact_score.index)


def scale_if_variable(series: pd.Series, invert: bool = False, default_value: float = 50.0) -> pd.Series:
    numeric = pd.to_numeric(series, errors="coerce").fillna(default_value).astype(float)
    if len(numeric) == 0 or np.isclose(float(numeric.max()), float(numeric.min())):
        return pd.Series(np.full(len(numeric), default_value), index=numeric.index)
    return min_max_scale(numeric, invert=invert)


def classify_suppliers(
    component_supplier_detail: pd.DataFrame,
    component_summary: pd.DataFrame,
    supplier_summary: pd.DataFrame,
) -> pd.DataFrame:
    ranking = component_supplier_detail.copy()
    ranking["rank_metric"] = (
        0.35 * min_max_scale(ranking["defect_rate"], invert=True)
        + 0.20 * min_max_scale(ranking["avg_lead_time"], invert=True)
        + 0.15 * min_max_scale(ranking["avg_risk_score"], invert=True)
        + 0.20 * min_max_scale(ranking["supplier_share"])
        + 0.10 * min_max_scale(ranking["spend"])
    )
    ranking["supplier_rank_within_component"] = ranking.groupby("component")["rank_metric"].rank(method="dense", ascending=False)
    primary_component_winners = ranking.loc[ranking["supplier_rank_within_component"].eq(1), ["component", "supplier"]].copy()
    primary_component_winners["is_preferred_supplier"] = True

    supplier_primary = primary_component_winners.groupby("supplier", as_index=False).agg(preferred_component_count=("component", "nunique"))
    supplier_decision_frame = supplier_summary.merge(supplier_primary, on="supplier", how="left")
    supplier_decision_frame["preferred_component_count"] = supplier_decision_frame["preferred_component_count"].fillna(0)

    preferred_suppliers = set(primary_component_winners["supplier"].tolist())
    protected_suppliers = set(
        supplier_decision_frame.loc[
            supplier_decision_frame["supports_single_source"] | supplier_decision_frame["supports_high_risk"],
            "supplier",
        ].tolist()
    )

    high_spend_cutoff = supplier_decision_frame["spend"].quantile(0.65) if len(supplier_decision_frame) else 0
    strong_perf_cutoff = supplier_decision_frame["performance_score"].quantile(0.55) if len(supplier_decision_frame) else 0
    weak_perf_cutoff = supplier_decision_frame["performance_score"].quantile(0.30) if len(supplier_decision_frame) else 0
    replaceable_cutoff = supplier_decision_frame["replaceability_score"].quantile(0.60) if len(supplier_decision_frame) else 0

    decisions: List[Dict[str, object]] = []
    for row in supplier_decision_frame.to_dict(orient="records"):
        supplier_name = row["supplier"]
        is_protected = supplier_name in protected_suppliers
        is_preferred = supplier_name in preferred_suppliers
        high_spend = float(row["spend"]) >= float(high_spend_cutoff)
        strong_perf = float(row["performance_score"]) >= float(strong_perf_cutoff)
        weak_perf = float(row["performance_score"]) <= float(weak_perf_cutoff)
        highly_replaceable = float(row["replaceability_score"]) >= float(replaceable_cutoff)

        if is_protected or is_preferred or (high_spend and strong_perf):
            decision = "Keep / Consolidate To"
            reason = "Protected supply coverage or preferred performance across overlapping components."
        elif (not is_protected) and highly_replaceable and weak_perf and int(row["preferred_component_count"]) == 0:
            decision = "Eliminate / De-prioritize"
            reason = "Replaceable supplier with weaker quality, lead time, and limited strategic dependence."
        else:
            decision = "Keep and Monitor"
            reason = "Valuable enough to retain, but performance and risk should improve before expansion."

        savings_rate = 0.0
        if decision == "Eliminate / De-prioritize":
            savings_rate = 0.06
        elif decision == "Keep and Monitor":
            savings_rate = 0.03
        elif decision == "Keep / Consolidate To" and not is_protected:
            savings_rate = 0.02
        estimated_savings = round(float(row["spend"]) * savings_rate, 2)

        issues: List[str] = []
        defect_median = float(supplier_decision_frame["defect_rate"].median()) if len(supplier_decision_frame) else 0.0
        lead_time_median = float(supplier_decision_frame["avg_lead_time"].median()) if len(supplier_decision_frame) else 0.0
        risk_median = float(supplier_decision_frame["avg_risk_score"].median()) if len(supplier_decision_frame) else 0.0
        if float(row["defect_rate"]) > defect_median:
            issues.append("quality drift")
        if float(row["avg_lead_time"]) > lead_time_median:
            issues.append("long lead time")
        if float(row["avg_risk_score"]) > risk_median:
            issues.append("elevated external risk")
        if bool(row["supports_single_source"]):
            issues.append("single-source dependency")
        if not issues:
            issues.append("maintain competitiveness")

        if decision == "Keep / Consolidate To":
            action_plan = "Increase award share where overlap exists, lock in continuity plans, and pursue structured cost negotiations."
        elif decision == "Eliminate / De-prioritize":
            action_plan = "Reallocate demand toward stronger alternatives, restrict new awards, and exit after qualification and inventory safeguards."
        else:
            action_plan = "Keep active with a corrective scorecard focused on quality, lead time, and risk-mitigation checkpoints."

        decision_rank = {"Keep / Consolidate To": 1, "Keep and Monitor": 2, "Eliminate / De-prioritize": 3}[decision]
        decisions.append(
            {
                "supplier": supplier_name,
                "decision": decision,
                "decision_reason": reason,
                "issues": ", ".join(issues),
                "supplier_action_plan": action_plan,
                "estimated_savings": estimated_savings,
                "decision_rank": decision_rank,
            }
        )

    return pd.DataFrame(decisions)


@st.cache_data(show_spinner=False)
def build_analytics(df: pd.DataFrame) -> Dict[str, pd.DataFrame]:
    component_supplier_detail = (
        df.groupby(["component", "supplier"], as_index=False)
        .agg(
            spend=("spend", "sum"),
            units=("units", "sum"),
            defects=("defects", "sum"),
            avg_lead_time=("lead_time", "mean"),
            avg_risk_score=("risk_score", "mean"),
            avg_criticality=("criticality", "mean"),
        )
    )
    component_supplier_detail["defect_rate"] = safe_divide(component_supplier_detail["defects"], component_supplier_detail["units"], fill_value=0.0)
    component_supplier_detail["component_total_spend"] = component_supplier_detail.groupby("component")["spend"].transform("sum")
    component_supplier_detail["supplier_share"] = safe_divide(component_supplier_detail["spend"], component_supplier_detail["component_total_spend"], fill_value=0.0)
    component_supplier_detail["effective_supplier_share"] = component_supplier_detail["supplier_share"]
    component_supplier_detail["supplier_count"] = component_supplier_detail.groupby("component")["supplier"].transform("nunique")
    component_supplier_detail["single_source_flag"] = component_supplier_detail["supplier_count"].eq(1)
    component_supplier_detail["backup_supplier_flag"] = component_supplier_detail["supplier_count"].gt(1)
    component_supplier_detail["concentration_gap"] = component_supplier_detail.groupby("component")["supplier_share"].transform("max") - component_supplier_detail["supplier_share"]

    component_summary = (
        component_supplier_detail.groupby("component", as_index=False)
        .agg(
            spend=("spend", "sum"),
            units=("units", "sum"),
            defects=("defects", "sum"),
            supplier_count=("supplier", "nunique"),
            largest_supplier_share=("supplier_share", "max"),
            worst_supplier_defect_rate=("defect_rate", "max"),
            avg_lead_time=("avg_lead_time", "mean"),
            risk_score=("avg_risk_score", "mean"),
            criticality=("avg_criticality", "mean"),
        )
    )
    component_summary["defect_rate"] = safe_divide(component_summary["defects"], component_summary["units"], fill_value=0.0)
    component_summary["single_source_flag"] = component_summary["supplier_count"].eq(1)
    component_summary["backup_supplier_flag"] = component_summary["supplier_count"].gt(1)
    component_summary["supplier_concentration"] = component_summary["largest_supplier_share"].clip(0, 1)
    dominant_supplier = (
        component_supplier_detail.sort_values(["component", "supplier_share", "spend"], ascending=[True, False, False])
        .drop_duplicates("component")[["component", "supplier"]]
        .rename(columns={"supplier": "dominant_supplier"})
    )
    component_summary = component_summary.merge(dominant_supplier, on="component", how="left")

    risk_supplier_count = min_max_scale(1 / component_summary["supplier_count"].replace(0, np.nan).fillna(1))
    risk_concentration = min_max_scale(component_summary["largest_supplier_share"])
    risk_defects = min_max_scale(component_summary["worst_supplier_defect_rate"])
    risk_lead_time = min_max_scale(component_summary["avg_lead_time"])
    risk_external = min_max_scale(component_summary["risk_score"])
    risk_criticality = min_max_scale(component_summary["criticality"])
    component_summary["supply_risk_score"] = (
        0.28 * risk_supplier_count
        + 0.22 * risk_concentration
        + 0.20 * risk_defects
        + 0.12 * risk_lead_time
        + 0.08 * risk_external
        + 0.10 * risk_criticality
    ).round(2)
    component_summary["sourcing_risk_level"] = classify_sourcing_risk_level(
        component_summary["supplier_count"], component_summary["supply_risk_score"]
    )

    leverage_spend = min_max_scale(component_summary["spend"])
    leverage_supplier_count = min_max_scale(component_summary["supplier_count"])
    leverage_inverse_concentration = min_max_scale(1 - component_summary["largest_supplier_share"])
    component_summary["negotiation_leverage_score"] = (
        0.45 * leverage_spend
        + 0.25 * leverage_supplier_count
        + 0.30 * leverage_inverse_concentration
    ).round(2)

    impact_spend = min_max_scale(component_summary["spend"])
    impact_defects = min_max_scale(component_summary["defect_rate"])
    impact_quality_cost = min_max_scale(component_summary["risk_adjusted_spend"] if "risk_adjusted_spend" in component_summary else component_summary["spend"])
    impact_criticality = min_max_scale(component_summary["criticality"])
    component_summary["profit_impact_score"] = (
        0.45 * impact_spend
        + 0.20 * impact_defects
        + 0.20 * impact_criticality
        + 0.15 * impact_quality_cost
    ).round(2)
    component_summary["kraljic_quadrant"] = classify_kraljic_quadrant(
        component_summary["profit_impact_score"], component_summary["supply_risk_score"]
    )
    strategic_risk = min_max_scale(component_summary["supply_risk_score"])
    strategic_impact = min_max_scale(component_summary["profit_impact_score"])
    strategic_single = component_summary["single_source_flag"].astype(float) * 100.0
    strategic_leverage = min_max_scale(component_summary["negotiation_leverage_score"])
    component_summary["strategic_priority_score"] = (
        0.35 * strategic_impact
        + 0.25 * strategic_risk
        + 0.20 * strategic_single
        + 0.20 * strategic_leverage
    ).round(2)
    component_summary["risk_adjusted_spend"] = (
        component_summary["spend"]
        * (1 + component_summary["defect_rate"] * 4)
        * (1 + component_summary["single_source_flag"].astype(float) * 0.35)
        * (1 + component_summary["supply_risk_score"] / 100)
    ).round(2)
    component_summary["high_risk_flag"] = component_summary["sourcing_risk_level"].eq("High")

    spend_pareto = assign_abc_categories(component_summary.sort_values("spend", ascending=False).reset_index(drop=True), "spend", "spend")
    risk_pareto = assign_abc_categories(component_summary.sort_values("risk_adjusted_spend", ascending=False).reset_index(drop=True), "risk_adjusted_spend", "risk")
    strategic_pareto = assign_abc_categories(component_summary.sort_values("strategic_priority_score", ascending=False).reset_index(drop=True), "strategic_priority_score", "strategic")

    component_summary = component_summary.merge(spend_pareto[["component", "spend_cum_share", "spend_abc"]], on="component", how="left")
    component_summary = component_summary.merge(risk_pareto[["component", "risk_cum_share", "risk_abc"]], on="component", how="left")
    component_summary = component_summary.merge(strategic_pareto[["component", "strategic_cum_share", "strategic_abc"]], on="component", how="left")

    protected_components = component_summary.loc[
        component_summary["single_source_flag"] | component_summary["high_risk_flag"],
        ["component", "single_source_flag", "high_risk_flag", "strategic_priority_score", "supply_risk_score"],
    ].copy()

    supplier_base = (
        component_supplier_detail.groupby("supplier", as_index=False)
        .agg(
            spend=("spend", "sum"),
            units=("units", "sum"),
            defects=("defects", "sum"),
            avg_lead_time=("avg_lead_time", "mean"),
            avg_risk_score=("avg_risk_score", "mean"),
            component_count=("component", "nunique"),
        )
    )
    supplier_base["defect_rate"] = safe_divide(supplier_base["defects"], supplier_base["units"], fill_value=0.0)
    supplier_base["portfolio_share"] = safe_divide(supplier_base["spend"], supplier_base["spend"].sum(), fill_value=0.0)

    component_flags = component_summary[
        ["component", "single_source_flag", "high_risk_flag", "strategic_priority_score", "supply_risk_score", "largest_supplier_share"]
    ].copy()
    supplier_component_map = component_supplier_detail[
        ["component", "supplier", "supplier_share"]
    ].merge(component_flags, on="component", how="left")
    supplier_protection = (
        supplier_component_map.groupby("supplier", as_index=False)
        .agg(
            supports_single_source=("single_source_flag", "max"),
            supports_high_risk=("high_risk_flag", "max"),
            protected_component_count=("component", lambda x: int(x.nunique())),
            strategic_priority_supported=("strategic_priority_score", "sum"),
            avg_component_risk=("supply_risk_score", "mean"),
            overlap_components=("component", "nunique"),
            max_component_share=("supplier_share", "max"),
        )
    )
    supplier_summary = supplier_base.merge(supplier_protection, on="supplier", how="left")
    for col in ["supports_single_source", "supports_high_risk"]:
        supplier_summary[col] = supplier_summary[col].fillna(False).astype(bool)
    for col in ["protected_component_count", "strategic_priority_supported", "avg_component_risk", "overlap_components", "max_component_share"]:
        supplier_summary[col] = supplier_summary[col].fillna(0.0)

    supplier_summary["supplier_risk_score"] = (
        0.30 * min_max_scale(supplier_summary["defect_rate"])
        + 0.20 * min_max_scale(supplier_summary["avg_lead_time"])
        + 0.20 * min_max_scale(supplier_summary["avg_component_risk"])
        + 0.15 * min_max_scale(supplier_summary["portfolio_share"])
        + 0.10 * supplier_summary["supports_single_source"].astype(float) * 100.0
        + 0.05 * supplier_summary["supports_high_risk"].astype(float) * 100.0
    ).round(2)

    supplier_summary["replaceability_score"] = (
        0.35 * min_max_scale(1 - supplier_summary["max_component_share"])
        + 0.25 * min_max_scale(1 - supplier_summary["portfolio_share"])
        + 0.20 * min_max_scale(1 - supplier_summary["component_count"] / max(1, supplier_summary["component_count"].max()))
        + 0.20 * (100 - min_max_scale(supplier_summary["strategic_priority_supported"]))
    ).round(2)
    supplier_summary["performance_score"] = (
        0.40 * min_max_scale(supplier_summary["defect_rate"], invert=True)
        + 0.25 * min_max_scale(supplier_summary["avg_lead_time"], invert=True)
        + 0.20 * min_max_scale(supplier_summary["avg_risk_score"], invert=True)
        + 0.15 * min_max_scale(supplier_summary["component_count"])
    ).round(2)

    supplier_decisions = classify_suppliers(component_supplier_detail, component_summary, supplier_summary)
    supplier_summary = supplier_summary.merge(supplier_decisions, on="supplier", how="left")

    component_supplier_detail = component_supplier_detail.merge(
        component_summary[
            [
                "component",
                "high_risk_flag",
                "sourcing_risk_level",
                "supply_risk_score",
                "profit_impact_score",
                "dominant_supplier",
                "kraljic_quadrant",
                "strategic_priority_score",
                "risk_adjusted_spend",
            ]
        ],
        on="component",
        how="left",
    )
    component_supplier_detail = component_supplier_detail.merge(
        supplier_summary[["supplier", "decision", "decision_reason", "estimated_savings"]],
        on="supplier",
        how="left",
    )

    return {
        "supplier_summary": supplier_summary.sort_values(["decision_rank", "spend"], ascending=[True, False]).reset_index(drop=True),
        "component_summary": component_summary.sort_values("strategic_priority_score", ascending=False).reset_index(drop=True),
        "component_supplier_detail": component_supplier_detail.sort_values(["strategic_priority_score", "spend"], ascending=[False, False]).reset_index(drop=True),
        "protected_components": protected_components.sort_values("strategic_priority_score", ascending=False).reset_index(drop=True),
        "spend_pareto": spend_pareto,
        "risk_pareto": risk_pareto,
        "strategic_pareto": strategic_pareto,
    }


@st.cache_data(show_spinner=False)
def build_executive_summary(
    analytics: Dict[str, pd.DataFrame], scenario_applied: bool = False
) -> Tuple[str, pd.DataFrame, pd.DataFrame]:
    supplier_summary = analytics["supplier_summary"]
    component_summary = analytics["component_summary"]
    detail = analytics["component_supplier_detail"]

    top_supplier = supplier_summary.sort_values("spend", ascending=False).iloc[0]
    top_component = component_summary.sort_values("spend", ascending=False).iloc[0]
    highest_risk_component = component_summary.sort_values("supply_risk_score", ascending=False).iloc[0]
    highest_priority_component = component_summary.sort_values("strategic_priority_score", ascending=False).iloc[0]
    quadrant_counts = component_summary["kraljic_quadrant"].value_counts()

    single_source_components = component_summary.loc[component_summary["single_source_flag"], "component"].tolist()
    high_risk_components = component_summary.loc[component_summary["high_risk_flag"], "component"].tolist()
    medium_risk_components = component_summary.loc[component_summary["sourcing_risk_level"].eq("Medium"), "component"].tolist()
    low_risk_components = component_summary.loc[component_summary["sourcing_risk_level"].eq("Low"), "component"].tolist()
    component_supplier_lookup = detail.groupby("component")["supplier"].apply(lambda x: ", ".join(sorted(pd.unique(x)))).to_dict()

    keep_suppliers = supplier_summary.loc[supplier_summary["decision"] == "Keep / Consolidate To", "supplier"].tolist()
    eliminate_suppliers = supplier_summary.loc[supplier_summary["decision"] == "Eliminate / De-prioritize", "supplier"].tolist()
    monitor_suppliers = supplier_summary.loc[supplier_summary["decision"] == "Keep and Monitor", "supplier"].tolist()
    total_savings = supplier_summary["estimated_savings"].sum()

    reason_bits = []
    if highest_priority_component["single_source_flag"]:
        reason_bits.append("single-source exposure")
    if highest_priority_component["defect_rate"] > component_summary["defect_rate"].median():
        reason_bits.append("above-median defects")
    if highest_priority_component["supply_risk_score"] >= component_summary["supply_risk_score"].quantile(0.75):
        reason_bits.append("elevated supply risk")
    if highest_priority_component["spend"] >= component_summary["spend"].quantile(0.75):
        reason_bits.append("top-tier spend")
    priority_reason = ", ".join(reason_bits) if reason_bits else "combined spend and supply exposure"

    exposure_components = sorted(set(single_source_components + high_risk_components))
    exposure_text = "; ".join([f"{comp}: {component_supplier_lookup.get(comp, 'Unknown')}" for comp in exposure_components]) if exposure_components else "no elevated exposure clusters"
    single_source_text = ", ".join(single_source_components) if single_source_components else "none found"
    high_risk_text = ", ".join(high_risk_components) if high_risk_components else "none found"
    medium_risk_text = ", ".join(medium_risk_components) if medium_risk_components else "none found"
    high_risk_count = len(high_risk_components)
    medium_risk_count = len(medium_risk_components)
    low_risk_count = len(low_risk_components)
    quadrant_text = ", ".join([f"{name}: {count}" for name, count in quadrant_counts.items()]) if not quadrant_counts.empty else "no quadrant data available"
    keep_count = len(keep_suppliers)
    eliminate_count = len(eliminate_suppliers)
    monitor_count = len(monitor_suppliers)

    summary_text = (
        f"Kraljic quadrant counts are {quadrant_text}. High sourcing-risk components total {high_risk_count} ({high_risk_text}), medium sourcing-risk components total {medium_risk_count} ({medium_risk_text}), and low sourcing-risk components total {low_risk_count}. "
        f"Single-source components are {single_source_text}. Suppliers covering those exposed components include {exposure_text}. "
        f"The highest-risk component is {highest_risk_component['component']} with a supply risk score of {highest_risk_component['supply_risk_score']:.1f}. "
        f"The highest strategic priority component is {highest_priority_component['component']} because of {priority_reason}. "
        f"Top supplier spend sits with {top_supplier['supplier']} at ${top_supplier['spend']:,.0f}, while {top_component['component']} is the largest-spend component at ${top_component['spend']:,.0f}. "
        f"At this stage, the base analysis indicates a supplier landscape with {keep_count} stronger consolidation candidates, {eliminate_count} potential exit or de-prioritization candidates, and {monitor_count} suppliers that may warrant closer review. "
        f"The next step is to run scenarios to test whether single-source exposure, high-risk components, and supplier fragmentation can be reduced without creating new uncovered demand or unacceptable savings tradeoffs. "
        f"Specific supplier actions, economic tradeoffs, and rationale should be confirmed only after an applied scenario is evaluated."
    )

    executive_actions_source = supplier_summary.copy()
    supplier_action_plan_source = supplier_summary.copy()
    if scenario_applied:
        executive_actions_source["decision_display"] = executive_actions_source["decision"]
        executive_actions_source["action_display"] = executive_actions_source["supplier_action_plan"]
        executive_actions_source["savings_display"] = executive_actions_source["estimated_savings"]
        supplier_action_plan_source["decision_display"] = supplier_action_plan_source["decision"]
        supplier_action_plan_source["action_display"] = supplier_action_plan_source["supplier_action_plan"]
        supplier_action_plan_source["justification_display"] = supplier_action_plan_source["decision_reason"]
        supplier_action_plan_source["savings_display"] = supplier_action_plan_source["estimated_savings"]
    else:
        executive_actions_source["decision_display"] = np.select(
            [
                executive_actions_source["decision"].eq("Eliminate / De-prioritize"),
                executive_actions_source["decision"].eq("Keep / Consolidate To"),
            ],
            ["Evaluate Exit Scenario", "Evaluate Consolidation Scenario"],
            default="Evaluate Monitoring Need",
        )
        executive_actions_source["action_display"] = np.select(
            [
                executive_actions_source["decision"].eq("Eliminate / De-prioritize"),
                executive_actions_source["decision"].eq("Keep / Consolidate To"),
            ],
            [
                "Run scenarios to test whether demand can move away from this supplier without creating new risk.",
                "Run scenarios to test whether consolidating more volume here improves the score without increasing exposure.",
            ],
            default="Review quality, lead time, and exposure trends before deciding whether intervention is needed.",
        )
        executive_actions_source["savings_display"] = np.nan
        supplier_action_plan_source["decision_display"] = executive_actions_source["decision_display"]
        supplier_action_plan_source["action_display"] = executive_actions_source["action_display"]
        supplier_action_plan_source["justification_display"] = np.select(
            [
                supplier_action_plan_source["decision"].eq("Eliminate / De-prioritize"),
                supplier_action_plan_source["decision"].eq("Keep / Consolidate To"),
            ],
            [
                "Current data suggests this supplier may be replaceable, but scenario testing should confirm whether exits are practical.",
                "Current data suggests this supplier may be a strong consolidation candidate, but scenario testing should confirm whether expanding share is beneficial.",
            ],
            default="Current data suggests this supplier may need closer review, but scenario testing should confirm whether action is required.",
        )
        supplier_action_plan_source["savings_display"] = np.nan

    executive_actions = executive_actions_source[
        ["supplier", "decision_display", "issues", "action_display", "savings_display"]
    ].rename(
        columns={
            "supplier": "Supplier",
            "decision_display": "Decision",
            "issues": "Issues",
            "action_display": "Action Plan",
            "savings_display": "Estimated Savings",
        }
    )

    supplier_action_plan = supplier_action_plan_source[
        ["supplier", "decision_display", "issues", "action_display", "savings_display", "justification_display"]
    ].rename(
        columns={
            "supplier": "Supplier",
            "decision_display": "Decision",
            "issues": "Issues",
            "action_display": "Action Plan",
            "savings_display": "Estimated Savings",
            "justification_display": "Justification",
        }
    )
    if scenario_applied and "scenario_role_rank" in supplier_summary.columns and "decision_rank_display" in supplier_summary.columns:
        order_lookup = supplier_summary.set_index("supplier")[["scenario_role_rank", "decision_rank_display", "spend"]].to_dict(orient="index")
        executive_actions["__role_rank"] = executive_actions["Supplier"].map(lambda name: order_lookup.get(name, {}).get("scenario_role_rank", 9))
        executive_actions["__decision_rank"] = executive_actions["Supplier"].map(lambda name: order_lookup.get(name, {}).get("decision_rank_display", 9))
        executive_actions["__spend"] = executive_actions["Supplier"].map(lambda name: order_lookup.get(name, {}).get("spend", 0.0))
        executive_actions = executive_actions.sort_values(["__role_rank", "__decision_rank", "__spend", "Supplier"], ascending=[True, True, False, True]).drop(columns=["__role_rank", "__decision_rank", "__spend"]).reset_index(drop=True)
        supplier_action_plan["__role_rank"] = supplier_action_plan["Supplier"].map(lambda name: order_lookup.get(name, {}).get("scenario_role_rank", 9))
        supplier_action_plan["__decision_rank"] = supplier_action_plan["Supplier"].map(lambda name: order_lookup.get(name, {}).get("decision_rank_display", 9))
        supplier_action_plan["__spend"] = supplier_action_plan["Supplier"].map(lambda name: order_lookup.get(name, {}).get("spend", 0.0))
        supplier_action_plan = supplier_action_plan.sort_values(["__role_rank", "__decision_rank", "__spend", "Supplier"], ascending=[True, True, False, True]).drop(columns=["__role_rank", "__decision_rank", "__spend"]).reset_index(drop=True)
    return summary_text, executive_actions, supplier_action_plan


@st.cache_data(show_spinner=False)
def build_step_plan(analytics: Dict[str, pd.DataFrame], scenario_applied: bool = False) -> pd.DataFrame:
    supplier_summary = analytics["supplier_summary"]
    component_summary = analytics["component_summary"]

    protected = component_summary.loc[component_summary["single_source_flag"] | component_summary["high_risk_flag"]]
    elimination_targets = supplier_summary.loc[supplier_summary["decision"] == "Eliminate / De-prioritize"]
    monitor_targets = supplier_summary.loc[supplier_summary["decision"] == "Keep and Monitor"]
    keep_targets = supplier_summary.loc[supplier_summary["decision"] == "Keep / Consolidate To"]

    rows = (
        [
            {
                "step": 1,
                "action": "Stabilize exposed components",
                "justification": "Single-source and high-risk components must be protected before any sourcing shifts.",
                "supporting evidence": ", ".join(protected["component"].tolist()) if not protected.empty else "No protected components identified",
                "expected outcome": "Continuity risk reduced and elimination actions gated behind supply protection.",
            },
            {
                "step": 2,
                "action": "Consolidate award share to retained suppliers",
                "justification": "Preferred suppliers combine stronger performance, overlap coverage, and strategic importance.",
                "supporting evidence": ", ".join(keep_targets["supplier"].tolist()) if not keep_targets.empty else "No consolidation targets identified",
                "expected outcome": "Higher leverage, cleaner supplier base, and fewer fragmented awards.",
            },
            {
                "step": 3,
                "action": "De-prioritize replaceable low-performing suppliers",
                "justification": "Suppliers marked for elimination have weaker performance and low protected dependency.",
                "supporting evidence": ", ".join(elimination_targets["supplier"].tolist()) if not elimination_targets.empty else "No elimination targets identified",
                "expected outcome": "Cost savings, reduced quality leakage, and tighter supplier governance.",
            },
            {
                "step": 4,
                "action": "Launch corrective scorecards for monitored suppliers",
                "justification": "Monitoring targets still matter to coverage, but they need measurable performance improvement.",
                "supporting evidence": ", ".join(monitor_targets["supplier"].tolist()) if not monitor_targets.empty else "No monitor targets identified",
                "expected outcome": "Better lead time and quality without unnecessary churn.",
            },
            {
                "step": 5,
                "action": "Prioritize sourcing work by strategic Pareto",
                "justification": "The highest strategic-priority components combine spend, risk, defects, and source exposure.",
                "supporting evidence": ", ".join(component_summary.sort_values("strategic_priority_score", ascending=False).head(5)["component"].tolist()),
                "expected outcome": "Management attention is focused where value and risk are both highest.",
            },
        ]
        if scenario_applied
        else [
            {
                "step": 1,
                "action": "Test exposure-reduction scenarios",
                "justification": "Single-source and high-risk components should be the first constraints tested in scenario design.",
                "supporting evidence": ", ".join(protected["component"].tolist()) if not protected.empty else "No protected components identified",
                "expected outcome": "A short list of scenario options that reduce exposure without creating uncovered demand.",
            },
            {
                "step": 2,
                "action": "Test consolidation upside",
                "justification": "Stronger suppliers may offer cleaner consolidation paths, but score and resilience tradeoffs should be validated.",
                "supporting evidence": ", ".join(keep_targets["supplier"].tolist()) if not keep_targets.empty else "No consolidation candidates identified",
                "expected outcome": "A comparison of whether consolidation improves overall score, coverage, and economics.",
            },
            {
                "step": 3,
                "action": "Test exit sensitivity",
                "justification": "More replaceable suppliers should be evaluated through scenarios before any exit recommendation is treated as actionable.",
                "supporting evidence": ", ".join(elimination_targets["supplier"].tolist()) if not elimination_targets.empty else "No exit candidates identified",
                "expected outcome": "Visibility into whether supplier exits create new risk, uncovered components, or weak net savings.",
            },
            {
                "step": 4,
                "action": "Review monitoring candidates",
                "justification": "Some suppliers may need watchlist treatment, but scenario testing should confirm whether active intervention is required.",
                "supporting evidence": ", ".join(monitor_targets["supplier"].tolist()) if not monitor_targets.empty else "No monitor candidates identified",
                "expected outcome": "A clear distinction between suppliers that need action and suppliers that simply need observation.",
            },
            {
                "step": 5,
                "action": "Prioritize scenario work by strategic Pareto",
                "justification": "The highest strategic-priority components combine spend, risk, defects, and source exposure.",
                "supporting evidence": ", ".join(component_summary.sort_values("strategic_priority_score", ascending=False).head(5)["component"].tolist()),
                "expected outcome": "Scenario analysis focused where value and risk are both highest.",
            },
        ]
    )
    return pd.DataFrame(rows)


@st.cache_data(show_spinner=False)
def build_supplier_consolidation_plan(
    analytics: Dict[str, pd.DataFrame], scenario_applied: bool = False
) -> Tuple[pd.DataFrame, List[str]]:
    supplier_summary = analytics["supplier_summary"].copy()
    plan = supplier_summary[
        [
            "supplier",
            "decision",
            "spend",
            "component_count",
            "protected_component_count",
            "performance_score",
            "replaceability_score",
            "estimated_savings",
            "decision_reason",
        ]
    ].copy()
    if scenario_applied:
        plan["decision_display"] = plan["decision"]
        plan["reason_display"] = plan["decision_reason"]
        plan["savings_display"] = plan["estimated_savings"]
    else:
        plan["decision_display"] = np.select(
            [
                plan["decision"].eq("Eliminate / De-prioritize"),
                plan["decision"].eq("Keep / Consolidate To"),
            ],
            ["Evaluate Exit Scenario", "Evaluate Consolidation Scenario"],
            default="Evaluate Monitoring Need",
        )
        plan["reason_display"] = np.select(
            [
                plan["decision"].eq("Eliminate / De-prioritize"),
                plan["decision"].eq("Keep / Consolidate To"),
            ],
            [
                "Current data suggests this supplier may be replaceable, but scenario testing should confirm whether exits are practical.",
                "Current data suggests this supplier may be a strong consolidation candidate, but scenario testing should confirm whether expanding share is beneficial.",
            ],
            default="Current data suggests this supplier may need closer review, but scenario testing should confirm whether action is required.",
        )
        plan["savings_display"] = np.nan
    if scenario_applied:
        plan["priority"] = np.select(
            [
                plan["decision"].eq("Eliminate / De-prioritize"),
                plan["decision"].eq("Keep / Consolidate To"),
            ],
            ["Exit / Reallocate", "Expand / Consolidate"],
            default="Monitor / Improve",
        )
        plan["owner_action"] = np.select(
            [
                plan["decision"].eq("Eliminate / De-prioritize"),
                plan["decision"].eq("Keep / Consolidate To"),
            ],
            [
                "Shift volume to stronger alternatives after qualification and safety-stock checks.",
                "Increase award share, negotiate pricing, and lock continuity protections.",
            ],
            default="Run corrective scorecard on quality, lead time, and resilience metrics.",
        )
    else:
        plan["priority"] = np.select(
            [
                plan["decision"].eq("Eliminate / De-prioritize"),
                plan["decision"].eq("Keep / Consolidate To"),
            ],
            ["Evaluate Exit Exposure", "Evaluate Consolidation Potential"],
            default="Review Performance Risk",
        )
        plan["owner_action"] = np.select(
            [
                plan["decision"].eq("Eliminate / De-prioritize"),
                plan["decision"].eq("Keep / Consolidate To"),
            ],
            [
                "Run scenarios to test whether demand can move away from this supplier without creating new risk.",
                "Run scenarios to test whether consolidating more volume here improves the score without increasing exposure.",
            ],
            default="Review quality, lead time, and exposure trends before deciding whether intervention is needed.",
        )
    plan = plan.rename(
        columns={
            "supplier": "Supplier",
            "decision_display": ("Recommendation" if scenario_applied else "Scenario Test"),
            "spend": "Spend",
            "component_count": "Components Covered",
            "protected_component_count": "Protected Components",
            "performance_score": "Performance Score",
            "replaceability_score": "Replaceability Score",
            "savings_display": "Estimated Savings",
            "reason_display": "Why",
            "priority": "Priority",
            "owner_action": "Action",
        }
    )
    plan = plan[
        [
            "Supplier",
            ("Recommendation" if scenario_applied else "Scenario Test"),
            "Spend",
            "Components Covered",
            "Protected Components",
            "Performance Score",
            "Replaceability Score",
            "Estimated Savings",
            "Why",
            "Priority",
            "Action",
        ]
    ].copy()
    if scenario_applied and "scenario_role_rank" in supplier_summary.columns and "decision_rank_display" in supplier_summary.columns:
        order_lookup = supplier_summary.set_index("supplier")[["scenario_role_rank", "decision_rank_display", "spend"]].to_dict(orient="index")
        plan["__role_rank"] = plan["Supplier"].map(lambda name: order_lookup.get(name, {}).get("scenario_role_rank", 9))
        plan["__decision_rank"] = plan["Supplier"].map(lambda name: order_lookup.get(name, {}).get("decision_rank_display", 9))
        plan["__spend"] = plan["Supplier"].map(lambda name: order_lookup.get(name, {}).get("spend", 0.0))
        plan = plan.sort_values(["__role_rank", "__decision_rank", "__spend", "Supplier"], ascending=[True, True, False, True]).drop(columns=["__role_rank", "__decision_rank", "__spend"]).reset_index(drop=True)
    else:
        plan = plan.sort_values(["Priority", "Estimated Savings"], ascending=[True, False]).reset_index(drop=True)

    assumptions = [
        (
            "Base-view consolidation outputs are directional signals to investigate through scenarios, not committed supplier moves."
            if not scenario_applied
            else "Consolidation recommendations favor suppliers with stronger quality, lead time, overlap coverage, and lower replacement risk."
        ),
        "Single-source and high sourcing-risk coverage is protected before any supplier exit or volume shift.",
        "Estimated savings assume a directional savings rate tied to the recommendation class, not contracted pricing.",
        "Supplier transitions assume qualification, inventory, and continuity safeguards are completed before reallocation.",
    ]
    return plan, assumptions


@st.cache_data(show_spinner=False)
def build_supplier_risk_assessment(
    analytics: Dict[str, pd.DataFrame], scenario_applied: bool = False
) -> Tuple[pd.DataFrame, List[str]]:
    supplier_summary = analytics["supplier_summary"].copy()
    defect_cutoff = float(supplier_summary["defect_rate"].quantile(0.75)) if len(supplier_summary) else 0.0
    lead_cutoff = float(supplier_summary["avg_lead_time"].quantile(0.75)) if len(supplier_summary) else 0.0
    external_cutoff = float(supplier_summary["avg_risk_score"].quantile(0.75)) if len(supplier_summary) else 0.0

    supplier_summary["supplier_risk_tier"] = np.select(
        [
            supplier_summary["supports_single_source"] | supplier_summary["supports_high_risk"],
            (supplier_summary["defect_rate"] >= defect_cutoff)
            | (supplier_summary["avg_lead_time"] >= lead_cutoff)
            | (supplier_summary["avg_risk_score"] >= external_cutoff),
        ],
        ["High", "Medium"],
        default="Low",
    )

    supplier_summary["risk_driver_summary"] = supplier_summary.apply(
        lambda row: ", ".join(
            [
                label
                for condition, label in [
                    (bool(row["supports_single_source"]), "supports single-source component"),
                    (bool(row["supports_high_risk"]), "supports high-risk component"),
                    (float(row["defect_rate"]) >= defect_cutoff, "higher defect rate"),
                    (float(row["avg_lead_time"]) >= lead_cutoff, "longer lead time"),
                    (float(row["avg_risk_score"]) >= external_cutoff, "elevated external risk"),
                ]
                if condition
            ]
        )
        or "stable relative to peer suppliers",
        axis=1,
    )

    risk_table = supplier_summary[
        [
            "supplier",
            "supplier_risk_tier",
            "spend",
            "supplier_risk_score",
            "defect_rate",
            "avg_lead_time",
            "avg_risk_score",
            "protected_component_count",
            "risk_driver_summary",
        ]
    ].rename(
        columns={
            "supplier": "Supplier",
            "supplier_risk_tier": "Risk Tier",
            "spend": "Spend",
            "supplier_risk_score": "Supplier Risk Score",
            "defect_rate": "Defect Rate",
            "avg_lead_time": "Avg Lead Time",
            "avg_risk_score": "External Risk Score",
            "protected_component_count": "Protected Components",
            "risk_driver_summary": "Risk Drivers",
        }
    )
    risk_table = risk_table.sort_values(["Risk Tier", "Spend"], ascending=[True, False]).reset_index(drop=True)

    assumptions = [
        "Supplier risk is assessed from operational quality, lead time, external risk score, and whether the supplier supports protected components.",
        "Any supplier supporting a single-source or high sourcing-risk component is elevated because disruption would hit exposed demand first.",
        "Medium risk indicates peer-relative weakness rather than imminent failure; it signals watchlist and mitigation work.",
        "Risk tiers are comparative within the uploaded dataset and should be complemented by financial, geopolitical, and contract review outside the file.",
        (
            "Before a scenario is applied, this table should be read as a risk-screening view to inform scenario design rather than a final action list."
            if not scenario_applied
            else "Applied-scenario risk tiers reflect the chosen supplier structure and should be used to monitor the implemented scenario."
        ),
    ]
    return risk_table, assumptions


@st.cache_data(show_spinner=False)
def build_strategic_sourcing_plan(
    analytics: Dict[str, pd.DataFrame], scenario_applied: bool = False
) -> Tuple[pd.DataFrame, List[str]]:
    component_summary = analytics["component_summary"].copy()
    strategy_map = (
        {
            "Strategic": "Executive supplier partnership, joint business planning, and resilience investment.",
            "Leverage": "Competitive bidding, should-cost negotiation, and disciplined volume consolidation.",
            "Bottleneck": "Supply assurance, backup qualification, inventory buffering, and spec flexibility review.",
            "Non-Critical": "Transaction efficiency, automation, and simplified ordering governance.",
        }
        if scenario_applied
        else {
            "Strategic": "Test scenarios that protect critical supply while preserving leverage where possible.",
            "Leverage": "Test whether sourcing competition or consolidation could improve cost without raising supply risk.",
            "Bottleneck": "Test backup qualification, inventory buffering, or spec flexibility options.",
            "Non-Critical": "Review whether transactional simplification or automation opportunities exist.",
        }
    )
    component_summary["sourcing_strategy"] = component_summary["kraljic_quadrant"].map(strategy_map)
    component_summary["next_step"] = np.select(
        [
            component_summary["kraljic_quadrant"].eq("Strategic"),
            component_summary["kraljic_quadrant"].eq("Leverage"),
            component_summary["kraljic_quadrant"].eq("Bottleneck"),
        ],
        (
            [
                "Build supplier executive review, continuity plan, and performance roadmap.",
                "Run sourcing event and negotiate across retained suppliers.",
                "Qualify alternatives and protect service with safety stock or dual-source actions.",
            ]
            if scenario_applied
            else [
                "Run scenarios to see which supplier structures reduce strategic exposure without hurting continuity.",
                "Run scenarios to compare consolidation, competition, and savings tradeoffs.",
                "Run scenarios to test backup options and service-protection levers.",
            ]
        ),
        default=(
            "Reduce touch labor and standardize replenishment controls."
            if scenario_applied
            else "Review whether simplification opportunities exist before changing the supply structure."
        ),
    )
    plan = component_summary[
        [
            "component",
            "kraljic_quadrant",
            "sourcing_risk_level",
            "supplier_count",
            "spend",
            "defect_rate",
            "sourcing_strategy",
            "next_step",
        ]
    ].rename(
        columns={
            "component": "Component",
            "kraljic_quadrant": "Kraljic Quadrant",
            "sourcing_risk_level": "Sourcing Risk",
            "supplier_count": "Supplier Count",
            "spend": "Spend",
            "defect_rate": "Defect Rate",
            "sourcing_strategy": "Strategy",
            "next_step": "Next Step",
        }
    )
    plan = plan.sort_values(["Kraljic Quadrant", "Spend"], ascending=[True, False]).reset_index(drop=True)

    assumptions = [
        "Kraljic positioning uses profit impact from spend and quality burden, and supply risk from sourcing concentration, defects, lead time, and external risk.",
        "Single-source exposure is treated as structurally high sourcing risk even when the blended score is moderate.",
        "Strategic components require relationship and resilience investment; leverage components favor commercial pressure; bottleneck components favor continuity protection.",
        "The plan assumes current demand mix is representative and that supplier coverage in the file reflects the practical sourcing options available.",
        (
            "In the base view, these strategy lines are prompts for scenario exploration rather than committed sourcing actions."
            if not scenario_applied
            else "In the applied view, these strategy lines reflect the selected scenario and can be used as execution guidance."
        ),
    ]
    return plan, assumptions


@st.cache_data(show_spinner=False)
def build_visual_pack(analytics: Dict[str, pd.DataFrame]) -> List[Dict[str, object]]:
    supplier_summary = analytics["supplier_summary"]
    component_summary = analytics["component_summary"]
    component_supplier_detail = analytics["component_supplier_detail"]
    spend_pareto = analytics["spend_pareto"]
    risk_pareto = analytics["risk_pareto"]
    strategic_pareto = analytics["strategic_pareto"]
    top_supplier = supplier_summary.sort_values("spend", ascending=False).iloc[0] if len(supplier_summary) else None
    top_component = component_summary.sort_values("spend", ascending=False).iloc[0] if len(component_summary) else None
    top_concentration = component_summary.sort_values("largest_supplier_share", ascending=False).iloc[0] if len(component_summary) else None
    top_supplier_risk = supplier_summary.sort_values("supplier_risk_score", ascending=False).iloc[0] if len(supplier_summary) else None
    top_risk_adjusted = component_summary.sort_values("risk_adjusted_spend", ascending=False).iloc[0] if len(component_summary) else None
    keep_count = int(supplier_summary["decision"].eq("Keep / Consolidate To").sum()) if len(supplier_summary) else 0
    monitor_count = int(supplier_summary["decision"].eq("Keep and Monitor").sum()) if len(supplier_summary) else 0
    exit_count = int(supplier_summary["decision"].eq("Eliminate / De-prioritize").sum()) if len(supplier_summary) else 0
    return [
        {
            "title": "Spend by Supplier",
            "data": supplier_summary.set_index("supplier")[["spend"]].sort_values("spend", ascending=False),
            "summary": (
                "This visual shows where the supplier spend base is concentrated so teams can identify leverage points and dependency risk. "
                f"Major takeaway: {top_supplier['supplier']} carries the largest supplier spend at ${top_supplier['spend']:,.0f}."
                if top_supplier is not None
                else "This visual shows where the supplier spend base is concentrated so teams can identify leverage points and dependency risk."
            ),
            "talking_points": [
                "Shows where the spend base is concentrated across suppliers.",
                "Use retained suppliers at the top of the chart to anchor negotiation waves.",
            ],
        },
        {
            "title": "Spend by Component",
            "data": component_summary.set_index("component")[["spend"]].sort_values("spend", ascending=False),
            "summary": (
                "This visual highlights the components with the biggest commercial impact so sourcing teams know where price, quality, and continuity matter most. "
                f"Major takeaway: {top_component['component']} is the highest-spend component at ${top_component['spend']:,.0f}."
                if top_component is not None
                else "This visual highlights the components with the biggest commercial impact so sourcing teams know where price, quality, and continuity matter most."
            ),
            "talking_points": [
                "Shows which components drive the largest commercial impact.",
                "Use alongside quadrant placement to separate leverage opportunities from bottleneck exposure.",
            ],
        },
        {
            "title": "Component Analysis Bubble",
            "data": component_summary[
                ["component", "spend", "supplier_count", "sourcing_risk_level", "supply_risk_score", "strategic_priority_score", "kraljic_quadrant"]
            ].set_index("component"),
            "summary": build_component_analysis_summary(component_summary),
            "talking_points": [
                "Bubble size shows strategic priority, while color shows sourcing risk.",
                "Use it to spot components that are simultaneously high spend, high risk, and strategically important.",
            ],
        },
        {
            "title": "Component-Supplier Detail",
            "data": component_supplier_detail[
                ["component", "supplier", "spend", "supplier_share", "effective_supplier_share", "supplier_count", "single_source_flag"]
            ],
            "summary": build_component_supplier_detail_summary(component_supplier_detail),
            "talking_points": [
                "Shows how component coverage is distributed across suppliers.",
                "Use it to spot concentration, backup depth, and supplier overlap by component.",
            ],
        },
        {
            "title": "Spend Pareto (ABC)",
            "data": spend_pareto[["component", "spend", "spend_cum_share", "spend_abc"]].set_index("component"),
            "summary": build_pareto_summary(spend_pareto, risk_pareto, strategic_pareto),
            "talking_points": [
                "Bars show each component's spend; the line shows cumulative share.",
                "Use it to identify the A-items that deserve the most commercial attention.",
            ],
        },
        {
            "title": "Supplier Concentration by Component",
            "data": component_summary.set_index("component")[["largest_supplier_share"]].sort_values("largest_supplier_share", ascending=False),
            "summary": (
                "This visual shows which components are most dependent on one supplier, which is a direct resilience concern. "
                f"Major takeaway: {top_concentration['component']} has the highest supplier concentration at {top_concentration['largest_supplier_share']:.0%}."
                if top_concentration is not None
                else "This visual shows which components are most dependent on one supplier, which is a direct resilience concern."
            ),
            "talking_points": [
                "High largest-supplier share indicates concentration and lower resilience.",
                "Components near 100% share need backup qualification before supplier exits.",
            ],
        },
        {
            "title": "Supplier Risk Score",
            "data": supplier_summary.set_index("supplier")[["supplier_risk_score"]].sort_values("supplier_risk_score", ascending=False),
            "summary": (
                "This visual ranks suppliers by operational and exposure-driven risk so teams can see where supplier management effort is most needed. "
                f"Major takeaway: {top_supplier_risk['supplier']} has the highest supplier risk score at {top_supplier_risk['supplier_risk_score']:.1f}."
                if top_supplier_risk is not None
                else "This visual ranks suppliers by operational and exposure-driven risk so teams can see where supplier management effort is most needed."
            )
            + " "
            + build_supplier_risk_methodology_note(),
            "talking_points": [
                "Highlights suppliers with elevated operational and exposure-driven risk.",
                "Cross-check these suppliers against protected components before consolidation decisions.",
            ],
        },
        {
            "title": "Risk-Adjusted Pareto",
            "data": component_summary.set_index("component")[["risk_adjusted_spend"]].sort_values("risk_adjusted_spend", ascending=False),
            "summary": (
                "This visual combines spend with quality burden, supply concentration, and sourcing risk so leaders can focus on the parts that matter most in practice, not just on paper. "
                f"Major takeaway: {top_risk_adjusted['component']} is the largest risk-adjusted burden in the portfolio."
                if top_risk_adjusted is not None
                else "This visual combines spend with quality burden, supply concentration, and sourcing risk so leaders can focus on the parts that matter most in practice, not just on paper."
            ),
            "talking_points": [
                "This chart includes spend, defects, single-source exposure, and supply risk in one view.",
                "Use it instead of raw spend when sequencing executive interventions.",
            ],
        },
        {
            "title": "Strategic Priority Pareto",
            "data": strategic_pareto[["component", "strategic_priority_score", "strategic_cum_share", "strategic_abc"]].set_index("component"),
            "summary": "This visual ranks components by overall strategic importance so teams can focus scenario work where business impact and supply exposure are both highest.",
            "talking_points": [
                "Bars show strategic priority score; the line shows cumulative share.",
                "Use it to identify the first components to review in sourcing strategy discussions.",
            ],
        },
        {
            "title": "Kraljic Positioning",
            "data": component_summary[
                ["component", "supplier_count", "sourcing_risk_level", "supply_risk_score", "profit_impact_score", "kraljic_quadrant"]
            ].set_index("component"),
            "summary": build_kraljic_positioning_summary(component_summary) + " " + build_risk_score_methodology_note(),
            "talking_points": [
                "Shows which components are Strategic, Leverage, Bottleneck, or Non-Critical.",
                "Use it to align sourcing strategy to component risk and business impact.",
            ],
        },
        {
            "title": "Supply Risk Score",
            "data": component_summary[["component", "supply_risk_score", "sourcing_risk_level", "supplier_count"]].set_index("component"),
            "summary": build_risk_analysis_summary(component_summary, build_supplier_risk_assessment(analytics)[0])
            + " "
            + build_risk_score_methodology_note(),
            "talking_points": [
                "Ranks components by modeled supply risk.",
                "Use it to identify where mitigation, backup qualification, or closer supplier management is needed first.",
            ],
        },
        {
            "title": "Supplier Count by Component",
            "data": component_summary[["component", "supplier_count", "sourcing_risk_level"]].set_index("component"),
            "summary": "This visual shows how many suppliers support each component so the tightest resilience constraints are easy to spot.",
            "talking_points": [
                "Components near one supplier are the most structurally constrained.",
                "Use it alongside concentration and risk charts to prioritize backup qualification work.",
            ],
        },
        {
            "title": "Strategic Sourcing Outcomes",
            "data": supplier_summary[["supplier", "decision", "spend", "estimated_savings"]].sort_values("spend", ascending=False),
            "summary": (
                "This visual summarizes the current supplier outcome model so teams can see which suppliers look stronger for consolidation, which need monitoring, and which appear more replaceable. "
                f"Major takeaway: the current portfolio suggests {keep_count} keep/consolidate suppliers, {monitor_count} monitor suppliers, and {exit_count} potential exit suppliers."
            ),
            "talking_points": [
                "Shows the recommended sourcing outcome for each supplier against its spend base.",
                "Use it to see which suppliers are consolidation targets, watchlist suppliers, and exit candidates.",
            ],
        },
        {
            "title": "Supplier Spend by Component Mix",
            "data": component_supplier_detail[
                [col for col in ["supplier", "component", "spend", "kraljic_quadrant", "single_source_flag", "sourcing_risk_level"] if col in component_supplier_detail.columns]
            ].copy(),
            "summary": build_supplier_component_mix_summary(component_supplier_detail, component_summary, supplier_summary),
            "talking_points": [
                "Shows which components make up each supplier's spend base.",
                "Use it to see whether a supplier's importance comes from broad portfolio coverage or dependence on a few concentrated components.",
            ],
        },
        {
            "title": "Supplier Spend by Kraljic Mix",
            "data": component_supplier_detail[
                [col for col in ["supplier", "component", "spend", "kraljic_quadrant"] if col in component_supplier_detail.columns]
            ].copy(),
            "summary": build_supplier_quadrant_mix_summary(component_supplier_detail, component_summary, supplier_summary),
            "talking_points": [
                "Shows whether supplier spend is concentrated in Strategic, Bottleneck, Leverage, or Non-Critical components.",
                "Use it to tell the difference between a large supplier and a strategically exposed supplier.",
            ],
        },
    ]


def build_kraljic_chart(component_summary: pd.DataFrame) -> alt.Chart:
    chart_data = component_summary.copy()
    x_mid = float(chart_data["supply_risk_score"].median()) if len(chart_data) else 0.0
    y_mid = float(chart_data["profit_impact_score"].median()) if len(chart_data) else 0.0
    x_max = float(chart_data["supply_risk_score"].max()) if len(chart_data) else 100.0
    y_max = float(chart_data["profit_impact_score"].max()) if len(chart_data) else 100.0
    x_span = max(x_max - x_mid, x_mid, 1.0)
    y_span = max(y_max - y_mid, y_mid, 1.0)
    quadrant_labels = pd.DataFrame(
        [
            {"x": x_mid + x_span * 0.45, "y": y_mid + y_span * 0.45, "label": "Strategic"},
            {"x": max(x_mid - x_span * 0.45, 0.0), "y": y_mid + y_span * 0.45, "label": "Leverage"},
            {"x": x_mid + x_span * 0.45, "y": max(y_mid - y_span * 0.45, 0.0), "label": "Bottleneck"},
            {"x": max(x_mid - x_span * 0.45, 0.0), "y": max(y_mid - y_span * 0.45, 0.0), "label": "Non-Critical"},
        ]
    )

    points = (
        alt.Chart(chart_data)
        .mark_circle(opacity=0.75, stroke="white", strokeWidth=1)
        .encode(
            x=alt.X("supply_risk_score:Q", title="Supply Risk Score"),
            y=alt.Y("profit_impact_score:Q", title="Profit Impact Score"),
            size=alt.Size("spend:Q", title="Spend", scale=alt.Scale(range=[180, 1600])),
            color=alt.Color(
                "sourcing_risk_level:N",
                title="Sourcing Risk",
                scale=alt.Scale(domain=RISK_LEVEL_DOMAIN, range=RISK_LEVEL_RANGE),
            ),
            tooltip=["component", "kraljic_quadrant", "sourcing_risk_level", "spend", "supplier_count", "defect_rate"],
        )
    )
    labels = (
        alt.Chart(chart_data)
        .mark_text(fontSize=11, baseline="middle", align="left", dx=8, color="#111")
        .encode(
            x="supply_risk_score:Q",
            y="profit_impact_score:Q",
            text="component:N",
        )
    )
    vertical = alt.Chart(pd.DataFrame({"x": [x_mid]})).mark_rule(strokeDash=[6, 6], color="#666").encode(x="x:Q")
    horizontal = alt.Chart(pd.DataFrame({"y": [y_mid]})).mark_rule(strokeDash=[6, 6], color="#666").encode(y="y:Q")
    quadrant_text = (
        alt.Chart(quadrant_labels)
        .mark_text(fontSize=14, fontWeight="bold", color="#555", opacity=0.65)
        .encode(
            x="x:Q",
            y="y:Q",
            text="label:N",
        )
    )
    return (
        (points + labels + vertical + horizontal + quadrant_text)
        .properties(height=560)
        .configure_axis(labelFontSize=12, titleFontSize=13)
        .configure_legend(labelFontSize=12, titleFontSize=13)
    )


def build_component_risk_bar_chart(
    component_summary: pd.DataFrame,
    metric: str,
    title: str,
    top_n: Optional[int] = 12,
    ascending: bool = False,
) -> alt.Chart:
    sorted_data = component_summary.sort_values(metric, ascending=ascending).copy()
    chart_data = sorted_data.head(top_n).copy() if top_n is not None else sorted_data
    return (
        alt.Chart(chart_data)
        .mark_bar()
        .encode(
            x=alt.X(f"{metric}:Q", title=title),
            y=alt.Y("component:N", sort="-x", title="Component"),
            color=alt.Color(
                "sourcing_risk_level:N",
                title="Sourcing Risk",
                scale=alt.Scale(domain=RISK_LEVEL_DOMAIN, range=RISK_LEVEL_RANGE),
            ),
            tooltip=["component", "sourcing_risk_level", metric, "supplier_count", "largest_supplier_share"],
        )
        .properties(height=360)
    )


def build_supplier_metric_chart(
    supplier_summary: pd.DataFrame,
    metric: str,
    title: str,
    top_n: Optional[int] = None,
    ascending: bool = False,
) -> alt.Chart:
    sorted_data = supplier_summary.sort_values(metric, ascending=ascending).copy()
    chart_data = sorted_data.head(top_n).copy() if top_n is not None else sorted_data
    color_scale = alt.Scale(
        domain=["Eliminate / De-prioritize", "Keep and Monitor", "Keep / Consolidate To"],
        range=["#d73027", "#f39c12", "#2e8b57"],
    )
    return (
        alt.Chart(chart_data)
        .mark_bar()
        .encode(
            x=alt.X("supplier:N", title="Supplier", axis=alt.Axis(labelAngle=0, labelLimit=220)),
            y=alt.Y(f"{metric}:Q", title=title),
            color=alt.Color("decision:N", title="Recommendation", scale=color_scale),
            tooltip=["supplier", "decision", metric, "spend", "estimated_savings"],
        )
        .properties(height=420)
        .configure_axis(labelFontSize=11, titleFontSize=13)
        .configure_legend(labelFontSize=12, titleFontSize=13)
    )


def build_pareto_chart(
    df: pd.DataFrame,
    category_col: str,
    value_col: str,
    cum_share_col: str,
    abc_col: str,
    x_title: str,
    y_title: str,
    top_n: Optional[int] = None,
    color_col: Optional[str] = None,
    color_title: Optional[str] = None,
    color_domain: Optional[List[str]] = None,
    color_range: Optional[List[str]] = None,
) -> alt.Chart:
    chart_data = df.head(top_n).copy() if top_n is not None else df.copy()
    category_axis = alt.Axis(labelAngle=-45, labelLimit=220, labelOverlap=False)
    tooltip_fields = [category_col, value_col, cum_share_col, abc_col]
    for extra_field in ["sourcing_risk_level", "supply_risk_score", "single_source_flag", "supplier_count"]:
        if extra_field in chart_data.columns and extra_field not in tooltip_fields:
            tooltip_fields.append(extra_field)
    abc_label_map = {
        "A": "A - Highest value",
        "B": "B - Mid-tier",
        "C": "C - Lower value",
    }
    bar_color_col = color_col or abc_col
    bar_color_title = color_title or "ABC Class"
    bar_color_domain = color_domain or ABC_CLASS_DOMAIN
    bar_color_range = color_range or ABC_CLASS_RANGE
    if color_col is None and abc_col in chart_data.columns:
        labeled_col = f"{abc_col}_label"
        chart_data[labeled_col] = chart_data[abc_col].astype(str).map(abc_label_map).fillna(chart_data[abc_col].astype(str))
        bar_color_col = labeled_col
        bar_color_domain = [abc_label_map.get(value, value) for value in ABC_CLASS_DOMAIN]
    bars = (
        alt.Chart(chart_data)
        .mark_bar()
        .encode(
            x=alt.X(f"{category_col}:N", title=x_title, axis=category_axis),
            y=alt.Y(f"{value_col}:Q", title=y_title),
            color=alt.Color(
                f"{bar_color_col}:N",
                title=bar_color_title,
                scale=alt.Scale(domain=bar_color_domain, range=bar_color_range),
            ),
            tooltip=tooltip_fields,
        )
    )
    line = (
        alt.Chart(chart_data)
        .mark_line(color="#1f4e79", point=True)
        .encode(
            x=alt.X(f"{category_col}:N", axis=category_axis),
            y=alt.Y(f"{cum_share_col}:Q", title="Cumulative Share", axis=alt.Axis(format=".0%")),
        )
    )
    return (
        alt.layer(bars, line)
        .resolve_scale(y="independent")
        .properties(height=460)
        .configure_axis(labelFontSize=11, titleFontSize=13)
        .configure_legend(labelFontSize=12, titleFontSize=13)
    )


def build_component_analysis_bubble_chart(component_summary: pd.DataFrame, show_legend: bool = True) -> alt.Chart:
    chart_data = component_summary.copy()
    chart_data["single_source_highlight"] = np.where(
        chart_data["single_source_flag"].fillna(False),
        "Single Source",
        "Not Single Source",
    )
    points = (
        alt.Chart(chart_data)
        .mark_point(filled=True, opacity=0.82, stroke="white", strokeWidth=1)
        .encode(
            x=alt.X("supply_risk_score:Q", title="Supply Risk Score"),
            y=alt.Y("spend:Q", title="Spend"),
            size=alt.Size("strategic_priority_score:Q", title="Strategic Priority", scale=alt.Scale(range=[120, 1800])),
            color=alt.Color(
                "single_source_highlight:N",
                title="Source Type",
                scale=alt.Scale(domain=["Single Source", "Not Single Source"], range=["#dc2626", "#2563eb"]),
                legend=None if not show_legend else alt.Legend(),
            ),
            tooltip=[
                "component",
                "spend",
                "supplier_count",
                "sourcing_risk_level",
                "supply_risk_score",
                "strategic_priority_score",
                "kraljic_quadrant",
                "single_source_highlight",
                "largest_supplier_share",
                "defect_rate",
            ],
        )
    )
    labels = (
        alt.Chart(chart_data)
        .mark_text(fontSize=10, baseline="middle", align="left", dx=8, color="#111")
        .encode(
            x="supply_risk_score:Q",
            y="spend:Q",
            text="component:N",
        )
    )
    return (
        (points + labels)
        .properties(height=560)
        .configure_axis(labelFontSize=12, titleFontSize=13)
        .configure_legend(labelFontSize=12, titleFontSize=13)
    )


def build_component_supplier_detail_chart(
    component_supplier_detail: pd.DataFrame, top_n: Optional[int] = None, show_legend: bool = True
) -> alt.Chart:
    chart_data = component_supplier_detail.copy()
    component_order = (
        chart_data.groupby("component", as_index=False)["spend"].sum().sort_values("spend", ascending=False)["component"].tolist()
    )
    if top_n is not None:
        component_order = component_order[:top_n]
        chart_data = chart_data[chart_data["component"].isin(component_order)].copy()
    return (
        alt.Chart(chart_data)
        .mark_bar()
        .encode(
            x=alt.X("effective_supplier_share:Q", title="Effective Supplier Coverage Share", axis=alt.Axis(format=".0%")),
            y=alt.Y("component:N", sort=component_order, title="Component"),
            color=alt.Color("supplier:N", title="Supplier", legend=None if not show_legend else alt.Legend()),
            tooltip=[
                "component",
                "supplier",
                alt.Tooltip("spend:Q", title="Spend", format=",.0f"),
                alt.Tooltip("supplier_share:Q", title="Current Spend Share", format=".0%"),
                alt.Tooltip("effective_supplier_share:Q", title="Effective Coverage Share", format=".0%"),
                alt.Tooltip("supplier_count:Q", title="Effective Supplier Count"),
                alt.Tooltip("single_source_flag:N", title="Single Source"),
            ],
        )
        .properties(height=max(360, min(840, 24 * max(len(component_order), 1))))
        .configure_axis(labelFontSize=11, titleFontSize=13)
        .configure_legend(labelFontSize=12, titleFontSize=13)
    )


def build_supplier_concentration_chart(component_summary: pd.DataFrame, top_n: Optional[int] = 12) -> alt.Chart:
    sorted_data = component_summary.sort_values("largest_supplier_share", ascending=False).copy()
    chart_data = sorted_data.head(top_n).copy() if top_n is not None else sorted_data
    bars = (
        alt.Chart(chart_data)
        .mark_bar()
        .encode(
            x=alt.X("largest_supplier_share:Q", title="Largest Supplier Share"),
            y=alt.Y("component:N", sort="-x", title="Component"),
            color=alt.Color(
                "sourcing_risk_level:N",
                title="Sourcing Risk",
                scale=alt.Scale(domain=RISK_LEVEL_DOMAIN, range=RISK_LEVEL_RANGE),
            ),
            tooltip=["component", "dominant_supplier", "largest_supplier_share", "supplier_count", "sourcing_risk_level"],
        )
    )
    labels = (
        alt.Chart(chart_data)
        .mark_text(align="left", baseline="middle", dx=8, fontSize=11, color="#111")
        .encode(
            x="largest_supplier_share:Q",
            y=alt.Y("component:N", sort="-x"),
            text="dominant_supplier:N",
        )
    )
    return (bars + labels).properties(height=400)


def build_strategic_outcomes_chart(supplier_summary: pd.DataFrame, scenario_applied: bool = False) -> alt.Chart:
    chart_data = supplier_summary.sort_values("spend", ascending=False).copy()
    if scenario_applied:
        chart_data["outcome_display"] = chart_data["decision"]
        legend_title = "Recommendation"
        domain = ["Eliminate / De-prioritize", "Keep and Monitor", "Keep / Consolidate To"]
    else:
        chart_data["outcome_display"] = np.select(
            [
                chart_data["decision"].eq("Eliminate / De-prioritize"),
                chart_data["decision"].eq("Keep / Consolidate To"),
            ],
            ["Evaluate Exit Scenario", "Evaluate Consolidation Scenario"],
            default="Evaluate Monitoring Need",
        )
        legend_title = "Scenario Test"
        domain = ["Evaluate Exit Scenario", "Evaluate Monitoring Need", "Evaluate Consolidation Scenario"]
    color_scale = alt.Scale(
        domain=domain,
        range=["#d73027", "#f39c12", "#2e8b57"],
    )
    return (
        alt.Chart(chart_data)
        .mark_bar()
        .encode(
            x=alt.X("supplier:N", title="Supplier", axis=alt.Axis(labelAngle=0, labelLimit=220)),
            y=alt.Y("spend:Q", title="Spend"),
            color=alt.Color("outcome_display:N", title=legend_title, scale=color_scale),
            tooltip=["supplier", "outcome_display", "decision", "spend", "estimated_savings"],
        )
        .properties(height=420)
        .configure_axis(labelFontSize=11, titleFontSize=13)
        .configure_legend(labelFontSize=12, titleFontSize=13)
    )


def build_supplier_component_mix_chart(
    component_supplier_detail: pd.DataFrame, component_summary: pd.DataFrame, top_n_suppliers: Optional[int] = None
) -> alt.Chart:
    chart_data = component_supplier_detail.copy()
    missing_cols = [col for col in ["kraljic_quadrant", "single_source_flag", "sourcing_risk_level"] if col not in chart_data.columns]
    if missing_cols:
        chart_data = chart_data.merge(
            component_summary[["component"] + missing_cols].drop_duplicates(subset=["component"]),
            on="component",
            how="left",
        )
    supplier_order = (
        chart_data.groupby("supplier", as_index=False)["spend"].sum().sort_values("spend", ascending=False)["supplier"].tolist()
    )
    if top_n_suppliers is not None:
        supplier_order = supplier_order[:top_n_suppliers]
        chart_data = chart_data.loc[chart_data["supplier"].isin(supplier_order)].copy()
    supplier_totals = chart_data.groupby("supplier")["spend"].transform("sum")
    chart_data["component_share_within_supplier"] = np.where(
        supplier_totals > 0,
        chart_data["spend"] / supplier_totals,
        0.0,
    )
    return (
        alt.Chart(chart_data)
        .mark_bar()
        .encode(
            x=alt.X("supplier:N", sort=supplier_order, title="Supplier", axis=alt.Axis(labelAngle=0, labelLimit=220)),
            y=alt.Y("spend:Q", title="Spend"),
            color=alt.Color("component:N", title="Component"),
            tooltip=[
                "supplier",
                "component",
                alt.Tooltip("spend:Q", title="Spend", format=",.0f"),
                alt.Tooltip("component_share_within_supplier:Q", title="Share of Supplier Spend", format=".0%"),
                "kraljic_quadrant",
                "sourcing_risk_level",
                alt.Tooltip("single_source_flag:N", title="Single Source"),
            ],
        )
        .properties(height=430)
        .configure_axis(labelFontSize=11, titleFontSize=13)
        .configure_legend(labelFontSize=10, titleFontSize=12, orient="bottom", columns=4)
    )


def build_supplier_quadrant_mix_chart(
    component_supplier_detail: pd.DataFrame, component_summary: pd.DataFrame, top_n_suppliers: Optional[int] = None
) -> alt.Chart:
    chart_data = component_supplier_detail.copy()
    missing_cols = [col for col in ["kraljic_quadrant", "single_source_flag", "sourcing_risk_level"] if col not in chart_data.columns]
    if missing_cols:
        chart_data = chart_data.merge(
            component_summary[["component"] + missing_cols].drop_duplicates(subset=["component"]),
            on="component",
            how="left",
        )
    supplier_order = (
        chart_data.groupby("supplier", as_index=False)["spend"].sum().sort_values("spend", ascending=False)["supplier"].tolist()
    )
    if top_n_suppliers is not None:
        supplier_order = supplier_order[:top_n_suppliers]
        chart_data = chart_data.loc[chart_data["supplier"].isin(supplier_order)].copy()
    quadrant_order = ["Strategic", "Bottleneck", "Leverage", "Non-Critical"]
    return (
        alt.Chart(chart_data)
        .mark_bar()
        .encode(
            x=alt.X("supplier:N", sort=supplier_order, title="Supplier", axis=alt.Axis(labelAngle=0, labelLimit=220)),
            y=alt.Y("spend:Q", title="Spend"),
            color=alt.Color(
                "kraljic_quadrant:N",
                title="Kraljic Quadrant",
                scale=alt.Scale(
                    domain=quadrant_order,
                    range=["#b45309", "#dc2626", "#2563eb", "#65a30d"],
                ),
            ),
            tooltip=[
                "supplier",
                "kraljic_quadrant",
                alt.Tooltip("spend:Q", title="Spend", format=",.0f"),
                alt.Tooltip("single_source_flag:N", title="Contains Single-Source Component"),
            ],
        )
        .properties(height=430)
        .configure_axis(labelFontSize=11, titleFontSize=13)
        .configure_legend(labelFontSize=11, titleFontSize=12)
    )


def get_dynamic_single_source_candidates(
    analytics: Dict[str, pd.DataFrame], selected_suppliers: List[str]
) -> List[Dict[str, object]]:
    if not selected_suppliers:
        return []

    detail = analytics["component_supplier_detail"].copy()
    component_spend = analytics["component_summary"][["component", "spend"]].copy()
    filtered = detail.loc[detail["supplier"].isin(set(selected_suppliers))].copy()
    if filtered.empty:
        return []

    candidate_frame = (
        filtered.groupby("component", as_index=False)
        .agg(
            selected_supplier_count=("supplier", "nunique"),
            selected_suppliers=("supplier", lambda x: ", ".join(sorted(pd.unique(x)))),
        )
    )
    candidate_frame = candidate_frame.loc[candidate_frame["selected_supplier_count"].eq(1)].copy()
    if candidate_frame.empty:
        return []

    candidate_frame = candidate_frame.merge(component_spend, on="component", how="left")
    candidate_frame = candidate_frame.rename(columns={"selected_suppliers": "dominant_supplier"})
    candidate_frame = candidate_frame.sort_values("spend", ascending=False)
    return candidate_frame[["component", "dominant_supplier", "spend"]].to_dict(orient="records")


def get_uncovered_candidates(
    analytics: Dict[str, pd.DataFrame], selected_suppliers: List[str]
) -> List[Dict[str, object]]:
    if not selected_suppliers:
        return analytics["component_summary"][["component", "dominant_supplier", "spend"]].sort_values("spend", ascending=False).to_dict(orient="records")

    detail = analytics["component_supplier_detail"].copy()
    component_summary = analytics["component_summary"][["component", "dominant_supplier", "spend"]].copy()
    covered_components = set(detail.loc[detail["supplier"].isin(set(selected_suppliers)), "component"])
    uncovered = component_summary.loc[~component_summary["component"].isin(covered_components)].copy()
    uncovered = uncovered.sort_values("spend", ascending=False)
    return uncovered.to_dict(orient="records")


def estimate_mitigation_assignment_cost(
    analytics: Dict[str, pd.DataFrame], component_name: str, supplier_name: str, mitigation_kind: str
) -> float:
    component_summary = analytics["component_summary"]
    supplier_summary = analytics["supplier_summary"]

    component_match = component_summary.loc[component_summary["component"].eq(component_name)]
    supplier_match = supplier_summary.loc[supplier_summary["supplier"].eq(supplier_name)]
    component_spend = float(component_match.iloc[0]["spend"]) if not component_match.empty else 0.0

    if supplier_match.empty:
        supplier_risk = 50.0
        supplier_lead_time = 0.0
        performance_score = 50.0
    else:
        supplier_risk = float(supplier_match.iloc[0]["supplier_risk_score"])
        supplier_lead_time = float(supplier_match.iloc[0]["avg_lead_time"])
        performance_score = float(supplier_match.iloc[0]["performance_score"])

    if mitigation_kind == "uncovered":
        base_cost = max(4000.0, component_spend * 0.025)
    else:
        base_cost = max(2000.0, component_spend * 0.012)

    risk_multiplier = 1.0 + (supplier_risk / 100.0) * 0.20
    lead_time_multiplier = 1.0 + min(max(supplier_lead_time, 0.0), 90.0) / 300.0
    performance_multiplier = 1.0 + max(0.0, 70.0 - performance_score) / 250.0
    return round(base_cost * risk_multiplier * lead_time_multiplier * performance_multiplier, 2)


@st.cache_data(show_spinner=False)
def build_consolidation_scenario(
    analytics: Dict[str, pd.DataFrame], selected_suppliers: Tuple[str, ...], mitigation_assignments: Tuple[str, ...]
) -> Tuple[Dict[str, float], pd.DataFrame, List[str]]:
    component_summary = analytics["component_summary"].copy()
    detail = analytics["component_supplier_detail"].copy()
    selected_set = set(selected_suppliers)
    mitigation_map: Dict[str, List[str]] = {}
    for assignment in mitigation_assignments:
        if "|||" not in assignment:
            continue
        component_name, supplier_name = assignment.split("|||", 1)
        mitigation_map.setdefault(component_name, []).append(supplier_name)
    mitigation_supplier_set = {supplier for suppliers in mitigation_map.values() for supplier in suppliers}

    scenario_rows: List[Dict[str, object]] = []
    total_mitigation_cost = 0.0
    for component_row in component_summary.to_dict(orient="records"):
        component_name = component_row["component"]
        component_detail = detail.loc[detail["component"] == component_name].copy()
        selected_detail = component_detail.loc[component_detail["supplier"].isin(selected_set)].copy()

        selected_count = int(selected_detail["supplier"].nunique())
        covered = selected_count > 0
        component_spend = float(component_row["spend"])
        current_selected_spend = float(selected_detail["spend"].sum()) if covered else 0.0
        coverage_share = current_selected_spend / component_spend if component_spend > 0 else 0.0
        planned_mitigation = ""
        effective_supplier_count = selected_count
        base_effective_supplier_count = selected_count
        baseline_scenario_supply_risk = np.nan
        scenario_risk_reduction = 0.0
        component_mitigation_cost = 0.0

        if covered:
            total_selected_spend = float(selected_detail["spend"].sum())
            selected_detail["scenario_weight"] = (
                selected_detail["spend"] / total_selected_spend if total_selected_spend > 0 else 0.0
            )
            scenario_defect_rate = float((selected_detail["defect_rate"] * selected_detail["scenario_weight"]).sum())
            scenario_lead_time = float((selected_detail["avg_lead_time"] * selected_detail["scenario_weight"]).sum())
            scenario_external_risk = float((selected_detail["avg_risk_score"] * selected_detail["scenario_weight"]).sum())
            scenario_largest_share = float(selected_detail["scenario_weight"].max())
            base_scenario_largest_share = scenario_largest_share
            baseline_weighted_supplier_inverse = 1 / max(base_effective_supplier_count, 1)
            baseline_scenario_supply_risk = round(
                0.28 * float(scale_if_variable(pd.Series([baseline_weighted_supplier_inverse]), default_value=100.0).iloc[0])
                + 0.22 * float(scale_if_variable(pd.Series([base_scenario_largest_share]), default_value=base_scenario_largest_share * 100).iloc[0])
                + 0.20 * float(scale_if_variable(pd.Series([scenario_defect_rate]), default_value=scenario_defect_rate * 100).iloc[0])
                + 0.12 * float(scale_if_variable(pd.Series([scenario_lead_time]), default_value=scenario_lead_time * 10).iloc[0])
                + 0.08 * float(scale_if_variable(pd.Series([scenario_external_risk]), default_value=scenario_external_risk).iloc[0])
                + 0.10 * float(scale_if_variable(pd.Series([component_row["criticality"]]), default_value=float(component_row["criticality"])).iloc[0]),
                2,
            )
            if bool(component_row["single_source_flag"]) and component_name in mitigation_map:
                incumbent_suppliers = set(selected_detail["supplier"].tolist())
                incremental_mitigation = sorted(set(mitigation_map[component_name]) - incumbent_suppliers)
                if incremental_mitigation:
                    planned_mitigation = ", ".join(incremental_mitigation)
                    effective_supplier_count = selected_count + len(incremental_mitigation)
                    scenario_largest_share = 1 / effective_supplier_count
                    component_mitigation_cost = sum(
                        estimate_mitigation_assignment_cost(analytics, component_name, supplier_name, "single_source")
                        for supplier_name in incremental_mitigation
                    )
                    total_mitigation_cost += component_mitigation_cost

            weighted_supplier_inverse = 1 / max(effective_supplier_count, 1)
            scenario_supply_risk = round(
                0.28 * float(scale_if_variable(pd.Series([weighted_supplier_inverse]), default_value=100.0).iloc[0])
                + 0.22 * float(scale_if_variable(pd.Series([scenario_largest_share]), default_value=scenario_largest_share * 100).iloc[0])
                + 0.20 * float(scale_if_variable(pd.Series([scenario_defect_rate]), default_value=scenario_defect_rate * 100).iloc[0])
                + 0.12 * float(scale_if_variable(pd.Series([scenario_lead_time]), default_value=scenario_lead_time * 10).iloc[0])
                + 0.08 * float(scale_if_variable(pd.Series([scenario_external_risk]), default_value=scenario_external_risk).iloc[0])
                + 0.10 * float(scale_if_variable(pd.Series([component_row["criticality"]]), default_value=float(component_row["criticality"])).iloc[0]),
                2,
            )
            scenario_risk_reduction = max(0.0, round(float(baseline_scenario_supply_risk) - float(scenario_supply_risk), 2))
            scenario_risk_level = "High" if effective_supplier_count <= 1 else "Medium" if effective_supplier_count <= 3 else "Low"
            status = "Covered"
            selected_supplier_names = ", ".join(sorted(selected_detail["supplier"].unique().tolist()))
        else:
            scenario_defect_rate = np.nan
            scenario_lead_time = np.nan
            scenario_external_risk = np.nan
            scenario_largest_share = np.nan
            scenario_supply_risk = np.nan
            assigned_pickups = sorted(set(mitigation_map.get(component_name, [])))
            if assigned_pickups:
                planned_mitigation = ", ".join(assigned_pickups)
                effective_supplier_count = len(assigned_pickups)
                scenario_supply_risk = np.nan
                component_mitigation_cost = sum(
                    estimate_mitigation_assignment_cost(analytics, component_name, supplier_name, "uncovered")
                    for supplier_name in assigned_pickups
                )
                total_mitigation_cost += component_mitigation_cost
                scenario_risk_level = "High" if effective_supplier_count <= 1 else "Medium" if effective_supplier_count <= 3 else "Low"
                status = "Mitigated Coverage"
                selected_supplier_names = ""
            else:
                scenario_risk_level = "Uncovered"
                status = "Not Covered"
                selected_supplier_names = ""

        scenario_rows.append(
            {
                "Component": component_name,
                "Current Spend": component_spend,
                "Selected Suppliers": selected_supplier_names,
                "Selected Supplier Count": selected_count,
                "Effective Supplier Count": effective_supplier_count,
                "Single-Source Mitigation Suppliers": planned_mitigation,
                "Base Scenario Supply Risk Score": baseline_scenario_supply_risk,
                "Scenario Risk Reduction": scenario_risk_reduction,
                "Mitigation Cost": component_mitigation_cost,
                "Current Coverage Share": coverage_share,
                "Scenario Status": status,
                "Scenario Risk Level": scenario_risk_level,
                "Scenario Supply Risk Score": scenario_supply_risk,
                "Scenario Defect Rate": scenario_defect_rate,
                "Scenario Avg Lead Time": scenario_lead_time,
                "Kraljic Quadrant": component_row["kraljic_quadrant"],
            }
        )

    scenario_df = pd.DataFrame(scenario_rows).sort_values(
        ["Scenario Status", "Current Spend"], ascending=[True, False]
    ).reset_index(drop=True)

    covered_df = scenario_df.loc[scenario_df["Scenario Status"].isin(["Covered", "Mitigated Coverage"])].copy()
    uncovered_df = scenario_df.loc[scenario_df["Scenario Status"].eq("Not Covered")].copy()
    base_single_source_components = set(
        component_summary.loc[component_summary["single_source_flag"], "component"].dropna().astype(str).tolist()
    )
    structural_single_source_df = covered_df.loc[
        covered_df["Component"].isin(base_single_source_components) & covered_df["Effective Supplier Count"].eq(1)
    ].copy() if not covered_df.empty else pd.DataFrame()
    scenario_created_single_source_df = covered_df.loc[
        covered_df["Effective Supplier Count"].eq(1) & covered_df["Selected Supplier Count"].gt(1)
    ].copy() if not covered_df.empty else pd.DataFrame()
    supplier_summary = analytics["supplier_summary"]
    excluded_supplier_summary = supplier_summary.loc[~supplier_summary["supplier"].isin(selected_set)].copy()
    estimated_savings = (
        float(excluded_supplier_summary["estimated_savings"].sum())
        if (selected_set and not excluded_supplier_summary.empty)
        else 0.0
    )
    net_savings = estimated_savings - total_mitigation_cost
    aggregate_risk_reduction = float(
        scenario_df["Scenario Risk Reduction"].fillna(0.0).sum()
    ) if not scenario_df.empty else 0.0

    metrics = {
        "selected_supplier_count": len(selected_set),
        "mitigation_supplier_count": len(mitigation_supplier_set),
        "covered_spend": float(covered_df["Current Spend"].sum()) if not covered_df.empty else 0.0,
        "covered_spend_share": float(covered_df["Current Spend"].sum() / component_summary["spend"].sum()) if len(component_summary) else 0.0,
        "fully_covered_components": int(covered_df["Current Coverage Share"].fillna(0).ge(1.0).sum()) if not covered_df.empty else 0,
        "single_source_components": int(covered_df["Effective Supplier Count"].eq(1).sum()) if not covered_df.empty else 0,
        "structural_single_source_components": int(len(structural_single_source_df)),
        "scenario_created_single_source_components": int(len(scenario_created_single_source_df)),
        "mitigated_single_source_components": int(
            covered_df["Single-Source Mitigation Suppliers"].ne("").sum()
        ) if not covered_df.empty else 0,
        "mitigated_uncovered_components": int(
            scenario_df["Scenario Status"].eq("Mitigated Coverage").sum()
        ) if not scenario_df.empty else 0,
        "medium_risk_components": int(covered_df["Effective Supplier Count"].between(2, 3).sum()) if not covered_df.empty else 0,
        "uncovered_components": int(len(uncovered_df)),
        "estimated_savings": estimated_savings,
        "mitigation_cost": total_mitigation_cost,
        "net_savings": net_savings,
        "aggregate_risk_reduction": aggregate_risk_reduction,
    }

    assumptions = [
        "You can choose any supplier for the general supplier scenario; the base-model keep recommendations remain visible elsewhere for comparison.",
        "A component is considered coverable only if at least one selected supplier already supplies it in the current data.",
        "Spend is held constant; covered spend is assumed transferable among selected incumbent suppliers, while uncovered components are flagged rather than auto-reassigned.",
        "Scenario savings are measured against today's full supplier set, so keeping all current suppliers returns zero modeled savings.",
        "Single-source mitigation assignments are component-specific planned future alternates, which lower structural risk only for the components you explicitly assign.",
        "Uncovered-component mitigation assignments are treated as planned supplier pickups for components not currently covered by the selected consolidation set.",
        "Scenario risk rises when a component would be left with one effective supplier and moderates when two or three effective suppliers remain.",
        "Mitigation cost is a modeled qualification and onboarding cost for each component-supplier backup assignment, adjusted for component spend, supplier risk, lead time, and performance.",
        "Net savings subtract modeled mitigation cost from gross consolidation savings, so adding more backups can improve resilience while reducing near-term savings.",
    ]
    return metrics, scenario_df, assumptions


def build_auto_mitigation_assignments(
    analytics: Dict[str, pd.DataFrame], selected_suppliers: Tuple[str, ...]
) -> Tuple[str, ...]:
    detail = analytics["component_supplier_detail"].copy()
    component_summary = analytics["component_summary"].copy()
    supplier_summary = analytics["supplier_summary"].copy()
    selected_set = set(selected_suppliers)
    assignments: List[str] = []
    single_source_suppliers = component_summary.loc[
        component_summary["single_source_flag"], ["component", "dominant_supplier"]
    ].set_index("component")["dominant_supplier"].to_dict()

    def best_supplier_for_component(component_name: str) -> Optional[str]:
        blocked_supplier = single_source_suppliers.get(component_name)
        if component_name in single_source_suppliers:
            options = supplier_summary.loc[~supplier_summary["supplier"].eq(blocked_supplier)].copy()
            if options.empty:
                return None
            options = options.sort_values(
                ["supplier_risk_score", "performance_score", "avg_lead_time", "spend"],
                ascending=[True, False, True, False],
            )
            return str(options.iloc[0]["supplier"])

        options = detail.loc[
            detail["component"].eq(component_name) & ~detail["supplier"].eq(blocked_supplier)
        ].copy()
        if options.empty:
            fallback = supplier_summary.loc[~supplier_summary["supplier"].eq(blocked_supplier)].copy()
            if fallback.empty:
                return None
            fallback = fallback.sort_values(
                ["supplier_risk_score", "performance_score", "avg_lead_time", "spend"],
                ascending=[True, False, True, False],
            )
            return str(fallback.iloc[0]["supplier"])
        options = options.sort_values(
            ["defect_rate", "avg_lead_time", "avg_risk_score", "spend"],
            ascending=[True, True, True, False],
        )
        return str(options.iloc[0]["supplier"])

    for row in get_dynamic_single_source_candidates(analytics, list(selected_suppliers)):
        component_name = str(row["component"])
        supplier_name = best_supplier_for_component(component_name)
        if supplier_name:
            assignments.append(f"{component_name}|||{supplier_name}")

    for row in get_uncovered_candidates(analytics, list(selected_suppliers)):
        component_name = str(row["component"])
        supplier_name = best_supplier_for_component(component_name)
        if supplier_name:
            assignments.append(f"{component_name}|||{supplier_name}")

    return tuple(sorted(set(assignments)))


def get_required_single_source_suppliers(analytics: Dict[str, pd.DataFrame]) -> Tuple[str, ...]:
    component_summary = analytics["component_summary"].copy()
    required_suppliers = (
        component_summary.loc[component_summary["single_source_flag"], "dominant_supplier"]
        .dropna()
        .astype(str)
        .unique()
        .tolist()
    )
    return tuple(sorted(required_suppliers))


def score_supplier_scenario(
    analytics: Dict[str, pd.DataFrame],
    selected_suppliers: Tuple[str, ...],
    mitigation_assignments: Tuple[str, ...],
    metrics: Dict[str, float],
    scenario_df: pd.DataFrame,
    optimization_objective: str = "Best Overall",
) -> Dict[str, object]:
    supplier_summary = analytics["supplier_summary"]
    selected_frame = supplier_summary.loc[supplier_summary["supplier"].isin(set(selected_suppliers))].copy()
    avg_supplier_risk = float(selected_frame["supplier_risk_score"].mean()) if not selected_frame.empty else 100.0
    avg_performance = float(selected_frame["performance_score"].mean()) if not selected_frame.empty else 0.0
    total_spend = float(analytics["component_summary"]["spend"].sum()) if len(analytics["component_summary"]) else 1.0

    high_risk_components = int(scenario_df["Scenario Risk Level"].eq("High").sum()) if not scenario_df.empty else 0
    medium_risk_components = int(scenario_df["Scenario Risk Level"].eq("Medium").sum()) if not scenario_df.empty else 0
    low_risk_components = int(scenario_df["Scenario Risk Level"].eq("Low").sum()) if not scenario_df.empty else 0
    mitigation_count = len({assignment.split("|||", 1)[1] for assignment in mitigation_assignments if "|||" in assignment})
    mitigation_assignment_count = len([assignment for assignment in mitigation_assignments if "|||" in assignment])
    aggregate_risk_reduction = float(metrics.get("aggregate_risk_reduction", 0.0))
    mitigation_cost = float(metrics.get("mitigation_cost", 0.0))
    net_savings = float(metrics.get("net_savings", 0.0))

    impacts = [
        ("Covered spend", float(metrics.get("covered_spend_share", 0.0)) * 1000.0, f"{float(metrics.get('covered_spend_share', 0.0)):.0%}"),
        ("Uncovered components", -float(metrics.get("uncovered_components", 0)) * 350.0, int(metrics.get("uncovered_components", 0))),
        ("High-risk components", -high_risk_components * 140.0, high_risk_components),
        ("Single-source components", -float(metrics.get("single_source_components", 0)) * 120.0, int(metrics.get("single_source_components", 0))),
        ("Medium-risk components", -medium_risk_components * 40.0, medium_risk_components),
        ("Mitigated uncovered", float(metrics.get("mitigated_uncovered_components", 0)) * 80.0, int(metrics.get("mitigated_uncovered_components", 0))),
        ("Mitigated single-source", float(metrics.get("mitigated_single_source_components", 0)) * 70.0, int(metrics.get("mitigated_single_source_components", 0))),
        ("Aggregate risk reduction", (aggregate_risk_reduction / max(len(scenario_df), 1)) * 9.0, round(aggregate_risk_reduction, 1)),
        ("Mitigation assignments", min(mitigation_assignment_count, 6) * 8.0, mitigation_assignment_count),
        ("Net savings", (net_savings / max(total_spend, 1.0)) * 240.0, round(net_savings, 0)),
        ("Mitigation cost", -(mitigation_cost / max(total_spend, 1.0)) * 75.0, round(mitigation_cost, 0)),
        ("Low-risk components", low_risk_components * 10.0, low_risk_components),
        ("Average supplier performance", (avg_performance / 100.0) * 70.0, round(avg_performance, 1)),
        ("Average supplier risk", -(avg_supplier_risk / 100.0) * 65.0, round(avg_supplier_risk, 1)),
        ("Distinct mitigation suppliers", -mitigation_count * 4.0, mitigation_count),
    ]
    objective_bonus = 0.0
    if optimization_objective == "Best Net Savings":
        objective_bonus = (
            (net_savings / max(total_spend, 1.0)) * 820.0
            - high_risk_components * 35.0
            - float(metrics.get("single_source_components", 0)) * 35.0
            - float(metrics.get("uncovered_components", 0)) * 120.0
        )
    elif optimization_objective == "Best Risk Reduction":
        objective_bonus = (
            aggregate_risk_reduction * 5.5
            - high_risk_components * 110.0
            - float(metrics.get("single_source_components", 0)) * 90.0
            - medium_risk_components * 25.0
            - float(metrics.get("uncovered_components", 0)) * 180.0
        )
    impacts.append((f"Objective emphasis ({optimization_objective})", objective_bonus, optimization_objective))
    score = float(sum(item[1] for item in impacts))
    breakdown = pd.DataFrame(impacts, columns=["Factor", "Impact", "Value"])
    breakdown["Impact"] = breakdown["Impact"].round(2)
    breakdown["Value"] = breakdown["Value"].astype(str)
    return {
        "score": round(score, 2),
        "high_risk_components": high_risk_components,
        "medium_risk_components": medium_risk_components,
        "low_risk_components": low_risk_components,
        "avg_supplier_risk": round(avg_supplier_risk, 2),
        "avg_performance": round(avg_performance, 2),
        "optimization_objective": optimization_objective,
        "breakdown": breakdown,
    }


@st.cache_data(show_spinner=False)
def recommend_best_supplier_scenario(
    analytics: Dict[str, pd.DataFrame],
    optimization_objective: str = "Best Overall",
) -> Dict[str, object]:
    supplier_summary = analytics["supplier_summary"].copy()
    all_suppliers = supplier_summary["supplier"].dropna().astype(str).tolist()
    if not all_suppliers:
        return {
            "selected_suppliers": tuple(),
            "mitigation_assignments": tuple(),
            "metrics": {},
            "scenario_df": pd.DataFrame(),
            "scorecard": {},
            "tested_scenarios": 0,
            "rationale": "No suppliers were available to evaluate.",
        }

    ranked_suppliers = supplier_summary.sort_values(
        ["decision_rank", "performance_score", "supplier_risk_score", "spend"],
        ascending=[True, False, True, False],
    )["supplier"].tolist()
    required_suppliers = get_required_single_source_suppliers(analytics)
    required_set = set(required_suppliers)
    optional_suppliers = [supplier for supplier in ranked_suppliers if supplier not in required_set]

    best_result: Optional[Dict[str, object]] = None
    tested_scenarios = 0
    min_selected_count = max(2 if len(all_suppliers) >= 2 else 1, len(required_suppliers))
    max_selected_count = len(all_suppliers)
    max_pool_size = min(len(ranked_suppliers), 8)
    for selected_count in range(min_selected_count, max_selected_count + 1):
        extra_slots = selected_count - len(required_suppliers)
        if extra_slots < 0:
            continue
        optional_pool_size = min(len(optional_suppliers), max(max_pool_size - len(required_suppliers), extra_slots))
        optional_pool_size = min(optional_pool_size, max(extra_slots + 3, 4))
        while optional_pool_size > extra_slots and math.comb(optional_pool_size, extra_slots) > 120:
            optional_pool_size -= 1
        candidate_optional_pool = optional_suppliers[:optional_pool_size]

        for combo in itertools.combinations(candidate_optional_pool, extra_slots):
            selected_suppliers = tuple(sorted(required_set.union(combo)))
            mitigation_assignments = build_auto_mitigation_assignments(analytics, selected_suppliers)
            metrics, scenario_df, _ = build_consolidation_scenario(analytics, selected_suppliers, mitigation_assignments)
            scorecard = score_supplier_scenario(
                analytics,
                selected_suppliers,
                mitigation_assignments,
                metrics,
                scenario_df,
                optimization_objective=optimization_objective,
            )
            tested_scenarios += 1

            result = {
                "selected_suppliers": selected_suppliers,
                "mitigation_assignments": mitigation_assignments,
                "metrics": metrics,
                "scenario_df": scenario_df,
                "scorecard": scorecard,
            }
            if best_result is None or scorecard["score"] > best_result["scorecard"]["score"]:
                best_result = result

    if best_result is None:
        return {
            "selected_suppliers": tuple(),
            "mitigation_assignments": tuple(),
            "metrics": {},
            "scenario_df": pd.DataFrame(),
            "scorecard": {},
            "tested_scenarios": 0,
            "rationale": "No viable scenarios were identified from the candidate pool.",
        }

    metrics = best_result["metrics"]
    scorecard = best_result["scorecard"]
    rationale = (
        f"Recommended under the `{optimization_objective}` objective with {len(best_result.get('selected_suppliers', ()))} selected suppliers because it delivers the strongest scenario score of {float(scorecard.get('score', 0.0)):,.1f}, "
        f"covers {metrics.get('covered_spend_share', 0.0):.0%} of spend, "
        f"leaves {int(metrics.get('uncovered_components', 0))} uncovered components, "
        f"results in {scorecard.get('high_risk_components', 0)} high-risk components, "
        f"and lands at net savings of ${float(metrics.get('net_savings', 0.0)):,.0f} after mitigation cost."
    )
    best_result["tested_scenarios"] = tested_scenarios
    best_result["rationale"] = rationale
    return best_result


@st.cache_data(show_spinner=False)
def build_applied_scenario_analytics(
    analytics: Dict[str, pd.DataFrame], selected_suppliers: Tuple[str, ...], mitigation_assignments: Tuple[str, ...]
) -> Tuple[Dict[str, pd.DataFrame], Dict[str, float], pd.DataFrame, List[str]]:
    base_component_summary = analytics["component_summary"].copy()
    base_detail = analytics["component_supplier_detail"].copy()
    base_supplier_summary = analytics["supplier_summary"].copy()
    scenario_metrics, scenario_df, scenario_assumptions = build_consolidation_scenario(
        analytics, selected_suppliers, mitigation_assignments
    )
    base_total_spend = float(base_component_summary["spend"].sum()) if len(base_component_summary) else 0.0
    scenario_total_spend = max(0.0, base_total_spend - float(scenario_metrics["estimated_savings"]))
    spend_scale = (scenario_total_spend / base_total_spend) if base_total_spend > 0 else 1.0

    scenario_component_summary = base_component_summary.merge(
        scenario_df[
            [
                "Component",
                "Selected Supplier Count",
                "Effective Supplier Count",
                "Scenario Status",
                "Scenario Risk Level",
                "Scenario Supply Risk Score",
                "Single-Source Mitigation Suppliers",
            ]
        ],
        left_on="component",
        right_on="Component",
        how="left",
    )
    scenario_component_summary["spend"] = (scenario_component_summary["spend"] * spend_scale).round(2)
    scenario_component_summary["supplier_count"] = scenario_component_summary["Effective Supplier Count"].fillna(0).astype(int)
    scenario_component_summary["single_source_flag"] = scenario_component_summary["supplier_count"].eq(1)
    scenario_component_summary["backup_supplier_flag"] = scenario_component_summary["supplier_count"].gt(1)
    scenario_component_summary["sourcing_risk_level"] = scenario_component_summary["Scenario Risk Level"].fillna("Uncovered")
    scenario_component_summary["high_risk_flag"] = scenario_component_summary["sourcing_risk_level"].eq("High")
    scenario_component_summary["supply_risk_score"] = scenario_component_summary["Scenario Supply Risk Score"].fillna(
        scenario_component_summary["supply_risk_score"]
    )
    scenario_component_summary["largest_supplier_share"] = np.where(
        scenario_component_summary["supplier_count"] > 0,
        1 / scenario_component_summary["supplier_count"].replace(0, np.nan),
        0.0,
    )
    scenario_component_summary["supplier_concentration"] = scenario_component_summary["largest_supplier_share"].fillna(0.0)
    leverage_spend = min_max_scale(scenario_component_summary["spend"])
    leverage_supplier_count = min_max_scale(scenario_component_summary["supplier_count"])
    leverage_inverse_concentration = min_max_scale(1 - scenario_component_summary["largest_supplier_share"].fillna(1.0))
    scenario_component_summary["negotiation_leverage_score"] = (
        0.45 * leverage_spend
        + 0.25 * leverage_supplier_count
        + 0.30 * leverage_inverse_concentration
    ).round(2)
    strategic_risk = min_max_scale(scenario_component_summary["supply_risk_score"])
    strategic_impact = min_max_scale(scenario_component_summary["profit_impact_score"])
    strategic_single = scenario_component_summary["single_source_flag"].astype(float) * 100.0
    strategic_leverage = min_max_scale(scenario_component_summary["negotiation_leverage_score"])
    scenario_component_summary["strategic_priority_score"] = (
        0.35 * strategic_impact
        + 0.25 * strategic_risk
        + 0.20 * strategic_single
        + 0.20 * strategic_leverage
    ).round(2)
    scenario_component_summary["risk_adjusted_spend"] = (
        scenario_component_summary["spend"]
        * (1 + scenario_component_summary["defect_rate"] * 4)
        * (1 + scenario_component_summary["single_source_flag"].astype(float) * 0.35)
        * (1 + scenario_component_summary["supply_risk_score"] / 100)
    ).round(2)

    spend_pareto = assign_abc_categories(
        scenario_component_summary.sort_values("spend", ascending=False).reset_index(drop=True), "spend", "spend"
    )
    risk_pareto = assign_abc_categories(
        scenario_component_summary.sort_values("risk_adjusted_spend", ascending=False).reset_index(drop=True),
        "risk_adjusted_spend",
        "risk",
    )
    strategic_pareto = assign_abc_categories(
        scenario_component_summary.sort_values("strategic_priority_score", ascending=False).reset_index(drop=True),
        "strategic_priority_score",
        "strategic",
    )
    scenario_component_summary = scenario_component_summary.drop(
        columns=[col for col in ["Component", "spend_cum_share", "spend_abc", "risk_cum_share", "risk_abc", "strategic_cum_share", "strategic_abc"] if col in scenario_component_summary.columns]
    )
    scenario_component_summary = scenario_component_summary.merge(
        spend_pareto[["component", "spend_cum_share", "spend_abc"]], on="component", how="left"
    )
    scenario_component_summary = scenario_component_summary.merge(
        risk_pareto[["component", "risk_cum_share", "risk_abc"]], on="component", how="left"
    )
    scenario_component_summary = scenario_component_summary.merge(
        strategic_pareto[["component", "strategic_cum_share", "strategic_abc"]], on="component", how="left"
    )

    selected_set = set(selected_suppliers)
    scenario_detail = base_detail.loc[base_detail["supplier"].isin(selected_set)].copy()
    scenario_detail["spend"] = (scenario_detail["spend"] * spend_scale).round(2)
    mitigation_rows: List[Dict[str, object]] = []
    mitigation_map: Dict[str, List[str]] = {}
    for assignment in mitigation_assignments:
        if "|||" not in assignment:
            continue
        component_name, supplier_name = assignment.split("|||", 1)
        mitigation_map.setdefault(component_name, []).append(supplier_name)
    supplier_lookup = base_supplier_summary.set_index("supplier").to_dict(orient="index")
    component_lookup = scenario_component_summary.set_index("component").to_dict(orient="index")
    existing_pairs = set(zip(scenario_detail["component"], scenario_detail["supplier"]))
    for component_name, suppliers in mitigation_map.items():
        for supplier_name in suppliers:
            if (component_name, supplier_name) in existing_pairs:
                continue
            supplier_info = supplier_lookup.get(supplier_name, {})
            component_info = component_lookup.get(component_name, {})
            mitigation_rows.append(
                {
                    "component": component_name,
                    "supplier": supplier_name,
                    "spend": 0.0,
                    "units": 0.0,
                    "defects": 0.0,
                    "avg_lead_time": float(supplier_info.get("avg_lead_time", 0.0)),
                    "avg_risk_score": float(supplier_info.get("avg_risk_score", 0.0)),
                    "avg_criticality": float(component_info.get("criticality", 50.0)),
                    "defect_rate": 0.0,
                    "component_total_spend": float(component_info.get("spend", 0.0)),
                    "supplier_share": 0.0,
                    "supplier_count": int(component_info.get("supplier_count", 0)),
                    "single_source_flag": bool(component_info.get("single_source_flag", False)),
                    "backup_supplier_flag": bool(component_info.get("backup_supplier_flag", False)),
                    "concentration_gap": 0.0,
                }
            )
    if mitigation_rows:
        scenario_detail = pd.concat([scenario_detail, pd.DataFrame(mitigation_rows)], ignore_index=True, sort=False)

    detail_total_spend = float(scenario_detail["spend"].sum()) if len(scenario_detail) else 0.0
    if detail_total_spend > 0 and scenario_total_spend > 0:
        detail_scale = scenario_total_spend / detail_total_spend
        scenario_detail["spend"] = (scenario_detail["spend"] * detail_scale).round(2)

    scenario_detail = scenario_detail.drop(
        columns=[
            col
            for col in [
                "high_risk_flag",
                "sourcing_risk_level",
                "supply_risk_score",
                "profit_impact_score",
                "dominant_supplier",
                "kraljic_quadrant",
                "strategic_priority_score",
                "risk_adjusted_spend",
                "decision",
                "decision_reason",
                "estimated_savings",
                "supplier_count",
                "single_source_flag",
                "backup_supplier_flag",
            ]
            if col in scenario_detail.columns
        ]
    )
    scenario_detail = scenario_detail.merge(
        scenario_component_summary[
            [
                "component",
                "high_risk_flag",
                "sourcing_risk_level",
                "supply_risk_score",
                "profit_impact_score",
                "dominant_supplier",
                "kraljic_quadrant",
                "strategic_priority_score",
                "risk_adjusted_spend",
                "supplier_count",
                "single_source_flag",
                "backup_supplier_flag",
            ]
        ],
        on="component",
        how="left",
        suffixes=("", "_scenario"),
    )
    scenario_detail["component_total_spend"] = scenario_detail["component"].map(
        scenario_component_summary.set_index("component")["spend"].to_dict()
    ).fillna(0.0)
    detail_component_totals = scenario_detail.groupby("component")["spend"].transform("sum")
    scenario_detail["supplier_share"] = safe_divide(
        scenario_detail["spend"], detail_component_totals.replace(0, np.nan), fill_value=0.0
    )
    scenario_detail["effective_supplier_share"] = np.where(
        scenario_detail["supplier_count"].fillna(0).gt(0),
        1 / scenario_detail["supplier_count"].replace(0, np.nan),
        0.0,
    )
    scenario_supplier_summary = (
        scenario_detail.groupby("supplier", as_index=False)
        .agg(
            spend=("spend", "sum"),
            units=("units", "sum"),
            defects=("defects", "sum"),
            avg_lead_time=("avg_lead_time", "mean"),
            avg_risk_score=("avg_risk_score", "mean"),
            component_count=("component", "nunique"),
        )
    )
    scenario_supplier_summary["defect_rate"] = safe_divide(
        scenario_supplier_summary["defects"], scenario_supplier_summary["units"], fill_value=0.0
    )
    scenario_supplier_summary["portfolio_share"] = safe_divide(
        scenario_supplier_summary["spend"], scenario_supplier_summary["spend"].sum(), fill_value=0.0
    )

    scenario_component_flags = scenario_component_summary[
        ["component", "single_source_flag", "high_risk_flag", "strategic_priority_score", "supply_risk_score"]
    ].copy()
    scenario_supplier_component_map = scenario_detail[
        ["component", "supplier", "supplier_share"]
    ].merge(scenario_component_flags, on="component", how="left")
    scenario_supplier_protection = (
        scenario_supplier_component_map.groupby("supplier", as_index=False)
        .agg(
            supports_single_source=("single_source_flag", "max"),
            supports_high_risk=("high_risk_flag", "max"),
            protected_component_count=("component", lambda x: int(x.nunique())),
            strategic_priority_supported=("strategic_priority_score", "sum"),
            avg_component_risk=("supply_risk_score", "mean"),
            overlap_components=("component", "nunique"),
            max_component_share=("supplier_share", "max"),
        )
    )
    scenario_supplier_summary = scenario_supplier_summary.merge(
        scenario_supplier_protection, on="supplier", how="left"
    )
    for col in ["supports_single_source", "supports_high_risk"]:
        scenario_supplier_summary[col] = scenario_supplier_summary[col].fillna(False).astype(bool)
    for col in ["protected_component_count", "strategic_priority_supported", "avg_component_risk", "overlap_components", "max_component_share"]:
        scenario_supplier_summary[col] = scenario_supplier_summary[col].fillna(0.0)

    scenario_supplier_summary["supplier_risk_score"] = (
        0.30 * min_max_scale(scenario_supplier_summary["defect_rate"])
        + 0.20 * min_max_scale(scenario_supplier_summary["avg_lead_time"])
        + 0.20 * min_max_scale(scenario_supplier_summary["avg_component_risk"])
        + 0.15 * min_max_scale(scenario_supplier_summary["portfolio_share"])
        + 0.10 * scenario_supplier_summary["supports_single_source"].astype(float) * 100.0
        + 0.05 * scenario_supplier_summary["supports_high_risk"].astype(float) * 100.0
    ).round(2)
    scenario_supplier_summary["replaceability_score"] = (
        0.35 * min_max_scale(1 - scenario_supplier_summary["max_component_share"])
        + 0.25 * min_max_scale(1 - scenario_supplier_summary["portfolio_share"])
        + 0.20 * min_max_scale(
            1 - scenario_supplier_summary["component_count"] / max(1, scenario_supplier_summary["component_count"].max())
        )
        + 0.20 * (100 - min_max_scale(scenario_supplier_summary["strategic_priority_supported"]))
    ).round(2)
    scenario_supplier_summary["performance_score"] = (
        0.40 * min_max_scale(scenario_supplier_summary["defect_rate"], invert=True)
        + 0.25 * min_max_scale(scenario_supplier_summary["avg_lead_time"], invert=True)
        + 0.20 * min_max_scale(scenario_supplier_summary["avg_risk_score"], invert=True)
        + 0.15 * min_max_scale(scenario_supplier_summary["component_count"])
    ).round(2)

    scenario_supplier_decisions = classify_suppliers(
        scenario_detail,
        scenario_component_summary,
        scenario_supplier_summary,
    )
    scenario_supplier_summary = scenario_supplier_summary.merge(
        scenario_supplier_decisions, on="supplier", how="left"
    )
    scenario_supplier_summary["estimated_savings"] = 0.0
    scenario_supplier_summary["scenario_role"] = np.where(
        scenario_supplier_summary["supplier"].isin(selected_set), "Selected Supplier", "Mitigation Supplier"
    )

    scenario_detail = scenario_detail.merge(
        scenario_supplier_summary[["supplier", "decision", "decision_reason", "estimated_savings"]],
        on="supplier",
        how="left",
    )

    protected_components = scenario_component_summary.loc[
        scenario_component_summary["single_source_flag"] | scenario_component_summary["high_risk_flag"],
        ["component", "single_source_flag", "high_risk_flag", "strategic_priority_score", "supply_risk_score"],
    ].copy()

    scenario_analytics = {
        "supplier_summary": scenario_supplier_summary.sort_values(["decision_rank", "spend"], ascending=[True, False]).reset_index(drop=True),
        "component_summary": scenario_component_summary.sort_values("strategic_priority_score", ascending=False).reset_index(drop=True),
        "component_supplier_detail": scenario_detail.sort_values(["strategic_priority_score", "spend"], ascending=[False, False]).reset_index(drop=True),
        "protected_components": protected_components.sort_values("strategic_priority_score", ascending=False).reset_index(drop=True),
        "spend_pareto": spend_pareto,
        "risk_pareto": risk_pareto,
        "strategic_pareto": strategic_pareto,
    }
    return scenario_analytics, scenario_metrics, scenario_df, scenario_assumptions


@st.cache_data(show_spinner=False)
def build_applied_executive_summary(
    base_analytics: Dict[str, pd.DataFrame],
    scenario_analytics: Dict[str, pd.DataFrame],
    scenario_metrics: Dict[str, float],
) -> str:
    base_summary_text, _, _ = build_executive_summary(base_analytics)

    base_components = base_analytics["component_summary"]
    scenario_components = scenario_analytics["component_summary"]
    base_suppliers = base_analytics["supplier_summary"]
    scenario_suppliers = scenario_analytics["supplier_summary"]
    scenario_detail = scenario_analytics["component_supplier_detail"]

    base_spend = float(base_components["spend"].sum()) if len(base_components) else 0.0
    scenario_spend = float(scenario_components["spend"].sum()) if len(scenario_components) else 0.0
    savings = float(scenario_metrics.get("estimated_savings", 0.0))
    base_high_risk = int(base_components["high_risk_flag"].sum()) if len(base_components) else 0
    scenario_high_risk = int(scenario_components["high_risk_flag"].sum()) if len(scenario_components) else 0
    base_supplier_count = int(base_suppliers["supplier"].nunique()) if len(base_suppliers) else 0
    scenario_supplier_count = int(scenario_suppliers["supplier"].nunique()) if len(scenario_suppliers) else 0
    uncovered = int(scenario_metrics.get("uncovered_components", 0))
    mitigated_uncovered = int(scenario_metrics.get("mitigated_uncovered_components", 0))
    mitigated_single = int(scenario_metrics.get("mitigated_single_source_components", 0))
    selected_suppliers = (
        scenario_suppliers.loc[scenario_suppliers.get("scenario_role", pd.Series(dtype=str)).eq("Selected Supplier"), "supplier"]
        .dropna()
        .astype(str)
        .tolist()
        if len(scenario_suppliers)
        else []
    )
    mitigation_suppliers = (
        scenario_suppliers.loc[scenario_suppliers.get("scenario_role", pd.Series(dtype=str)).eq("Mitigation Supplier"), "supplier"]
        .dropna()
        .astype(str)
        .tolist()
        if len(scenario_suppliers)
        else []
    )
    gross_savings = float(scenario_metrics.get("estimated_savings", 0.0))
    mitigation_cost = float(scenario_metrics.get("mitigation_cost", 0.0))
    net_savings = float(scenario_metrics.get("net_savings", gross_savings - mitigation_cost))
    risk_reduction = float(scenario_metrics.get("aggregate_risk_reduction", 0.0))
    coverage_share = float(scenario_metrics.get("covered_spend_share", 0.0))
    scenario_top_supplier = scenario_suppliers.sort_values("spend", ascending=False).iloc[0] if len(scenario_suppliers) else None
    scenario_top_component = scenario_components.sort_values("spend", ascending=False).iloc[0] if len(scenario_components) else None
    scenario_highest_risk_component = scenario_components.sort_values("supply_risk_score", ascending=False).iloc[0] if len(scenario_components) else None
    scenario_highest_priority_component = scenario_components.sort_values("strategic_priority_score", ascending=False).iloc[0] if len(scenario_components) else None
    scenario_quadrant_counts = scenario_components["kraljic_quadrant"].value_counts() if len(scenario_components) else pd.Series(dtype=int)
    scenario_single_source_components = scenario_components.loc[scenario_components["single_source_flag"], "component"].tolist() if len(scenario_components) else []
    scenario_high_risk_components = scenario_components.loc[scenario_components["high_risk_flag"], "component"].tolist() if len(scenario_components) else []
    scenario_medium_risk_components = scenario_components.loc[scenario_components["sourcing_risk_level"].eq("Medium"), "component"].tolist() if len(scenario_components) else []
    scenario_low_risk_count = int(scenario_components["sourcing_risk_level"].eq("Low").sum()) if len(scenario_components) else 0
    scenario_component_supplier_lookup = (
        scenario_detail.groupby("component")["supplier"].apply(lambda x: ", ".join(sorted(pd.unique(x)))).to_dict()
        if len(scenario_detail)
        else {}
    )
    mitigation_assignment_rows = []
    if len(scenario_components) and "Single-Source Mitigation Suppliers" in scenario_components.columns:
        mitigation_assignment_rows = (
            scenario_components.loc[
                scenario_components["Single-Source Mitigation Suppliers"].fillna("").ne(""),
                ["component", "Single-Source Mitigation Suppliers"],
            ]
            .to_dict(orient="records")
        )
    scenario_exposure_components = sorted(set(scenario_single_source_components + scenario_high_risk_components))
    scenario_exposure_text = (
        "; ".join([f"{comp}: {scenario_component_supplier_lookup.get(comp, 'Unknown')}" for comp in scenario_exposure_components])
        if scenario_exposure_components
        else "no elevated exposure clusters"
    )
    scenario_quadrant_text = ", ".join([f"{name}: {count}" for name, count in scenario_quadrant_counts.items()]) if not scenario_quadrant_counts.empty else "no quadrant data available"
    scenario_single_source_text = ", ".join(scenario_single_source_components) if scenario_single_source_components else "none found"
    scenario_high_risk_text = ", ".join(scenario_high_risk_components) if scenario_high_risk_components else "none found"
    scenario_medium_risk_text = ", ".join(scenario_medium_risk_components) if scenario_medium_risk_components else "none found"

    mitigation_assignment_text = (
        "; ".join(
            [
                f"{row['component']}: {row['Single-Source Mitigation Suppliers']}"
                for row in mitigation_assignment_rows
            ]
        )
        if mitigation_assignment_rows
        else ""
    )
    what_changed_text = (
        f"What changed: the applied scenario keeps {', '.join(selected_suppliers) if selected_suppliers else 'no selected suppliers'}"
        + (
            f", adds separate mitigation suppliers {', '.join(mitigation_suppliers)}"
            if mitigation_suppliers
            else ""
        )
        + (
            f", and assigns mitigation coverage by component as {mitigation_assignment_text}"
            if mitigation_assignment_text
            else ", and does not assign explicit mitigation coverage"
        )
        + "."
    )
    why_changed_text = (
        f"Why this scenario was applied: it is designed to cover {coverage_share:.0%} of spend, reduce high-risk exposure,"
        f" mitigate {mitigated_single} single-source components and {mitigated_uncovered} uncovered components,"
        f" deliver gross savings of ${gross_savings:,.0f}, incur modeled mitigation cost of ${mitigation_cost:,.0f},"
        f" and land at net savings of ${net_savings:,.0f} with aggregate risk reduction of {risk_reduction:.1f} points."
    )

    scenario_change_text = (
        f"Applied scenario impact: supplier count moves from {base_supplier_count} to {scenario_supplier_count}, "
        f"total spend moves from ${base_spend:,.0f} to ${scenario_spend:,.0f}, "
        f"reflecting modeled savings of ${savings:,.0f}, "
        f"and high sourcing-risk components move from {base_high_risk} to {scenario_high_risk}. "
        f"The scenario leaves {uncovered} uncovered components, mitigates {mitigated_uncovered} uncovered components, "
        f"and mitigates {mitigated_single} single-source components."
    )

    applied_summary_text = (
        f"Applied scenario summary: Kraljic quadrant counts are {scenario_quadrant_text}. "
        f"High sourcing-risk components total {len(scenario_high_risk_components)} ({scenario_high_risk_text}), "
        f"medium sourcing-risk components total {len(scenario_medium_risk_components)} ({scenario_medium_risk_text}), "
        f"and low sourcing-risk components total {scenario_low_risk_count}. "
        f"Single-source components are {scenario_single_source_text}. Suppliers covering any remaining exposed components include {scenario_exposure_text}. "
        f"The highest-risk component is {scenario_highest_risk_component['component']} with a supply risk score of {scenario_highest_risk_component['supply_risk_score']:.1f}. "
        f"The highest strategic priority component is {scenario_highest_priority_component['component']}. "
        f"Top supplier spend sits with {scenario_top_supplier['supplier']} at ${scenario_top_supplier['spend']:,.0f}, "
        f"while {scenario_top_component['component']} is the largest-spend component at ${scenario_top_component['spend']:,.0f}. "
        f"This applied scenario represents the chosen supplier set and mitigation structure, with gross savings of ${gross_savings:,.0f}, "
        f"modeled mitigation cost of ${mitigation_cost:,.0f}, and net savings of ${net_savings:,.0f}."
    )

    return (
        "Base case summary: " + base_summary_text + "\n\n"
        + what_changed_text + "\n\n"
        + why_changed_text + "\n\n"
        + scenario_change_text + "\n\n"
        + applied_summary_text
    )


@st.cache_data(show_spinner=False)
def build_applied_supplier_plan_analytics(
    base_analytics: Dict[str, pd.DataFrame],
    scenario_analytics: Dict[str, pd.DataFrame],
) -> Dict[str, pd.DataFrame]:
    base_supplier_summary = base_analytics["supplier_summary"].copy()
    scenario_supplier_summary = scenario_analytics["supplier_summary"].copy()
    scenario_detail = scenario_analytics["component_supplier_detail"].copy()

    scenario_lookup = scenario_supplier_summary.set_index("supplier").to_dict(orient="index")
    mitigation_component_map: Dict[str, List[str]] = {}
    if len(scenario_detail):
        mitigation_rows = scenario_detail.loc[scenario_detail["spend"].fillna(0.0).eq(0.0)].copy()
        if not mitigation_rows.empty:
            mitigation_component_map = (
                mitigation_rows.groupby("supplier")["component"]
                .apply(lambda values: sorted(pd.unique(values).tolist()))
                .to_dict()
            )
    combined_rows: List[Dict[str, object]] = []
    for row in base_supplier_summary.to_dict(orient="records"):
        supplier_name = row["supplier"]
        scenario_row = scenario_lookup.get(supplier_name)
        if scenario_row is None:
            merged_row = row.copy()
            merged_row["spend"] = 0.0
            merged_row["portfolio_share"] = 0.0
            risk_score = float(row.get("supplier_risk_score", 0.0))
            performance_score = float(row.get("performance_score", 0.0))
            replaceability_score = float(row.get("replaceability_score", 0.0))
            reason_bits = []
            if risk_score >= 60:
                reason_bits.append("higher supplier risk")
            if performance_score <= 45:
                reason_bits.append("weaker performance")
            if replaceability_score >= 60:
                reason_bits.append("higher replaceability")
            removal_basis = ", ".join(reason_bits) if reason_bits else "overall scenario score tradeoffs"
            merged_row["decision"] = "Removed"
            merged_row["decision_reason"] = (
                "This supplier was not retained in the applied scenario because "
                + removal_basis
                + " made it a less attractive fit than the retained supplier set, and its demand is modeled as reallocated elsewhere."
            )
            merged_row["issues"] = "removed"
            merged_row["supplier_action_plan"] = "No new awards under the applied scenario; verify transition, inventory, and continuity controls before full removal."
            merged_row["scenario_role"] = "Removed"
            combined_rows.append(merged_row)
        else:
            merged_row = row.copy()
            merged_row.update(scenario_row)
            mitigated_components = mitigation_component_map.get(supplier_name, [])
            is_selected_supplier = str(merged_row.get("scenario_role", "")) == "Selected Supplier"
            if is_selected_supplier and mitigated_components:
                merged_row["decision"] = "Retained + Mitigation"
                merged_row["decision_reason"] = (
                    "This supplier was retained for direct supply coverage and also assigned as a mitigation option for "
                    + ", ".join(mitigated_components)
                    + "."
                )
                merged_row["issues"] = "retained supplier with mitigation role"
                merged_row["supplier_action_plan"] = (
                    "Keep this supplier in the active award set and qualify or monitor its backup role for "
                    + ", ".join(mitigated_components)
                    + "."
                )
            elif is_selected_supplier:
                merged_row["decision"] = "Retained"
                merged_row["decision_reason"] = "This supplier remains in the applied scenario because it supports the chosen supply structure and coverage profile."
                merged_row["issues"] = "retained"
                merged_row["supplier_action_plan"] = "Retain this supplier in the active award set and manage performance against the applied scenario."
            elif mitigated_components:
                merged_row["decision"] = "Mitigation Only"
                merged_row["decision_reason"] = (
                    "This supplier was not retained for direct awards, but it was assigned as a mitigation option for "
                    + ", ".join(mitigated_components)
                    + "."
                )
                merged_row["issues"] = "mitigation-only supplier"
                merged_row["supplier_action_plan"] = (
                    "Maintain this supplier as a qualified mitigation source for "
                    + ", ".join(mitigated_components)
                    + "."
                )
            combined_rows.append(merged_row)

    combined_supplier_summary = pd.DataFrame(combined_rows)
    role_order = {
        "Selected Supplier": 1,
        "Mitigation Supplier": 2,
        "Removed": 3,
    }
    decision_order = {
        "Retained + Mitigation": 1,
        "Retained": 2,
        "Mitigation Only": 3,
        "Removed": 4,
    }
    combined_supplier_summary["scenario_role_rank"] = combined_supplier_summary["scenario_role"].map(role_order).fillna(9)
    combined_supplier_summary["decision_rank_display"] = combined_supplier_summary["decision"].map(decision_order).fillna(9)
    combined_supplier_summary = combined_supplier_summary.sort_values(
        ["scenario_role_rank", "decision_rank_display", "spend", "supplier"], ascending=[True, True, False, True]
    ).reset_index(drop=True)

    combined_analytics = dict(scenario_analytics)
    combined_analytics["supplier_summary"] = combined_supplier_summary
    return combined_analytics


@st.cache_data(show_spinner=False)
def make_download_bundle(data_map: Dict[str, pd.DataFrame]) -> bytes:
    output = io.StringIO()
    for name, frame in data_map.items():
        output.write(f"## {name}\n")
        output.write(frame.to_csv(index=False))
        output.write("\n")
    return output.getvalue().encode("utf-8")


@st.cache_data(show_spinner=False)
def make_powerpoint_export(
    summary_text: str,
    analytics: Dict[str, pd.DataFrame],
    scenario_applied: bool = False,
) -> bytes:
    try:
        from PIL import Image, ImageDraw, ImageFont
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
        from pptx.util import Inches, Pt
    except ModuleNotFoundError as exc:
        missing_package = str(exc).split("'")[1] if "'" in str(exc) else str(exc)
        raise RuntimeError(
            f"PowerPoint export is unavailable because the `{missing_package}` package is missing from the environment."
        ) from exc

    supplier_summary = analytics["supplier_summary"]
    component_summary = analytics["component_summary"]
    component_supplier_detail = analytics["component_supplier_detail"]
    spend_pareto = analytics["spend_pareto"]
    risk_pareto = analytics["risk_pareto"]
    strategic_pareto = analytics["strategic_pareto"]
    step_plan = build_step_plan(analytics, scenario_applied=scenario_applied)
    visual_pack = build_visual_pack(analytics)

    def add_bullets_slide(prs: Presentation, title: str, bullets: List[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()
        for idx, bullet in enumerate(bullets):
            paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.font.size = Pt(18)

    def add_step_plan_slide(prs: Presentation, title: str, plan_df: pd.DataFrame) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.2), Inches(12.4), Inches(0.55))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True

        intro_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.78), Inches(12.2), Inches(0.5))
        intro_frame = intro_box.text_frame
        intro_frame.text = (
            "This slide sequences the analyst's work from exposure diagnosis through scenario testing and implementation focus."
        )
        intro_frame.paragraphs[0].font.size = Pt(12)

        y = 1.25
        for _, row in plan_df.head(5).iterrows():
            step_box = slide.shapes.add_textbox(Inches(0.55), Inches(y), Inches(12.0), Inches(1.0))
            step_frame = step_box.text_frame
            step_frame.word_wrap = True
            step_frame.margin_left = Pt(2)
            step_frame.margin_right = Pt(2)
            step_frame.margin_top = Pt(1)
            step_frame.margin_bottom = Pt(1)

            step_title = step_frame.paragraphs[0]
            step_title.text = f"Step {int(row['step'])}: {row['action']}"
            step_title.font.size = Pt(14)
            step_title.font.bold = True
            step_title.font.name = "Segoe UI"

            evidence = str(row.get("supporting evidence", "")).strip()
            detail_text = (
                f"Why: {row['justification']} "
                f"Current data points: {evidence}. "
                f"Expected outcome: {row['expected outcome']}."
            )
            detail = step_frame.add_paragraph()
            detail.text = detail_text
            detail.font.size = Pt(10.5)
            detail.font.name = "Segoe UI"
            detail.space_before = Pt(0)
            detail.space_after = Pt(0)
            y += 1.15

    def compact_slide_text(text: str, max_chars: int = 180) -> str:
        cleaned = " ".join(str(text).split())
        if len(cleaned) <= max_chars:
            return cleaned
        sentence_parts = [part.strip() for part in cleaned.split(". ") if part.strip()]
        shortened = sentence_parts[0] if sentence_parts else cleaned[:max_chars]
        if len(shortened) > max_chars:
            shortened = shortened[: max_chars - 1].rstrip() + "..."
        return shortened

    def wrap_slide_text(text: str, max_chars: int = 180, line_width: int = 58) -> str:
        compacted = compact_slide_text(text, max_chars=max_chars)
        return textwrap.fill(compacted, width=line_width, break_long_words=False, break_on_hyphens=False)

    def normalize_slide_text(text: str) -> str:
        return " ".join(str(text).split())

    def wrap_text_to_pixels(draw: "ImageDraw.ImageDraw", text: str, font: "ImageFont.ImageFont", max_width: int) -> str:
        words = str(text).split()
        if not words:
            return ""
        lines: List[str] = []
        current_line = words[0]
        for word in words[1:]:
            candidate = f"{current_line} {word}"
            candidate_width = draw.textbbox((0, 0), candidate, font=font)[2]
            if candidate_width <= max_width:
                current_line = candidate
            else:
                lines.append(current_line)
                current_line = word
        lines.append(current_line)
        return "\n".join(lines)

    def build_wrapped_talking_points(points: List[str]) -> List[str]:
        font = None
        font_candidates = [
            Path("C:/Windows/Fonts/segoeui.ttf"),
            Path("C:/Windows/Fonts/arial.ttf"),
            Path("C:/Windows/Fonts/calibri.ttf"),
            Path("C:/Windows/Fonts/tahoma.ttf"),
            Path("DejaVuSans.ttf"),
            Path("LiberationSans-Regular.ttf"),
        ]
        for font_name in font_candidates:
            try:
                font = ImageFont.truetype(str(font_name), 40)
                break
            except Exception:
                continue
        if font is None:
            font = ImageFont.load_default()
        canvas = Image.new("RGB", (3400, 1800), "white")
        draw = ImageDraw.Draw(canvas)
        max_text_width = 1750
        wrapped_points: List[str] = []
        for point in points[:4]:
            compacted = normalize_slide_text(point)
            wrapped = wrap_text_to_pixels(draw, compacted, font, max_text_width)
            wrapped_points.append(wrapped)
        return wrapped_points

    def format_name_list(values: List[str], max_items: int = 5) -> str:
        items = [str(v) for v in values if str(v).strip()]
        if not items:
            return "none"
        items = items[:max_items]
        if len(items) == 1:
            return items[0]
        if len(items) == 2:
            return f"{items[0]} and {items[1]}"
        return ", ".join(items[:-1]) + f", and {items[-1]}"

    slide_subtitles = {
        "Spend by Supplier": "Supplier spend concentration and negotiation leverage",
        "Spend by Component": "Largest-spend components in the portfolio",
        "Component Analysis Bubble": "Components positioned by spend, risk, and strategic priority",
        "Component-Supplier Detail": "How supplier coverage is distributed across components",
        "Spend Pareto (ABC)": "A-items driving most of total spend",
        "Supplier Concentration by Component": "Components most exposed to single-supplier dependence",
        "Supplier Risk Score": "Suppliers with the highest operational and exposure risk",
        "Risk-Adjusted Pareto": "Components where cost, quality, and supply risk combine most heavily",
        "Strategic Priority Pareto": "Components ranked by a combined score of spend, quality, and supply exposure",
        "Kraljic Positioning": "Business impact plotted against supply risk",
        "Supply Risk Score": "Components with the highest modeled supply risk",
        "Supplier Count by Component": "Components with the thinnest supplier coverage",
        "Strategic Sourcing Outcomes": "Current supplier outcomes by consolidation logic",
        "Supplier Spend by Component Mix": "Supplier spend broken into the exact components driving that exposure",
        "Supplier Spend by Kraljic Mix": "Supplier spend broken into Strategic, Bottleneck, Leverage, and Non-Critical mix",
    }
    def build_slide_findings(title: str) -> List[str]:
        if title == "Spend by Supplier" and not supplier_summary.empty:
            ranked = supplier_summary.sort_values("spend", ascending=False).head(2)
            top_row = ranked.iloc[0]
            second_row = ranked.iloc[1] if len(ranked) > 1 else None
            return [
                (
                    f"{top_row['supplier']} is the largest supplier by spend, and {second_row['supplier']} is close behind at a very similar spend level, so both suppliers should be treated as major commercial and continuity relationships."
                    if second_row is not None
                    else f"{top_row['supplier']} is the largest supplier by spend, so it should be treated as a major commercial and continuity relationship."
                ),
                "From a supply-chain perspective, the suppliers at the top of this chart are the first ones to review for negotiation strategy, service performance, and contingency planning because they have the biggest portfolio impact.",
            ]
        if title == "Spend by Component" and not component_summary.empty:
            top_rows = component_summary.sort_values("spend", ascending=False).head(3)
            top_row = top_rows.iloc[0]
            return [
                f"{top_row['component']} is the highest-spend component, and the top spend group also includes {format_name_list(top_rows['component'].astype(str).tolist()[1:], max_items=2)}.",
                "These are the components where sourcing decisions, quality performance, and supply disruption would have the biggest immediate cost impact on the portfolio.",
            ]
        if title == "Component-Supplier Detail" and not component_supplier_detail.empty:
            phoenix_rows = component_supplier_detail.loc[
                component_supplier_detail["component"].astype(str).eq("Phoenix Feather")
            ]
            twine_rows = component_supplier_detail.loc[
                component_supplier_detail["component"].astype(str).eq("Twine")
            ]
            tied_rows = component_supplier_detail.sort_values(
                ["effective_supplier_share", "spend"], ascending=[False, False]
            ).head(3)
            phoenix_suppliers = phoenix_rows["supplier"].astype(str).tolist()
            twine_suppliers = twine_rows["supplier"].astype(str).tolist()
            return [
                f"Phoenix Feather is currently covered by {format_name_list(phoenix_suppliers, max_items=4)}, and Twine is currently covered by {format_name_list(twine_suppliers, max_items=4)}, so these are the first components to review when we want to confirm whether mitigation is truly creating backup depth.",
                f"The largest effective coverage positions on the chart belong to {format_name_list((tied_rows['supplier'].astype(str) + ' for ' + tied_rows['component'].astype(str)).tolist(), max_items=3)}, which helps show where one supplier still carries a large share of practical coverage.",
            ]
        if title == "Component Analysis Bubble" and not component_summary.empty:
            high_risk_components = component_summary.loc[component_summary["high_risk_flag"], "component"].astype(str).tolist()
            single_source_components = component_summary.loc[component_summary["single_source_flag"], "component"].astype(str).tolist()
            strategic_components = component_summary.loc[component_summary["kraljic_quadrant"].eq("Strategic"), "component"].astype(str).tolist()
            top_spend_component = component_summary.sort_values("spend", ascending=False).iloc[0]["component"]
            return [
                f"This slide combines several supply-chain conditions at once: the high-risk components are {format_name_list(high_risk_components)}, the single-source components are {format_name_list(single_source_components)}, and the Strategic items include {format_name_list(strategic_components)}.",
                f"{top_spend_component} stands out as the highest-spend component, while Twine and Phoenix Feather deserve explicit attention because they are the clearest examples of components where supply exposure can quickly become a continuity problem. The risk axis is the same weighted model built from supplier count, concentration, defects, lead time, external risk, and criticality.",
            ]
        if title == "Spend Pareto (ABC)" and not component_summary.empty:
            a_items = spend_pareto.loc[spend_pareto["spend_abc"].eq("A"), "component"].astype(str).tolist()
            top_row = component_summary.sort_values("spend", ascending=False).iloc[0]
            return [
                f"The A-components by spend are {format_name_list(a_items, max_items=8)}, which means these parts account for the largest share of the purchasing budget and deserve first attention in cost and continuity reviews.",
                f"In this ABC view, A means the highest-value items that drive most spend, B means the middle tier that still needs active management, and C means the lower-value tail that usually gets lighter-touch control; {top_row['component']} is the single largest spend item in the file.",
            ]
        if title == "Supplier Concentration by Component" and not component_summary.empty:
            sorted_rows = component_summary.sort_values(["largest_supplier_share", "supplier_count", "spend"], ascending=[False, True, False])
            top_components = sorted_rows.head(3)["component"].astype(str).tolist()
            tightest_row = sorted_rows.iloc[0]
            return [
                f"The most concentrated components are {format_name_list(top_components)}, and here 'most concentrated' means they have the fewest practical suppliers and the highest share controlled by one supplier.",
                f"{tightest_row['component']} sits at the top of this exposure, so the main supply-chain question is whether to add backup coverage, protect with inventory, or consciously accept that concentration with a clear rationale.",
            ]
        if title == "Supplier Risk Score" and not supplier_summary.empty:
            top_row = supplier_summary.sort_values("supplier_risk_score", ascending=False).iloc[0]
            high_risk_components = component_summary.loc[component_summary["high_risk_flag"], "component"].astype(str).tolist()
            single_source_components = component_summary.loc[component_summary["single_source_flag"], "component"].astype(str).tolist()
            supplier_actions = []
            for _, row in supplier_summary.sort_values(["decision", "supplier_risk_score"], ascending=[True, False]).iterrows():
                supplier_name = str(row["supplier"])
                decision = str(row.get("decision", ""))
                if decision == "Eliminate / De-prioritize":
                    action = "evaluate for exit or de-prioritization"
                elif decision == "Keep and Monitor":
                    action = "retain, but place under active monitoring"
                else:
                    action = "retain as part of the preferred supply base"
                supplier_actions.append(f"{supplier_name}: {action}")
            return [
                f"{top_row['supplier']} is currently the highest-risk supplier, and the main component drivers behind this view are {format_name_list(high_risk_components, max_items=3)} plus the fragile supply positions on {format_name_list(single_source_components, max_items=3)}.",
                "This supplier risk score is a weighted blend of defect rate, lead time, average component risk, portfolio share, single-source exposure, and high-risk component exposure. Recommended actions are: " + "; ".join(supplier_actions),
            ]
        if title == "Risk-Adjusted Pareto" and not component_summary.empty:
            top_rows = component_summary.sort_values("risk_adjusted_spend", ascending=False).head(4)
            top_row = top_rows.iloc[0]
            return [
                f"{top_row['component']} has the highest risk-adjusted burden, and the leading exposure group also includes {format_name_list(top_rows['component'].astype(str).tolist()[1:], max_items=3)}.",
                "In plain English, this chart shows the components where we are spending a lot of money and also carrying a lot of supply-chain pain at the same time, such as concentration, quality trouble, or sourcing risk.",
            ]
        if title == "Strategic Priority Pareto" and not component_summary.empty:
            a_items = strategic_pareto.loc[strategic_pareto["strategic_abc"].eq("A"), "component"].astype(str).tolist()
            top_row = component_summary.sort_values("strategic_priority_score", ascending=False).iloc[0]
            return [
                f"The A-items in this strategic view are {format_name_list(a_items, max_items=6)}, which means these are the first components leadership should look at in sourcing governance and mitigation planning.",
                f"{top_row['component']} sits at the top of the ranking, and this strategic priority score is derived from a blend of spend, defect burden, single-source exposure, and modeled supply risk rather than from spend alone.",
            ]
        if title == "Supply Risk Score" and not component_summary.empty:
            top_rows = component_summary.sort_values("supply_risk_score", ascending=False).head(3)
            top_row = top_rows.iloc[0]
            return [
                f"{top_row['component']} has the highest modeled supply risk, and the top risk group also includes {format_name_list(top_rows['component'].astype(str).tolist()[1:], max_items=2)}.",
                "This risk score is a weighted blend of supplier count, supplier concentration, defect rate, lead time, external risk, and criticality, with the heaviest emphasis placed on structural supply fragility.",
            ]
        if title == "Supplier Count by Component" and not component_summary.empty:
            lowest_rows = component_summary.sort_values(["supplier_count", "spend"], ascending=[True, False]).head(3)
            low_row = lowest_rows.iloc[0]
            return [
                f"{format_name_list(lowest_rows['component'].astype(str).tolist(), max_items=3)} sit at the low end of supplier depth, which means these parts have the least flexibility if there is a disruption or supplier exit.",
                "Low supplier count does not automatically mean failure is imminent, but it does show where fallback options are limited and where resilience could tighten fastest.",
            ]
        if title == "Kraljic Positioning" and not component_summary.empty:
            strategic_items = component_summary.loc[component_summary["kraljic_quadrant"].eq("Strategic"), "component"].astype(str).tolist()
            bottleneck_items = component_summary.loc[component_summary["kraljic_quadrant"].eq("Bottleneck"), "component"].astype(str).tolist()
            leverage_items = component_summary.loc[component_summary["kraljic_quadrant"].eq("Leverage"), "component"].astype(str).tolist()
            return [
                f"The Strategic components are {format_name_list(strategic_items)}, the Bottleneck components are {format_name_list(bottleneck_items)}, and the Leverage components are {format_name_list(leverage_items)}.",
                "The supply-risk side of this chart comes from the weighted risk model built on supplier count, concentration, defects, lead time, external risk, and criticality, so the upper half reflects both structural and operational exposure.",
            ]
        if title == "Strategic Sourcing Outcomes" and not supplier_summary.empty:
            keep_count_local = int(supplier_summary["decision"].eq("Keep / Consolidate To").sum())
            monitor_count_local = int(supplier_summary["decision"].eq("Keep and Monitor").sum())
            exit_count_local = int(supplier_summary["decision"].eq("Eliminate / De-prioritize").sum())
            keep_suppliers = supplier_summary.loc[supplier_summary["decision"].eq("Keep / Consolidate To"), "supplier"].astype(str).tolist()
            monitor_suppliers = supplier_summary.loc[supplier_summary["decision"].eq("Keep and Monitor"), "supplier"].astype(str).tolist()
            exit_suppliers = supplier_summary.loc[supplier_summary["decision"].eq("Eliminate / De-prioritize"), "supplier"].astype(str).tolist()
            action_parts = []
            if keep_suppliers:
                action_parts.append(f"retain and consolidate toward {format_name_list(keep_suppliers, max_items=7)}")
            if monitor_suppliers:
                action_parts.append(f"retain but monitor {format_name_list(monitor_suppliers, max_items=7)}")
            if exit_suppliers:
                action_parts.append(f"evaluate exit or de-prioritization for {format_name_list(exit_suppliers, max_items=7)}")
            return [
                f"The model points to {keep_count_local} stronger consolidation candidates ({format_name_list(keep_suppliers)}), {monitor_count_local} supplier(s) to monitor ({format_name_list(monitor_suppliers)}), and {exit_count_local} supplier(s) that look more replaceable ({format_name_list(exit_suppliers)}).",
                "Recommended actions by supplier group are to " + "; ".join(action_parts) + ".",
            ]
        if title == "Supplier Spend by Component Mix" and not component_supplier_detail.empty:
            mix_frame = component_supplier_detail.copy()
            supplier_totals = mix_frame.groupby("supplier", as_index=False)["spend"].sum().sort_values("spend", ascending=False)
            top_supplier = supplier_totals.iloc[0]["supplier"]
            top_components = (
                mix_frame.loc[mix_frame["supplier"].eq(top_supplier)]
                .sort_values("spend", ascending=False)["component"]
                .astype(str)
                .tolist()[:3]
            )
            potters_text = ""
            potters_mix = mix_frame.loc[mix_frame["supplier"].astype(str).eq("Potters Parts")].sort_values("spend", ascending=False)
            if not potters_mix.empty:
                potters_top = potters_mix.iloc[0]
                potters_text = (
                    f" Potters Parts also has a large concentration in {potters_top['component']}, which is worth calling out because a supplier with a big share tied to one component creates more focused continuity exposure."
                )
            return [
                f"{top_supplier} has the largest total supplier spend, and its current component mix is led by {format_name_list(top_components, max_items=3)}.",
                "This view helps show whether supplier importance comes from broad portfolio coverage or from heavy dependence on a few specific components, which is critical when testing consolidation scenarios."
                + potters_text,
            ]
        if title == "Supplier Spend by Kraljic Mix" and not component_supplier_detail.empty:
            mix_frame = component_supplier_detail.copy()
            if "kraljic_quadrant" not in mix_frame.columns:
                mix_frame = mix_frame.merge(
                    component_summary[["component", "kraljic_quadrant"]].drop_duplicates(subset=["component"]),
                    on="component",
                    how="left",
                )
            strategic_supplier = (
                mix_frame.loc[mix_frame["kraljic_quadrant"].eq("Strategic")]
                .groupby("supplier", as_index=False)["spend"]
                .sum()
                .sort_values("spend", ascending=False)
                .head(1)
            )
            bottleneck_supplier = (
                mix_frame.loc[mix_frame["kraljic_quadrant"].eq("Bottleneck")]
                .groupby("supplier", as_index=False)["spend"]
                .sum()
                .sort_values("spend", ascending=False)
                .head(1)
            )
            strategic_text = (
                f"The largest Strategic mix sits with {strategic_supplier.iloc[0]['supplier']}."
                if not strategic_supplier.empty
                else ""
            )
            bottleneck_text = (
                f" The largest Bottleneck mix sits with {bottleneck_supplier.iloc[0]['supplier']}."
                if not bottleneck_supplier.empty
                else ""
            )
            return [
                "This view separates plain spend concentration from strategic exposure by showing how much of each supplier's spend sits in Strategic, Bottleneck, Leverage, and Non-Critical components.",
                strategic_text
                + bottleneck_text
                + " For Leverage items, the supply chain manager should use competition, negotiation, and volume consolidation to improve cost. "
                + "For Bottleneck items, the priority is backup qualification, inventory protection, and service assurance. "
                + "For Non-Critical items, the focus should be transaction simplification, automation, and reducing administrative effort.",
            ]
        return []

    def add_altair_chart_slide(prs: Presentation, title: str, subtitle: str, chart: alt.Chart, talking_points: List[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.2), Inches(12.4), Inches(0.55))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        subtitle_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.78), Inches(12.4), Inches(0.55))
        subtitle_box.text_frame.text = compact_slide_text(subtitle)
        subtitle_box.text_frame.paragraphs[0].font.size = Pt(13)
        image_bytes = io.BytesIO()
        export_chart = (
            chart.properties(width=1180, height=300, background="white")
            .configure_axis(labelFontSize=11, titleFontSize=12, grid=False)
            .configure_legend(
                orient="bottom",
                direction="horizontal",
                columns=3,
                titleFontSize=9,
                labelFontSize=8,
                symbolSize=70,
                labelLimit=100,
            )
            .configure_view(strokeOpacity=0, fill="white")
        )
        export_chart.save(image_bytes, format="png", scale_factor=2)
        image_bytes.seek(0)
        slide.shapes.add_picture(image_bytes, Inches(0.55), Inches(1.1), width=Inches(12.2), height=Inches(3.0))
        if talking_points:
            wrapped_points = build_wrapped_talking_points(talking_points)
            note_tops = [4.42, 5.10, 5.78, 6.46]
            note_heights = [0.62, 0.62, 0.62, 0.62]
            for idx, wrapped_point in enumerate(wrapped_points[:4]):
                note_box = slide.shapes.add_textbox(Inches(0.72), Inches(note_tops[idx]), Inches(11.2), Inches(note_heights[idx]))
                note_frame = note_box.text_frame
                note_frame.clear()
                note_frame.word_wrap = True
                note_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                note_frame.margin_left = Pt(2)
                note_frame.margin_right = Pt(2)
                note_frame.margin_top = Pt(1)
                note_frame.margin_bottom = Pt(1)
                paragraph = note_frame.paragraphs[0]
                paragraph.text = f"\u2022 {wrapped_point}"
                paragraph.font.size = Pt(9.5)
                paragraph.font.name = "Segoe UI"
                paragraph.font.color.rgb = RGBColor(31, 41, 55)
                paragraph.alignment = PP_ALIGN.LEFT
                paragraph.space_before = Pt(0)
                paragraph.space_after = Pt(0)
                paragraph.level = 0
                paragraph.left_margin = Pt(14)
                paragraph.first_line_indent = Pt(-10)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Supplier and Sourcing Analysis"
    title_slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
    title_slide.placeholders[1].text = "Executive visuals export from the Streamlit planning tool"
    title_slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(18)

    summary_bits = [segment.strip() for segment in summary_text.split(". ") if segment.strip()]
    add_bullets_slide(prs, "Executive Summary", summary_bits[:4])
    for visual in visual_pack:
        title = str(visual["title"])
        subtitle = slide_subtitles.get(title, compact_slide_text(str(visual.get("summary", "")), max_chars=120))
        professor_notes = build_professor_notes(title, analytics, scenario_applied=scenario_applied)
        activity_note = build_professor_activity_note(title, analytics, scenario_applied=scenario_applied)
        talking_points = professor_notes or build_slide_findings(title) or [str(point) for point in visual.get("talking_points", [])]
        if activity_note:
            talking_points = talking_points + [activity_note]
        if title == "Spend by Supplier":
            chart = build_supplier_metric_chart(supplier_summary, "spend", "Spend")
        elif title == "Spend by Component":
            chart = build_component_risk_bar_chart(component_summary, "spend", "Spend", top_n=None)
        elif title == "Component Analysis Bubble":
            chart = build_component_analysis_bubble_chart(component_summary, show_legend=True)
        elif title == "Component-Supplier Detail":
            chart = build_component_supplier_detail_chart(component_supplier_detail, top_n=None, show_legend=False)
        elif title == "Spend Pareto (ABC)":
            chart = build_pareto_chart(spend_pareto, "component", "spend", "spend_cum_share", "spend_abc", "Component", "Spend")
        elif title == "Supplier Concentration by Component":
            chart = build_supplier_concentration_chart(component_summary, top_n=None)
        elif title == "Supplier Risk Score":
            chart = build_supplier_metric_chart(supplier_summary, "supplier_risk_score", "Supplier Risk Score")
        elif title == "Risk-Adjusted Pareto":
            chart = build_pareto_chart(
                risk_pareto,
                "component",
                "risk_adjusted_spend",
                "risk_cum_share",
                "risk_abc",
                "Component",
                "Risk-Adjusted Spend",
                color_col="sourcing_risk_level",
                color_title="Sourcing Risk",
                color_domain=RISK_LEVEL_DOMAIN,
                color_range=RISK_LEVEL_RANGE,
            )
        elif title == "Strategic Priority Pareto":
            chart = build_pareto_chart(
                strategic_pareto,
                "component",
                "strategic_priority_score",
                "strategic_cum_share",
                "strategic_abc",
                "Component",
                "Strategic Priority Score",
            )
        elif title == "Kraljic Positioning":
            chart = build_kraljic_chart(component_summary)
        elif title == "Supply Risk Score":
            chart = build_component_risk_bar_chart(component_summary, "supply_risk_score", "Supply Risk Score", top_n=None)
        elif title == "Supplier Count by Component":
            chart = build_component_risk_bar_chart(component_summary, "supplier_count", "Supplier Count", top_n=None, ascending=True)
        elif title == "Strategic Sourcing Outcomes":
            chart = build_strategic_outcomes_chart(supplier_summary, scenario_applied=scenario_applied)
        elif title == "Supplier Spend by Component Mix":
            chart = build_supplier_component_mix_chart(component_supplier_detail, component_summary, top_n_suppliers=None)
        elif title == "Supplier Spend by Kraljic Mix":
            chart = build_supplier_quadrant_mix_chart(component_supplier_detail, component_summary, top_n_suppliers=None)
        else:
            continue
        add_altair_chart_slide(prs, title, subtitle, chart, talking_points)
    add_step_plan_slide(prs, "Step-by-Step Analyst Action Plan", step_plan)

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def display_metric_row(component_summary: pd.DataFrame, supplier_summary: pd.DataFrame):
    total_spend = supplier_summary["spend"].sum()
    supplier_count = supplier_summary["supplier"].nunique()
    component_count = component_summary["component"].nunique()
    risk_level_counts = component_summary["sourcing_risk_level"].value_counts()
    high_risk_count = int(risk_level_counts.get("High", 0))
    medium_risk_count = int(risk_level_counts.get("Medium", 0))
    low_risk_count = int(risk_level_counts.get("Low", 0))
    c1, c2, c3, c4, c5, c6 = st.columns(6)
    c1.metric("Total Spend", f"${total_spend:,.0f}")
    c2.metric("Suppliers", f"{supplier_count}")
    c3.metric("Components", f"{component_count}")
    c4.metric("High-Risk Components", f"{high_risk_count}")
    c5.metric("Medium-Risk Components", f"{medium_risk_count}")
    c6.metric("Low-Risk Components", f"{low_risk_count}")


def display_assumptions(title: str, assumptions: List[str]):
    st.caption(title)
    for assumption in assumptions:
        st.write(f"- {assumption}")


def show_table(df: pd.DataFrame):
    st.dataframe(df, width="stretch", hide_index=True)


def build_normalized_input_preview(df: pd.DataFrame, diagnostics: List[str]) -> pd.DataFrame:
    preview = df.head(50).copy()
    status_map = {
        "risk_score": None,
        "criticality": None,
    }

    for diagnostic in diagnostics:
        for field_name in status_map:
            if f"`{field_name}`" not in diagnostic:
                continue
            if "was not found in the data source" in diagnostic:
                status_map[field_name] = "Not found in source"
            elif "no numeric values were usable" in diagnostic:
                status_map[field_name] = "No usable numeric values"

    for field_name, status_text in status_map.items():
        if status_text and field_name in preview.columns:
            preview[field_name] = status_text

    return preview


def source_field_missing(diagnostics: List[str], field_name: str) -> bool:
    field_token = f"`{field_name}`"
    for diagnostic in diagnostics:
        if field_token in diagnostic and (
            "was not found in the data source" in diagnostic or "no numeric values were usable" in diagnostic
        ):
            return True
    return False


def build_supplier_spend_summary(supplier_summary: pd.DataFrame, component_summary: pd.DataFrame) -> str:
    if supplier_summary.empty or component_summary.empty:
        return "Supplier spend analysis shows where commercial leverage and supplier concentration sit across the spend base."
    top_supplier = supplier_summary.sort_values("spend", ascending=False).iloc[0]
    top_three_share = float(
        supplier_summary.sort_values("spend", ascending=False).head(3)["spend"].sum() / max(supplier_summary["spend"].sum(), 1.0)
    )
    top_component = component_summary.sort_values("spend", ascending=False).iloc[0]
    return (
        "Supplier spend analysis shows where purchasing volume is concentrated so teams can spot leverage, dependency, and fragmentation. "
        f"Major takeaways: {top_supplier['supplier']} is the largest supplier at ${top_supplier['spend']:,.0f}, the top three suppliers account for {top_three_share:.0%} of supplier spend, "
        f"and {top_component['component']} is the largest-spend component at ${top_component['spend']:,.0f}."
    )


def build_component_analysis_summary(component_summary: pd.DataFrame) -> str:
    if component_summary.empty:
        return "Component analysis shows where supply risk, quality burden, and sourcing complexity sit at the part level."
    single_source_count = int(component_summary["single_source_flag"].sum())
    high_risk_count = int(component_summary["high_risk_flag"].sum())
    strategic_count = int(component_summary["kraljic_quadrant"].eq("Strategic").sum())
    single_source_components = component_summary.loc[component_summary["single_source_flag"], "component"].tolist()
    high_risk_components = component_summary.loc[component_summary["high_risk_flag"], "component"].tolist()
    strategic_components = component_summary.loc[
        component_summary["kraljic_quadrant"].eq("Strategic"), "component"
    ].tolist()
    top_priority = component_summary.sort_values("strategic_priority_score", ascending=False).iloc[0]
    single_source_text = (
        f"the {single_source_count} single-source components are {', '.join(single_source_components)}"
        if single_source_components
        else "no single-source components were found"
    )
    high_risk_text = (
        f"the {high_risk_count} high-risk components are {', '.join(high_risk_components)}"
        if high_risk_components
        else "no high-risk components were found"
    )
    strategic_text = (
        f"the {strategic_count} Strategic-quadrant components include {', '.join(strategic_components[:6])}"
        if strategic_components
        else "no Strategic-quadrant components were identified"
    )
    return (
        "Component analysis connects spend, supplier coverage, quality, and risk so teams can target the parts that most affect continuity and profit. "
        f"Major takeaways: {single_source_text}, {high_risk_text}, {strategic_text}, "
        f"and {top_priority['component']} is the top strategic-priority component with the highest strategic priority score."
    )


def build_component_spend_summary(component_summary: pd.DataFrame) -> str:
    if component_summary.empty:
        return "Component spend analysis shows which parts drive the largest share of purchasing cost."
    top_component = component_summary.sort_values("spend", ascending=False).iloc[0]
    top_three_share = float(
        component_summary.sort_values("spend", ascending=False).head(3)["spend"].sum() / max(component_summary["spend"].sum(), 1.0)
    )
    return (
        "Component spend analysis shows which parts drive the largest share of purchasing cost, so teams know where sourcing changes or quality issues have the biggest financial effect. "
        f"Major takeaways: {top_component['component']} is the highest-spend component at ${top_component['spend']:,.0f}, and the top three components account for {top_three_share:.0%} of total component spend."
    )


def build_component_supplier_detail_summary(component_supplier_detail: pd.DataFrame) -> str:
    if component_supplier_detail.empty:
        return "Component-supplier detail shows how each component's supply coverage is split across suppliers."
    top_component = (
        component_supplier_detail.groupby("component", as_index=False)["spend"].sum().sort_values("spend", ascending=False).iloc[0]
    )
    single_source_components = (
        component_supplier_detail.loc[component_supplier_detail["single_source_flag"], "component"].drop_duplicates().tolist()
    )
    max_effective_share = float(component_supplier_detail["effective_supplier_share"].max()) if len(component_supplier_detail) else 0.0
    most_concentrated_rows = component_supplier_detail.loc[
        component_supplier_detail["effective_supplier_share"].eq(max_effective_share),
        ["supplier", "component", "effective_supplier_share", "spend"],
    ].sort_values(["spend", "component", "supplier"], ascending=[False, True, True])
    concentration_pairs = [
        f"{row['supplier']} for {row['component']}"
        for row in most_concentrated_rows.head(4).to_dict(orient="records")
    ]
    concentration_text = (
        f"{', '.join(concentration_pairs[:-1])}, and {concentration_pairs[-1]} currently carry the largest effective coverage share at {max_effective_share:.0%}"
        if len(concentration_pairs) > 1
        else f"{concentration_pairs[0]} currently carries the largest effective coverage share at {max_effective_share:.0%}"
    )
    single_source_text = (
        "the single-source components are " + ", ".join(single_source_components)
        if single_source_components
        else "no single-source components were found"
    )
    return (
        "Component-supplier detail shows how each component's supply coverage is distributed across its current suppliers, which makes concentration, backup depth, and supplier overlap easier to see. "
        f"Major takeaways: {top_component['component']} is the largest component by spend, {concentration_text}, "
        f"and {single_source_text}."
    )


def build_pareto_summary(spend_pareto: pd.DataFrame, risk_pareto: pd.DataFrame, strategic_pareto: pd.DataFrame) -> str:
    if spend_pareto.empty or risk_pareto.empty or strategic_pareto.empty:
        return "ABC/Pareto analysis is used to separate the small number of items driving most value or risk from the long tail. In general, A items are the highest-priority few that drive most of the total, B items are the middle layer worth selective management attention, and C items are the long tail with lower relative impact."
    spend_a = spend_pareto.loc[spend_pareto["spend_abc"].eq("A"), "component"].tolist()
    spend_a_text = ", ".join(spend_a[:6]) if spend_a else "no A-items were identified"
    return (
        "ABC/Pareto analysis helps supply-chain teams focus scarce effort on the few items that drive the most spend. "
        "How to interpret it: A items are the highest-priority few that account for most of the cumulative total, B items are the middle tier that still merit active management, and C items are the long tail that usually need lighter-touch control. "
        f"Major takeaways: spend A-items are {spend_a_text}."
    )


def build_risk_adjusted_pareto_summary(risk_pareto: pd.DataFrame) -> str:
    if risk_pareto.empty:
        return "Risk-adjusted Pareto highlights the components that combine spend with supply exposure and quality burden."
    risk_a = risk_pareto.loc[risk_pareto["risk_abc"].eq("A"), "component"].tolist()
    top_risk_adjusted = risk_pareto.sort_values("risk_adjusted_spend", ascending=False).iloc[0]
    high_risk_components = risk_pareto.loc[risk_pareto["sourcing_risk_level"].eq("High"), "component"].tolist()
    high_risk_text = (
        f"the high-risk components in this ranking are {', '.join(high_risk_components[:6])}"
        if high_risk_components
        else "no components in this ranking are currently flagged high risk"
    )
    risk_a_text = ", ".join(risk_a[:6]) if risk_a else "no A-items were identified"
    return (
        "Risk-adjusted Pareto combines spend with quality burden, single-source exposure, and sourcing risk so the ranking reflects practical supply-chain impact rather than spend alone. "
        f"Major takeaways: the A-items in this ranking are {risk_a_text}, "
        f"{top_risk_adjusted['component']} has the largest risk-adjusted burden at {top_risk_adjusted['risk_adjusted_spend']:,.0f}, "
        f"and {high_risk_text}."
    )


def build_risk_analysis_summary(component_summary: pd.DataFrame, supplier_risk_assessment: pd.DataFrame) -> str:
    if component_summary.empty:
        return "Risk analysis is used to identify which components and suppliers create the biggest continuity exposure."
    high_risk_components = component_summary.loc[component_summary["high_risk_flag"], "component"].tolist()
    top_risk_component = component_summary.sort_values("supply_risk_score", ascending=False).iloc[0]
    high_risk_suppliers = (
        supplier_risk_assessment.loc[supplier_risk_assessment["Risk Tier"].eq("High"), "Supplier"].tolist()
        if not supplier_risk_assessment.empty and "Risk Tier" in supplier_risk_assessment.columns
        else []
    )
    flagged_supplier_explanation = ""
    if high_risk_suppliers and high_risk_components:
        flagged_supplier_explanation = (
            " They are flagged high risk largely because they support exposed components, especially where supply is single-sourced or otherwise structurally vulnerable."
        )
    return (
        "Risk analysis highlights where supply interruption, quality failure, or concentration could affect service and cost. "
        f"Major takeaways: the highest-risk component is {top_risk_component['component']} with a supply risk score of {top_risk_component['supply_risk_score']:.1f}, "
        f"{'high-risk components are ' + ', '.join(high_risk_components) if high_risk_components else 'no high-risk components were found'}, "
        f"and {'suppliers currently flagged high risk are ' + ', '.join(high_risk_suppliers) if high_risk_suppliers else 'no suppliers are currently flagged high risk'}."
        + flagged_supplier_explanation
    )


def build_risk_score_methodology_note() -> str:
    return (
        "The supply risk score is a weighted model built from supplier count, supplier concentration, defect rate, lead time, "
        "external risk score, and component criticality. Structural supply factors carry the most weight, so fewer suppliers and "
        "higher concentration move the score more than any single external rating."
    )


def build_supplier_risk_methodology_note() -> str:
    return (
        "The supplier risk score is a weighted model built from defect rate, lead time, the average risk of the components the supplier supports, "
        "portfolio share, single-source exposure, and high-risk component exposure. Quality and responsiveness carry the most weight, with extra "
        "risk added when a supplier supports fragile parts of the supply base."
    )


def build_abc_definition_note() -> str:
    return (
        "ABC definitions: A items are the highest-value items that drive most of the total and need the tightest control; "
        "B items are the mid-tier that still need active management; C items are the lower-value tail that usually need lighter-touch control."
    )


def format_name_list(values: List[str], max_items: int = 5) -> str:
    items = [str(v) for v in values if str(v).strip()]
    if not items:
        return "none"
    items = items[:max_items]
    if len(items) == 1:
        return items[0]
    if len(items) == 2:
        return f"{items[0]} and {items[1]}"
    return ", ".join(items[:-1]) + f", and {items[-1]}"


def build_supplier_concentration_summary(component_summary: pd.DataFrame) -> str:
    if component_summary.empty:
        return "Supplier concentration analysis shows where component demand depends too heavily on one supplier."
    top_concentration = component_summary.sort_values("largest_supplier_share", ascending=False).iloc[0]
    single_source_components = component_summary.loc[component_summary["single_source_flag"], "component"].tolist()
    return (
        "Supplier concentration analysis shows where demand is concentrated with one supplier, which is a direct resilience and continuity concern. "
        f"Major takeaways: {top_concentration['component']} has the highest concentration at {top_concentration['largest_supplier_share']:.0%}, "
        f"and {'the single-source components are ' + ', '.join(single_source_components) if single_source_components else 'no single-source components were found'}."
    )


def build_supplier_component_mix_summary(
    component_supplier_detail: pd.DataFrame, component_summary: pd.DataFrame, supplier_summary: pd.DataFrame
) -> str:
    if component_supplier_detail.empty or supplier_summary.empty:
        return "Supplier-component mix shows which components make up each supplier's spend base."
    chart_data = component_supplier_detail.copy()
    top_supplier = supplier_summary.sort_values("spend", ascending=False).iloc[0]["supplier"]
    top_supplier_mix = chart_data.loc[chart_data["supplier"].eq(top_supplier)].sort_values("spend", ascending=False)
    top_supplier_components = top_supplier_mix["component"].astype(str).tolist()[:3]
    share_frame = chart_data.copy()
    supplier_totals = share_frame.groupby("supplier")["spend"].transform("sum")
    share_frame["component_share_within_supplier"] = np.where(
        supplier_totals > 0,
        share_frame["spend"] / supplier_totals,
        0.0,
    )
    top_concentrated_row = share_frame.sort_values(["component_share_within_supplier", "spend"], ascending=[False, False]).iloc[0]
    strategic_suppliers = (
        share_frame.loc[share_frame["kraljic_quadrant"].eq("Strategic")]
        .groupby("supplier", as_index=False)["spend"]
        .sum()
        .sort_values("spend", ascending=False)
    )
    strategic_supplier_text = (
        f"The largest Strategic exposure currently sits with {strategic_suppliers.iloc[0]['supplier']}."
        if not strategic_suppliers.empty
        else ""
    )
    potters_parts_text = ""
    if "Potters Parts" in share_frame["supplier"].astype(str).tolist():
        potters_mix = share_frame.loc[share_frame["supplier"].astype(str).eq("Potters Parts")].sort_values("spend", ascending=False)
        if not potters_mix.empty:
            potters_top = potters_mix.iloc[0]
            potters_parts_text = (
                f" Potters Parts also has a notably large block in {potters_top['component']}, which matters because heavy concentration in one component "
                "can turn a large supplier relationship into a targeted continuity risk if that part is disrupted."
            )
    return (
        "This visual combines supplier spend with component mix so teams can see not just how much spend sits with each supplier, "
        "but which exact components create that exposure. "
        f"Major takeaways: {top_supplier} is the largest supplier by spend and its mix is led by {format_name_list(top_supplier_components, max_items=3)}, "
        f"while the most concentrated supplier-component position is {top_concentrated_row['supplier']} on {top_concentrated_row['component']} at "
        f"{top_concentrated_row['component_share_within_supplier']:.0%} of that supplier's spend. {strategic_supplier_text}{potters_parts_text}".strip()
    )


def build_supplier_quadrant_mix_summary(
    component_supplier_detail: pd.DataFrame, component_summary: pd.DataFrame, supplier_summary: pd.DataFrame
) -> str:
    if component_supplier_detail.empty or supplier_summary.empty:
        return "Supplier quadrant mix shows how each supplier's spend is split across Strategic, Bottleneck, Leverage, and Non-Critical components."
    chart_data = component_supplier_detail.copy()
    if "kraljic_quadrant" not in chart_data.columns:
        chart_data = chart_data.merge(
            component_summary[["component", "kraljic_quadrant"]].drop_duplicates(subset=["component"]),
            on="component",
            how="left",
        )
    quadrant_spend = (
        chart_data.groupby(["supplier", "kraljic_quadrant"], as_index=False)["spend"]
        .sum()
        .sort_values("spend", ascending=False)
    )
    strategic_supplier = quadrant_spend.loc[quadrant_spend["kraljic_quadrant"].eq("Strategic")].head(1)
    bottleneck_supplier = quadrant_spend.loc[quadrant_spend["kraljic_quadrant"].eq("Bottleneck")].head(1)
    leverage_supplier = quadrant_spend.loc[quadrant_spend["kraljic_quadrant"].eq("Leverage")].head(1)
    strategic_text = (
        f"The largest Strategic mix currently sits with {strategic_supplier.iloc[0]['supplier']}."
        if not strategic_supplier.empty
        else ""
    )
    bottleneck_text = (
        f" The largest Bottleneck mix sits with {bottleneck_supplier.iloc[0]['supplier']}."
        if not bottleneck_supplier.empty
        else ""
    )
    leverage_text = (
        f" The largest Leverage mix sits with {leverage_supplier.iloc[0]['supplier']}."
        if not leverage_supplier.empty
        else ""
    )
    return (
        "This visual shows how each supplier's spend is distributed across Strategic, Bottleneck, Leverage, and Non-Critical components, "
        "which helps separate simple spend concentration from true strategic exposure. "
        + strategic_text
        + bottleneck_text
        + leverage_text
    )


def build_professor_notes(title: str, analytics: Dict[str, pd.DataFrame], scenario_applied: bool = False) -> List[str]:
    supplier_summary = analytics["supplier_summary"]
    component_summary = analytics["component_summary"]
    component_supplier_detail = analytics["component_supplier_detail"]
    spend_pareto = analytics["spend_pareto"]
    risk_pareto = analytics["risk_pareto"]
    strategic_pareto = analytics["strategic_pareto"]

    if title == "Spend by Supplier" and not supplier_summary.empty:
        top_suppliers = supplier_summary.sort_values("spend", ascending=False).head(4)["supplier"].astype(str).tolist()
        top_supplier = supplier_summary.sort_values("spend", ascending=False).iloc[0]
        return [
            f"This chart shows where the money is concentrated across suppliers. In this dataset, {format_name_list(top_suppliers, max_items=4)} sit at the top, so these relationships matter most commercially.",
            f"The conclusion is that supplier size matters because large suppliers create leverage for negotiation, but they also create dependency if too much spend is concentrated in a small part of the supply base. Specifically, the leading suppliers are {format_name_list(top_suppliers, max_items=4)}, with {top_supplier['supplier']} carrying the largest supplier spend at ${top_supplier['spend']:,.0f}.",
            f"The action is to start with {format_name_list(top_suppliers, max_items=4)}: review their contracts, compare pricing and service performance, and confirm whether contingency coverage is strong enough if one of these large suppliers is disrupted.",
        ]
    if title == "Spend by Component" and not component_summary.empty:
        top_components = component_summary.sort_values("spend", ascending=False).head(5)["component"].astype(str).tolist()
        top_component = component_summary.sort_values("spend", ascending=False).iloc[0]
        return [
            f"This chart is showing which components drive the most money through the portfolio. Here, the leading spend components are {format_name_list(top_components, max_items=5)}.",
            f"The conclusion is that these components should receive disproportionate management attention because even small changes in price, quality, or supply continuity will move the overall outcome. In this file, the main spend drivers are {format_name_list(top_components, max_items=5)}, with {top_component['component']} as the largest item at ${top_component['spend']:,.0f}.",
            f"The action is to focus sourcing and planning work on {format_name_list(top_components, max_items=5)}: review current pricing, check supplier capacity and performance for these parts, and confirm whether any of them need backup sourcing or added inventory protection.",
        ]
    if title == "Component Analysis Bubble" and not component_summary.empty:
        high_risk = component_summary.loc[component_summary["high_risk_flag"], "component"].astype(str).tolist()
        single_source = component_summary.loc[component_summary["single_source_flag"], "component"].astype(str).tolist()
        strategic_items = component_summary.loc[component_summary["kraljic_quadrant"].eq("Strategic"), "component"].astype(str).tolist()
        return [
            f"This bubble chart is highlighting several signals at once: spend on the vertical axis, supply risk on the horizontal axis, and strategic priority through bubble size. In this file, {format_name_list(high_risk, max_items=3)} stand out as high-risk and {format_name_list(single_source, max_items=3)} stand out as single-source.",
            f"The conclusion is that components combining high spend, high risk, and strategic importance should move to the front of the sourcing agenda because they create the biggest operational and financial consequences. Here that especially points toward {format_name_list(high_risk, max_items=3)} and Strategic items such as {format_name_list(strategic_items, max_items=4)}.",
            f"The action is to use this chart to prioritize mitigation on {format_name_list(high_risk, max_items=3)} and to place Strategic items such as {format_name_list(strategic_items, max_items=4)} into active supplier review, backup qualification, or executive escalation.",
        ]
    if title == "Component-Supplier Detail" and not component_supplier_detail.empty:
        top_pairs = (
            component_supplier_detail.sort_values(["effective_supplier_share", "spend"], ascending=[False, False])
            .head(5)
        )
        pair_text = format_name_list((top_pairs["supplier"].astype(str) + " for " + top_pairs["component"].astype(str)).tolist(), max_items=5)
        return [
            "This chart is showing how each component's practical supply coverage is distributed across suppliers, which is how you see backup depth and dependence at the part level rather than only at the supplier level.",
            f"The conclusion is that components with one dominant bar or very few suppliers are structurally less resilient than components whose coverage is spread across several sources. The most concentrated visible positions here are {pair_text}.",
            f"The action is to start with {pair_text}: confirm whether each of those components has a viable backup source, assign mitigation suppliers where needed, and avoid removing suppliers tied to those parts until continuity protection is in place.",
        ]
    if title == "Spend Pareto (ABC)" and not spend_pareto.empty:
        a_items = spend_pareto.loc[spend_pareto["spend_abc"].eq("A"), "component"].astype(str).tolist()
        top_component = spend_pareto.sort_values("spend", ascending=False).iloc[0]
        return [
            f"This Pareto chart is highlighting the ABC concept. A-items are the small set of components that drive most spend, and in this file they are {format_name_list(a_items, max_items=8)}.",
            f"The conclusion is that not all items deserve the same level of management effort; a relatively small number of parts usually control most of the economic outcome. In this file, {top_component['component']} is the leading spend driver at ${top_component['spend']:,.0f}.",
            f"The action is to apply the tightest commercial and continuity discipline to A-items such as {format_name_list(a_items, max_items=8)}, while using lighter planning and review controls on B and C items.",
            build_abc_definition_note(),
        ]
    if title == "Supplier Concentration by Component" and not component_summary.empty:
        concentration_rows = component_summary.sort_values(["largest_supplier_share", "supplier_count"], ascending=[False, True]).head(5)
        top_components = concentration_rows["component"].astype(str).tolist()
        concentration_examples = format_name_list(
            (concentration_rows["component"].astype(str) + " via " + concentration_rows["dominant_supplier"].astype(str)).tolist(),
            max_items=5,
        )
        return [
            f"This chart is showing concentration risk at the component level. The components with bars farthest to the right, such as {format_name_list(top_components, max_items=5)}, are the ones most dependent on one supplier.",
            f"The conclusion is that concentration is a resilience issue because a disruption at the dominant supplier will hit those components harder and faster than a diversified item. In this file, the clearest concentration examples are {concentration_examples}, which is why parts like Phoenix Feather and Twine should be reviewed together when they show the same exposure pattern.",
            f"The action is to review the most concentrated components one by one, starting with {format_name_list(top_components, max_items=5)}, and decide whether to qualify another supplier, hold more safety stock, or explicitly accept the current dependence on {concentration_rows.iloc[0]['dominant_supplier']} and similar dominant suppliers.",
        ]
    if title == "Supplier Risk Score" and not supplier_summary.empty:
        top_supplier = supplier_summary.sort_values("supplier_risk_score", ascending=False).iloc[0]["supplier"]
        high_suppliers = supplier_summary.sort_values("supplier_risk_score", ascending=False).head(5)["supplier"].astype(str).tolist()
        return [
            f"This chart is highlighting supplier risk as a portfolio concept, not just a vendor scorecard. {top_supplier} ranks highest because the model blends quality, lead time, portfolio exposure, and support for fragile components.",
            f"The conclusion is that a supplier can become important and risky at the same time, especially when it supports single-source or high-risk parts. The suppliers most exposed on this view are {format_name_list(high_suppliers, max_items=5)}.",
            f"The action is to put {format_name_list(high_suppliers, max_items=5)} under closer review, inspect which exposed components they support, and decide whether to improve oversight, shift share, or add backup coverage on the affected parts.",
        ]
    if title == "Risk-Adjusted Pareto" and not risk_pareto.empty:
        top_items = risk_pareto.sort_values("risk_adjusted_spend", ascending=False).head(5)["component"].astype(str).tolist()
        top_row = risk_pareto.sort_values("risk_adjusted_spend", ascending=False).iloc[0]
        return [
            f"This chart is showing a more advanced concept than plain spend: it ranks components by a combined burden of cost, quality, and supply exposure. The top names here are {format_name_list(top_items, max_items=5)}.",
            f"The conclusion is that the most important items are not always the biggest spend items; some components become more urgent because risk and quality burden amplify their true management importance. In this file, the main risk-adjusted priorities are {format_name_list(top_items, max_items=5)}, with {top_row['component']} highest at {top_row['risk_adjusted_spend']:,.0f}.",
            f"The action is to start cross-functional work with {format_name_list(top_items, max_items=5)}, because these components combine commercial importance with real supply-chain burden and should be reviewed jointly by sourcing, planning, and quality teams.",
        ]
    if title == "Strategic Priority Pareto" and not strategic_pareto.empty:
        a_items = strategic_pareto.loc[strategic_pareto["strategic_abc"].eq("A"), "component"].astype(str).tolist()
        top_priority_items = strategic_pareto.sort_values("strategic_priority_score", ascending=False).head(5)["component"].astype(str).tolist()
        top_row = strategic_pareto.sort_values("strategic_priority_score", ascending=False).iloc[0]
        return [
            f"This chart is highlighting strategic prioritization. The A-items in this ranking are {format_name_list(a_items, max_items=6)}, meaning these are the first items leadership should discuss in sourcing governance.",
            f"The conclusion is that strategic attention should go where business impact and supply exposure combine, not just where spend is high. Here, the leading strategic priorities are {format_name_list(top_priority_items, max_items=5)}, with {top_row['component']} highest at {top_row['strategic_priority_score']:.1f}.",
            f"The action is to use {format_name_list(top_priority_items, max_items=5)} as the first components in scenario design, mitigation planning, and executive sourcing review.",
        ]
    if title == "Kraljic Positioning" and not component_summary.empty:
        strategic_items = component_summary.loc[component_summary["kraljic_quadrant"].eq("Strategic"), "component"].astype(str).tolist()
        bottleneck_items = component_summary.loc[component_summary["kraljic_quadrant"].eq("Bottleneck"), "component"].astype(str).tolist()
        leverage_items = component_summary.loc[component_summary["kraljic_quadrant"].eq("Leverage"), "component"].astype(str).tolist()
        non_critical_items = component_summary.loc[component_summary["kraljic_quadrant"].eq("Non-Critical"), "component"].astype(str).tolist()
        return [
            "This chart is highlighting the Kraljic framework by placing components according to business impact and supply risk. It helps separate Strategic, Bottleneck, Leverage, and Non-Critical items.",
            f"The conclusion is that different types of components require different sourcing behavior; a one-size-fits-all strategy is not appropriate across all quadrants. In this file, Strategic items include {format_name_list(strategic_items, max_items=6)}, Bottleneck items include {format_name_list(bottleneck_items, max_items=6)}, Leverage items include {format_name_list(leverage_items, max_items=6)}, and Non-Critical items include {format_name_list(non_critical_items, max_items=6)}.",
            f"The action is to negotiate and consolidate Leverage items such as {format_name_list(leverage_items, max_items=6)}, protect and back up Bottleneck items such as {format_name_list(bottleneck_items, max_items=6)}, invest in continuity on Strategic items such as {format_name_list(strategic_items, max_items=6)}, and simplify management of Non-Critical items such as {format_name_list(non_critical_items, max_items=6)}.",
        ]
    if title == "Supply Risk Score" and not component_summary.empty:
        top_items = component_summary.sort_values("supply_risk_score", ascending=False).head(5)["component"].astype(str).tolist()
        top_row = component_summary.sort_values("supply_risk_score", ascending=False).iloc[0]
        return [
            f"This chart is showing which components have the highest modeled supply risk. The most exposed items in this file are {format_name_list(top_items, max_items=5)}.",
            f"The conclusion is that these items are the most likely to drive continuity problems because the risk model weights supplier depth, concentration, quality, lead time, external risk, and criticality. Here, the highest-risk set is {format_name_list(top_items, max_items=5)}, with {top_row['component']} highest at {top_row['supply_risk_score']:.1f}.",
            f"The action is to start mitigation, backup qualification, and supplier-management effort with {format_name_list(top_items, max_items=5)} before lower-risk items are addressed.",
        ]
    if title == "Supplier Count by Component" and not component_summary.empty:
        lowest_items = component_summary.sort_values(["supplier_count", "spend"], ascending=[True, False]).head(5)["component"].astype(str).tolist()
        low_row = component_summary.sort_values(["supplier_count", "spend"], ascending=[True, False]).iloc[0]
        return [
            f"This chart is showing structural resilience in its simplest form: how many suppliers support each component. The thinnest-coverage items here are {format_name_list(lowest_items, max_items=5)}.",
            f"The conclusion is that low supplier count means less flexibility and fewer recovery options if a disruption occurs. In this file, the lowest-depth set includes {format_name_list(lowest_items, max_items=5)}, with {low_row['component']} at {int(low_row['supplier_count'])} supplier(s).",
            f"The action is to review low-count components such as {format_name_list(lowest_items, max_items=5)} for backup qualification, mitigation suppliers, or explicit risk acceptance.",
        ]
    if title == "Strategic Sourcing Outcomes" and not supplier_summary.empty:
        keep_suppliers = supplier_summary.loc[supplier_summary["decision"].eq("Keep / Consolidate To"), "supplier"].astype(str).tolist()
        exit_suppliers = supplier_summary.loc[supplier_summary["decision"].eq("Eliminate / De-prioritize"), "supplier"].astype(str).tolist()
        monitor_suppliers = supplier_summary.loc[supplier_summary["decision"].eq("Keep and Monitor"), "supplier"].astype(str).tolist()
        return [
            f"This chart is showing the current sourcing recommendation logic at the supplier level. In this file, the likely consolidation targets are {format_name_list(keep_suppliers, max_items=6)}, while potential exit candidates include {format_name_list(exit_suppliers, max_items=4)}.",
            f"The conclusion is that supplier strategy should combine economics with resilience, not pursue savings in isolation. Here the keep group is {format_name_list(keep_suppliers, max_items=6)}, the monitor group is {format_name_list(monitor_suppliers, max_items=5)}, and the exit group is {format_name_list(exit_suppliers, max_items=5)}.",
            f"The action is to test the keep, monitor, and exit logic supplier by supplier, especially for {format_name_list(exit_suppliers, max_items=5)} and {format_name_list(monitor_suppliers, max_items=5)}, before implementation.",
        ]
    if title == "Supplier Spend by Component Mix" and not component_supplier_detail.empty:
        mix_frame = component_supplier_detail.copy()
        top_blocks = mix_frame.sort_values("spend", ascending=False).head(5)
        block_text = format_name_list((top_blocks["supplier"].astype(str) + " on " + top_blocks["component"].astype(str)).tolist(), max_items=5)
        return [
            "This chart is highlighting that supplier spend should be broken into the components creating that spend. It helps distinguish a broadly important supplier from a supplier whose importance is concentrated in only one or two parts.",
            f"The conclusion is that a supplier can look large for very different reasons, and concentrated component exposure often creates more focused continuity risk than diversified exposure. In this file, the clearest large supplier-component blocks include {block_text}.",
            f"The action is to examine large supplier-component blocks such as {block_text} for concentration, single-source exposure, and whether backup coverage exists before changing supplier awards.",
        ]
    if title == "Supplier Spend by Kraljic Mix" and not component_supplier_detail.empty:
        mix_frame = component_supplier_detail.copy()
        if "kraljic_quadrant" not in mix_frame.columns:
            mix_frame = mix_frame.merge(
                component_summary[["component", "kraljic_quadrant"]].drop_duplicates(subset=["component"]),
                on="component",
                how="left",
            )
        strategic_supplier = (
            mix_frame.loc[mix_frame["kraljic_quadrant"].eq("Strategic")]
            .groupby("supplier", as_index=False)["spend"]
            .sum()
            .sort_values("spend", ascending=False)
            .head(3)
        )
        return [
            "This chart is highlighting how supplier spend maps into Strategic, Bottleneck, Leverage, and Non-Critical demand. That lets you see not just how big a supplier is, but how strategically sensitive its spend mix is.",
            "The conclusion is that two suppliers with similar spend can require very different management if one is loaded with Strategic and Bottleneck demand while the other mainly supports Leverage or Non-Critical items."
            + (
                f" In this file, the suppliers carrying the largest Strategic mix are {format_name_list(strategic_supplier['supplier'].astype(str).tolist(), max_items=3)}."
                if not strategic_supplier.empty
                else ""
            ),
            "The action is to segment suppliers by their quadrant mix, negotiate and consolidate Leverage-heavy relationships, protect Bottleneck-heavy relationships, and simplify the handling of suppliers dominated by Non-Critical demand.",
        ]
    return []


def build_professor_activity_note(title: str, analytics: Dict[str, pd.DataFrame], scenario_applied: bool = False) -> str:
    component_summary = analytics["component_summary"]
    supplier_summary = analytics["supplier_summary"]
    if title == "Spend by Supplier":
        return "Potential supply-chain activities include preparing scorecards for the largest suppliers, scheduling business reviews, reviewing contract leverage, and validating contingency coverage on the biggest spend relationships."
    if title == "Spend by Component":
        return "Potential supply-chain activities include launching cost reviews on the top-spend parts, checking supplier capacity, reviewing defect history, and validating whether alternate sources or inventory buffers are needed."
    if title == "Component Analysis Bubble":
        return "Potential supply-chain activities include opening dual-source qualification work, escalating the riskiest strategic parts into leadership review, testing mitigation suppliers, and prioritizing service-protection plans for the exposed components."
    if title == "Component-Supplier Detail":
        return "Potential supply-chain activities include identifying components with thin coverage, assigning backup suppliers, validating technical interchangeability, and setting inventory guardrails for parts where backup qualification will take time."
    if title == "Spend Pareto (ABC)":
        return "Potential supply-chain activities include setting review cadence by ABC class, concentrating negotiations and savings projects on A-items, and simplifying planning controls for lower-value B and C items."
    if title == "Supplier Concentration by Component":
        return "Potential supply-chain activities include creating a concentrated-part watchlist, prioritizing alternate source development, confirming safety-stock policies, and documenting which concentration risks are being actively accepted versus reduced."
    if title == "Supplier Risk Score":
        return "Potential supply-chain activities include tightening supplier scorecards, increasing review cadence, creating supplier-specific mitigation plans, and reducing award share where performance and exposure remain unacceptable."
    if title == "Risk-Adjusted Pareto":
        return "Potential supply-chain activities include creating a cross-functional priority list, aligning sourcing, quality, and planning teams around the highest-burden items, and sequencing mitigation and negotiation work from this ranking."
    if title == "Strategic Priority Pareto":
        return "Potential supply-chain activities include building an executive watchlist, assigning owners to the top strategic items, prioritizing scenario work for those parts, and reviewing whether commercial and resilience plans are aligned."
    if title == "Kraljic Positioning":
        return "Potential supply-chain activities include running sourcing events for Leverage items, qualifying alternates and buffers for Bottleneck items, assigning continuity plans to Strategic items, and simplifying ordering rules for Non-Critical items."
    if title == "Supply Risk Score":
        return "Potential supply-chain activities include launching mitigation workflows, qualifying alternate suppliers, increasing monitoring on the riskiest parts, and aligning planners and buyers on short-term protection measures."
    if title == "Supplier Count by Component":
        return "Potential supply-chain activities include building a single-source reduction plan, documenting approved backups, reviewing inventory strategies for low-depth parts, and escalating unresolved gaps into scenario planning."
    if title == "Strategic Sourcing Outcomes":
        return "Potential supply-chain activities include validating remove-versus-retain scenarios, preparing transition plans for exit candidates, setting monitoring actions for watchlist suppliers, and confirming that savings do not create new supply gaps."
    if title == "Supplier Spend by Component Mix":
        return "Potential supply-chain activities include mapping supplier-component dependencies, checking whether large blocks align to strategic or single-source parts, and prioritizing mitigation work where one supplier carries a disproportionate share of a critical component."
    if title == "Supplier Spend by Kraljic Mix":
        return "Potential supply-chain activities include segmenting suppliers by quadrant mix, aligning supplier-management playbooks to that mix, and making sure Strategic and Bottleneck-heavy suppliers have stronger continuity and governance controls."
    return ""


def build_kraljic_positioning_summary(component_summary: pd.DataFrame) -> str:
    if component_summary.empty:
        return "Kraljic positioning is used to sort components by profit impact and supply risk so teams can choose the right sourcing strategy."
    quadrant_counts = component_summary["kraljic_quadrant"].value_counts()
    strategic_components = component_summary.loc[component_summary["kraljic_quadrant"].eq("Strategic"), "component"].tolist()
    bottleneck_components = component_summary.loc[component_summary["kraljic_quadrant"].eq("Bottleneck"), "component"].tolist()
    leverage_components = component_summary.loc[component_summary["kraljic_quadrant"].eq("Leverage"), "component"].tolist()
    non_critical_components = component_summary.loc[component_summary["kraljic_quadrant"].eq("Non-Critical"), "component"].tolist()
    strategic_text = ", ".join(strategic_components[:5]) if strategic_components else "none"
    bottleneck_text = ", ".join(bottleneck_components[:5]) if bottleneck_components else "none"
    leverage_text = ", ".join(leverage_components[:5]) if leverage_components else "none"
    non_critical_text = ", ".join(non_critical_components[:5]) if non_critical_components else "none"
    return (
        "Kraljic positioning maps components by supply risk and profit impact so teams can separate Strategic, Bottleneck, "
        "Leverage, and Non-Critical items and align sourcing actions accordingly. "
        f"What it is telling us: Strategic components total {int(quadrant_counts.get('Strategic', 0))} "
        f"({strategic_text}), Bottleneck components total {int(quadrant_counts.get('Bottleneck', 0))} "
        f"({bottleneck_text}), Leverage components total {int(quadrant_counts.get('Leverage', 0))} "
        f"({leverage_text}), and Non-Critical components total {int(quadrant_counts.get('Non-Critical', 0))} "
        f"({non_critical_text})."
    )


def load_data(uploaded_file) -> Tuple[pd.DataFrame, str, List[str], pd.DataFrame]:
    if uploaded_file is None:
        return get_default_data()
    return load_uploaded_data(uploaded_file.name, uploaded_file.getvalue())


def load_persisted_scenario_state() -> Dict[str, object]:
    # Keep scenario drafts inside Streamlit session state so hosted deployments
    # do not rely on local filesystem writes that may be ephemeral or unavailable.
    snapshot = st.session_state.get("_scenario_state_snapshot", {})
    return snapshot if isinstance(snapshot, dict) else {}


def save_persisted_scenario_state(payload: Dict[str, object]) -> None:
    normalized_payload = payload if isinstance(payload, dict) else {}
    st.session_state["persisted_scenario_state"] = normalized_payload
    st.session_state["_scenario_state_snapshot"] = normalized_payload


def build_scenario_compare_snapshot(
    label: str,
    selected_suppliers: List[str],
    mitigation_assignments: List[str],
    metrics: Dict[str, object],
    scorecard: Dict[str, object],
) -> Dict[str, object]:
    breakdown = scorecard.get("breakdown", pd.DataFrame())
    return {
        "label": label,
        "signature": {
            "selected_suppliers": tuple(sorted(selected_suppliers)),
            "mitigation_assignments": tuple(sorted(mitigation_assignments)),
        },
        "selected_suppliers": list(selected_suppliers),
        "mitigation_assignments": list(mitigation_assignments),
        "selected_supplier_count": int(metrics.get("selected_supplier_count", 0)),
        "mitigation_supplier_count": int(metrics.get("mitigation_supplier_count", 0)),
        "covered_spend_share": float(metrics.get("covered_spend_share", 0.0)),
        "uncovered_components": int(metrics.get("uncovered_components", 0)),
        "mitigated_single_source_components": int(metrics.get("mitigated_single_source_components", 0)),
        "mitigated_uncovered_components": int(metrics.get("mitigated_uncovered_components", 0)),
        "estimated_savings": float(metrics.get("estimated_savings", 0.0)),
        "mitigation_cost": float(metrics.get("mitigation_cost", 0.0)),
        "net_savings": float(metrics.get("net_savings", 0.0)),
        "aggregate_risk_reduction": float(metrics.get("aggregate_risk_reduction", 0.0)),
        "score": float(scorecard.get("score", 0.0)),
        "optimization_objective": str(scorecard.get("optimization_objective", "Best Overall")),
        "score_breakdown": breakdown if isinstance(breakdown, pd.DataFrame) else pd.DataFrame(),
    }


def build_scenario_compare_table(current_snapshot: Dict[str, object], saved_snapshots: List[Dict[str, object]]) -> pd.DataFrame:
    metric_rows = [
        ("Selected suppliers", "selected_supplier_count", "number"),
        ("Mitigation suppliers", "mitigation_supplier_count", "number"),
        ("Covered spend %", "covered_spend_share", "percent"),
        ("Uncovered components", "uncovered_components", "number"),
        ("Mitigated single-source", "mitigated_single_source_components", "number"),
        ("Mitigated uncovered", "mitigated_uncovered_components", "number"),
        ("Gross savings", "estimated_savings", "currency"),
        ("Mitigation cost", "mitigation_cost", "currency"),
        ("Net savings", "net_savings", "currency"),
        ("Risk reduction", "aggregate_risk_reduction", "decimal"),
        ("Scenario score", "score", "decimal"),
    ]
    scenario_columns = [current_snapshot] + list(saved_snapshots)
    rows: List[Dict[str, object]] = []
    for display_name, key, value_type in metric_rows:
        row: Dict[str, object] = {"Metric": display_name}
        for snapshot in scenario_columns:
            raw_value = snapshot.get(key, 0)
            if value_type == "percent":
                formatted = f"{float(raw_value):.0%}"
            elif value_type == "currency":
                formatted = f"${float(raw_value):,.0f}"
            elif value_type == "decimal":
                formatted = f"{float(raw_value):,.1f}"
            else:
                formatted = f"{int(raw_value)}"
            row[str(snapshot.get("label", "Scenario"))] = formatted
        rows.append(row)
    return pd.DataFrame(rows)


def render_app():
    st.title("Supplier and Sourcing Analysis Tool")
    st.caption("Single-pipeline supplier decisions, sourcing priorities, and executive visuals.")
    st.caption("Build: 2026-04-07 cc0cf29")
    render_glossary_drawer()
    if "persisted_scenario_state" not in st.session_state:
        st.session_state["persisted_scenario_state"] = load_persisted_scenario_state()
    if "applied_scenario" not in st.session_state:
        st.session_state["applied_scenario"] = None
    if "scenario_comparisons" not in st.session_state:
        st.session_state["scenario_comparisons"] = {}

    uploaded_file = st.file_uploader("Upload purchasing data", type=["csv", "xlsx", "xls"])
    if uploaded_file is None:
        st.info("Using built-in sample data. Upload your own file to replace it.")

    with st.spinner("Loading and analyzing supplier data..."):
        try:
            normalized_df, data_source_label, input_diagnostics, input_field_status = load_data(uploaded_file)
            base_analytics = build_analytics(normalized_df)
        except Exception as exc:
            st.error(f"Unable to load file: {exc}")
            st.stop()

    with st.expander("Data Quality & Inferred Fields"):
        st.write(build_data_quality_summary(input_field_status))
        if input_diagnostics:
            for note in input_diagnostics:
                st.write(f"- {note}")
        show_table(input_field_status)

    previous_data_source_label = st.session_state.get("active_data_source_label")
    if previous_data_source_label != data_source_label:
        st.session_state["active_data_source_label"] = data_source_label
        st.session_state["applied_scenario"] = None
        st.session_state["scenario_builder"] = {}
        st.session_state["scenario_recommendation"] = {}
        scenario_comparisons = st.session_state.get("scenario_comparisons", {})
        if isinstance(scenario_comparisons, dict):
            scenario_comparisons[data_source_label] = []
            st.session_state["scenario_comparisons"] = scenario_comparisons
        st.session_state.pop("pending_scenario_recommendation", None)
        st.session_state.pop("scenario_selected_suppliers", None)
        st.session_state["scenario_builder_source_label"] = data_source_label
        for key in list(st.session_state.keys()):
            if key.startswith("mitigation_") or key.startswith("uncovered_"):
                del st.session_state[key]
        save_persisted_scenario_state(
            {
                "data_source_label": data_source_label,
                "scenario_builder": {},
                "applied_scenario": None,
            }
        )

    persisted_state = st.session_state.get("persisted_scenario_state", {})
    persisted_source = persisted_state.get("data_source_label")

    scenario_state = st.session_state.get("applied_scenario")
    if isinstance(scenario_state, dict) and scenario_state.get("data_source_label") not in (None, data_source_label):
        scenario_state = None
        st.session_state["applied_scenario"] = None
    applied_scenario_metrics = None
    applied_scenario_assumptions: List[str] = []
    if scenario_state:
        analytics, applied_scenario_metrics, _, applied_scenario_assumptions = build_applied_scenario_analytics(
            base_analytics,
            tuple(scenario_state.get("selected_suppliers", [])),
            tuple(scenario_state.get("mitigation_assignments", [])),
        )
    else:
        analytics = base_analytics

    if scenario_state and applied_scenario_metrics is not None:
        summary_text = build_applied_executive_summary(base_analytics, analytics, applied_scenario_metrics)
        applied_plan_analytics = build_applied_supplier_plan_analytics(base_analytics, analytics)
        _, executive_actions, supplier_action_plan = build_executive_summary(applied_plan_analytics, scenario_applied=True)
    else:
        summary_text, executive_actions, supplier_action_plan = build_executive_summary(analytics, scenario_applied=False)
    _, _, base_supplier_action_plan = build_executive_summary(base_analytics, scenario_applied=False)
    scenario_applied = bool(scenario_state and applied_scenario_metrics is not None)
    step_plan = build_step_plan(analytics, scenario_applied=scenario_applied)
    consolidation_plan, consolidation_assumptions = build_supplier_consolidation_plan(
        applied_plan_analytics if scenario_applied else analytics,
        scenario_applied=scenario_applied,
    )
    supplier_risk_assessment, supplier_risk_assumptions = build_supplier_risk_assessment(analytics, scenario_applied=scenario_applied)
    if source_field_missing(input_diagnostics, "risk_score") and "External Risk Score" in supplier_risk_assessment.columns:
        supplier_risk_assessment["External Risk Score"] = "Not in source"
    strategic_sourcing_plan, strategic_sourcing_assumptions = build_strategic_sourcing_plan(analytics, scenario_applied=scenario_applied)
    visual_pack = build_visual_pack(analytics)

    st.caption(data_source_label)
    if scenario_state and applied_scenario_metrics is not None:
        st.info(
            "Applied scenario dashboard view: "
            f"{applied_scenario_metrics['selected_supplier_count']} selected suppliers, "
            f"{applied_scenario_metrics['mitigation_supplier_count']} mitigation suppliers, "
            f"{applied_scenario_metrics['covered_spend_share']:.0%} covered spend."
        )
        if st.button("Revert To Base Dashboard", type="secondary"):
            st.session_state["applied_scenario"] = None
            save_persisted_scenario_state(
                {
                    "data_source_label": data_source_label,
                    "scenario_builder": st.session_state.get("scenario_builder", {}),
                    "applied_scenario": None,
                }
            )
            st.rerun()

    supplier_summary = analytics["supplier_summary"]
    component_summary = analytics["component_summary"]
    component_supplier_detail = analytics["component_supplier_detail"]
    spend_pareto = analytics["spend_pareto"]
    risk_pareto = analytics["risk_pareto"]
    strategic_pareto = analytics["strategic_pareto"]

    display_metric_row(component_summary, supplier_summary)

    tabs = st.tabs(
        [
            "Executive Summary",
            "Supplier Spend",
            "Component Analysis",
            "ABC / Pareto",
            "Risk Analysis",
            "Supplier Action Plans",
            "Executive Visuals & Downloads",
            "Supplier Scenarios",
            "Supplier Consolidation Plan",
        ]
    )

    with tabs[0]:
        st.subheader("Executive Summary")
        st.write(summary_text)
        st.subheader("Executive Actions")
        show_table(executive_actions)

    with tabs[1]:
        st.subheader("Supplier Spend Analysis")
        st.write(build_supplier_spend_summary(supplier_summary, component_summary))
        st.altair_chart(build_supplier_metric_chart(supplier_summary, "spend", "Spend"), width="stretch")
        st.subheader("Component Spend Analysis")
        st.write(build_component_spend_summary(component_summary))
        st.altair_chart(build_component_risk_bar_chart(component_summary, "spend", "Spend", top_n=None), width="stretch")
        show_table(
            supplier_summary[["supplier", "spend", "component_count", "defect_rate", "avg_lead_time", "supplier_risk_score", "decision"]].rename(
                columns={"supplier_risk_score": "Supplier Risk Score"}
            ),
        )

    with tabs[2]:
        st.subheader("Component Analysis")
        st.write(build_component_analysis_summary(component_summary))
        st.caption("Bubble size shows strategic priority, color shows sourcing risk, and diamond markers highlight single-source components.")
        st.altair_chart(build_component_analysis_bubble_chart(component_summary), width="stretch")
        st.subheader("Component-Supplier Detail")
        st.write(build_component_supplier_detail_summary(component_supplier_detail))
        st.caption("This visual shows how each component's effective supply coverage is split across its suppliers. In an applied scenario, mitigation suppliers count toward the visible coverage share even if current awarded spend has not shifted yet. Longer single-color bars indicate concentrated coverage, while more segmented bars indicate broader backup depth.")
        st.altair_chart(build_component_supplier_detail_chart(component_supplier_detail, top_n=None), width="stretch")
        show_table(
            component_summary[
                [
                    "component",
                    "spend",
                    "supplier_count",
                    "sourcing_risk_level",
                    "largest_supplier_share",
                    "defect_rate",
                    "single_source_flag",
                    "backup_supplier_flag",
                    "supply_risk_score",
                    "profit_impact_score",
                    "kraljic_quadrant",
                    "negotiation_leverage_score",
                    "strategic_priority_score",
                ]
            ]
        )
        show_table(
            component_supplier_detail[
                [
                    "component",
                    "supplier",
                    "spend",
                    "supplier_share",
                    "defect_rate",
                    "supplier_count",
                    "single_source_flag",
                    "backup_supplier_flag",
                    "sourcing_risk_level",
                    "supply_risk_score",
                    "profit_impact_score",
                    "kraljic_quadrant",
                    "strategic_priority_score",
                    "decision",
                ]
            ]
        )

    with tabs[3]:
        st.subheader("Spend Pareto (ABC)")
        st.write(build_pareto_summary(spend_pareto, risk_pareto, strategic_pareto))
        st.caption("This visual shows which components account for most of total spend. The bars show each component's individual spend, and the line shows the cumulative share so you can see where the A, B, and C breakpoints occur. A-items are the small set of components driving the largest cumulative share of spend and should get the most commercial attention.")
        st.altair_chart(
            build_pareto_chart(spend_pareto, "component", "spend", "spend_cum_share", "spend_abc", "Component", "Spend"),
            width="stretch",
        )
        show_table(spend_pareto[["component", "spend", "spend_cum_share", "spend_abc"]])

        st.subheader("Risk-Adjusted Pareto")
        st.write(build_risk_adjusted_pareto_summary(risk_pareto))
        st.caption("This visual shows risk-adjusted spend by component. The bars show each component's individual risk-adjusted burden, and the line shows the cumulative share so you can see where the A, B, and C breakpoints occur. The tooltip and table below show the actual sourcing risk tier for each component.")
        st.altair_chart(
            build_pareto_chart(
                risk_pareto,
                "component",
                "risk_adjusted_spend",
                "risk_cum_share",
                "risk_abc",
                "Component",
                "Risk-Adjusted Spend",
                color_col="sourcing_risk_level",
                color_title="Sourcing Risk",
                color_domain=RISK_LEVEL_DOMAIN,
                color_range=RISK_LEVEL_RANGE,
            ),
            width="stretch",
        )
        show_table(
            risk_pareto[["component", "sourcing_risk_level", "risk_adjusted_spend", "defect_rate", "single_source_flag", "supply_risk_score", "risk_cum_share", "risk_abc"]],
        )

        st.subheader("Strategic Priority Pareto")
        st.caption("This visual ranks components by overall strategic importance so teams can focus scenario work where business impact and supply exposure are both highest. The bars show each component's individual strategic priority score, and the line shows the cumulative share so you can see where the A, B, and C breakpoints occur. A-items in this chart are the first components to review in sourcing strategy discussions.")
        st.altair_chart(
            build_pareto_chart(
                strategic_pareto,
                "component",
                "strategic_priority_score",
                "strategic_cum_share",
                "strategic_abc",
                "Component",
                "Strategic Priority Score",
            ),
            width="stretch",
        )
        show_table(
            strategic_pareto[["component", "strategic_priority_score", "strategic_cum_share", "strategic_abc"]],
        )
        st.subheader("Pareto 4 Quadrant Chart")
        st.caption("This visual combines supply risk and profit impact so each component can be interpreted through the Kraljic lens. The upper-right Strategic area contains the components where continuity planning and executive supplier management matter most.")
        st.altair_chart(build_kraljic_chart(component_summary), width="stretch")

    with tabs[4]:
        st.subheader("Risk Analysis")
        st.write(build_risk_analysis_summary(component_summary, supplier_risk_assessment))
        st.caption(
            "This visual ranks components by modeled supply risk so teams can see which items are most exposed. "
            "The components at the top of the chart are the ones most likely to need mitigation, backup qualification, or closer supplier management. "
            + build_risk_score_methodology_note()
        )
        st.altair_chart(build_component_risk_bar_chart(component_summary, "supply_risk_score", "Supply Risk Score"), width="stretch")
        show_table(
            component_summary[
                ["component", "supply_risk_score", "defect_rate", "avg_lead_time", "risk_score", "criticality", "single_source_flag", "high_risk_flag"]
            ]
        )
        st.subheader("Kraljic Positioning")
        st.write(build_kraljic_positioning_summary(component_summary))
        st.caption(build_risk_score_methodology_note())
        st.altair_chart(build_kraljic_chart(component_summary), width="stretch")
        show_table(
            component_summary[
                ["component", "supplier_count", "sourcing_risk_level", "supply_risk_score", "profit_impact_score", "kraljic_quadrant"]
            ]
        )
        st.subheader("Supplier Concentration")
        st.write(build_supplier_concentration_summary(component_summary))
        left, right = st.columns(2)
        with left:
            st.caption("This visual shows the largest-supplier share for each component so concentration risk is easy to spot. Components with bars extending farthest to the right have the highest concentration and are the most exposed to supplier disruption.")
            st.altair_chart(build_supplier_concentration_chart(component_summary, top_n=None), width="stretch")
        with right:
            st.caption("This visual shows how many suppliers support each component. Components with supplier counts near one are the tightest resilience constraints in the portfolio.")
            st.altair_chart(
                build_component_risk_bar_chart(component_summary, "supplier_count", "Supplier Count", top_n=None, ascending=True),
                width="stretch",
            )

    with tabs[7]:
        st.subheader("Supplier Scenarios")
        scenario_supplier_summary = base_analytics["supplier_summary"]
        scenario_component_supplier_detail = base_analytics["component_supplier_detail"]
        pending_recommendation = st.session_state.pop("pending_scenario_recommendation", None)
        if isinstance(pending_recommendation, dict) and pending_recommendation.get("data_source_label") == data_source_label:
            recommended_suppliers = list(pending_recommendation.get("selected_suppliers", []))
            recommended_assignments = list(pending_recommendation.get("mitigation_assignments", []))
            assignment_map: Dict[str, List[str]] = {}
            for assignment in recommended_assignments:
                if "|||" not in assignment:
                    continue
                component_name, supplier_name = assignment.split("|||", 1)
                assignment_map.setdefault(component_name, []).append(supplier_name)

            st.session_state["scenario_selected_suppliers"] = recommended_suppliers
            for component_name in base_analytics["component_summary"]["component"].tolist():
                recommended_component_suppliers = assignment_map.get(component_name, [])
                st.session_state[f"mitigation_{component_name}"] = recommended_component_suppliers
                st.session_state[f"uncovered_{component_name}"] = recommended_component_suppliers

            st.session_state["scenario_builder"] = {
                "selected_suppliers": recommended_suppliers,
                "mitigation_assignments": recommended_assignments,
            }
            st.session_state["scenario_recommendation"] = {
                "data_source_label": data_source_label,
                "selected_suppliers": recommended_suppliers,
                "mitigation_assignments": recommended_assignments,
                "rationale": pending_recommendation.get("rationale", ""),
                "tested_scenarios": int(pending_recommendation.get("tested_scenarios", 0)),
                "score": float(pending_recommendation.get("score", 0.0)),
                "optimization_objective": str(pending_recommendation.get("optimization_objective", "Best Overall")),
                "covered_spend_share": float(pending_recommendation.get("covered_spend_share", 0.0)),
                "estimated_savings": float(pending_recommendation.get("estimated_savings", 0.0)),
                "net_savings": float(pending_recommendation.get("net_savings", pending_recommendation.get("estimated_savings", 0.0))),
                "high_risk_components": int(pending_recommendation.get("high_risk_components", 0)),
                "uncovered_components": int(pending_recommendation.get("uncovered_components", 0)),
                "score_breakdown": pending_recommendation.get("score_breakdown", pd.DataFrame()),
            }
            save_persisted_scenario_state(
                {
                    "data_source_label": data_source_label,
                    "scenario_builder": st.session_state["scenario_builder"],
                    "applied_scenario": None,
                }
            )
        supplier_single_source_components = (
            scenario_component_supplier_detail.loc[scenario_component_supplier_detail["single_source_flag"], ["supplier", "component"]]
            .groupby("supplier")["component"]
            .apply(lambda values: ", ".join(sorted(pd.unique(values))))
            .to_dict()
        )
        supplier_label_map = {
            row["supplier"]: (
                f"{row['supplier']} (Risk Score: {row['supplier_risk_score']:.1f}; "
                + (
                    f"Single-source provider: Yes - {supplier_single_source_components[row['supplier']]})"
                    if row["supplier"] in supplier_single_source_components
                    else "Single-source provider: No)"
                )
            )
            for row in scenario_supplier_summary[["supplier", "supplier_risk_score"]].to_dict(orient="records")
        }
        recommended_keep_suppliers = (
            scenario_supplier_summary.loc[scenario_supplier_summary["decision"].eq("Keep / Consolidate To"), ["supplier", "spend"]]
            .sort_values("spend", ascending=False)["supplier"]
            .tolist()
        )
        all_supplier_options = scenario_supplier_summary.sort_values("spend", ascending=False)["supplier"].tolist()
        default_count = min(4, len(all_supplier_options)) if all_supplier_options else 0
        base_selected_suppliers = list(all_supplier_options)
        persisted_builder = (
            persisted_state.get("scenario_builder", {})
            if persisted_source == data_source_label and isinstance(persisted_state.get("scenario_builder"), dict)
            else {}
        )
        applied_builder = scenario_state if isinstance(scenario_state, dict) else {}
        default_selected_suppliers = persisted_builder.get("selected_suppliers") or applied_builder.get("selected_suppliers") or (
            recommended_keep_suppliers[:default_count] if recommended_keep_suppliers else all_supplier_options[:default_count]
        )
        default_selected_suppliers = [supplier for supplier in default_selected_suppliers if supplier in all_supplier_options]
        if st.session_state.pop("pending_revert_to_base_scenario", False):
            st.session_state["applied_scenario"] = None
            st.session_state["scenario_builder"] = {}
            st.session_state["scenario_selected_suppliers"] = base_selected_suppliers
            st.session_state["scenario_assignment_supplier_signature"] = tuple(sorted(base_selected_suppliers))
            for key in list(st.session_state.keys()):
                if key.startswith("mitigation_") or key.startswith("uncovered_"):
                    del st.session_state[key]
            save_persisted_scenario_state(
                {
                    "data_source_label": data_source_label,
                    "scenario_builder": {},
                    "applied_scenario": None,
                }
            )
        if st.session_state.get("scenario_builder_source_label") != data_source_label:
            st.session_state["scenario_builder_source_label"] = data_source_label
            st.session_state["scenario_selected_suppliers"] = default_selected_suppliers
        st.caption(
            "Base-model keep recommendations: "
            + (", ".join(recommended_keep_suppliers) if recommended_keep_suppliers else "none identified")
        )
        scenario_component_summary = base_analytics["component_summary"]
        scenario_single_source_names = scenario_component_summary.loc[
            scenario_component_summary["single_source_flag"], "component"
        ].astype(str).tolist()
        scenario_high_risk_names = scenario_component_summary.loc[
            scenario_component_summary["high_risk_flag"], "component"
        ].astype(str).tolist()
        top_supplier_candidates = scenario_supplier_summary.sort_values("spend", ascending=False)["supplier"].astype(str).tolist()[:3]
        with st.container(border=True):
            st.markdown("**Suggested Workflow**")
            st.write(
                "1. Generate the recommended scenario as a starting point.\n"
                "2. Review which suppliers are being retained and which components still need explicit mitigation.\n"
                "3. Evaluate the draft after each material change so the metrics below reflect the current supplier and mitigation structure.\n"
                "4. Apply the scenario to the dashboard only after the coverage, savings, and risk tradeoffs look acceptable."
            )
            analysis_guidance_parts = []
            if scenario_single_source_names:
                analysis_guidance_parts.append(
                    "The first components to protect are the single-source items "
                    + format_name_list(scenario_single_source_names, max_items=6)
                    + "."
                )
            if scenario_high_risk_names:
                analysis_guidance_parts.append(
                    "The highest-risk components to watch during scenario testing are "
                    + format_name_list(scenario_high_risk_names, max_items=6)
                    + "."
                )
            if top_supplier_candidates:
                analysis_guidance_parts.append(
                    "The suppliers with the largest spend positions in the base case are "
                    + format_name_list(top_supplier_candidates, max_items=3)
                    + ", so changes involving them tend to move savings and coverage the most."
                )
            if not analysis_guidance_parts:
                analysis_guidance_parts.append(
                    "Use the scenario workflow to test whether supplier count can be reduced without creating uncovered demand or new high-risk exposure."
                )
            st.caption("Guidance from the current analysis: " + " ".join(analysis_guidance_parts))
        optimization_objective = st.selectbox(
            "Optimization objective",
            options=["Best Overall", "Best Net Savings", "Best Risk Reduction"],
            key="scenario_optimization_objective",
            help="Choose how strongly the recommendation engine should emphasize balanced performance, net savings, or risk reduction.",
        )
        evaluated_builder = (
            st.session_state.get("scenario_builder", {})
            if st.session_state.get("scenario_builder_source_label") == data_source_label
            else {}
        )
        with st.container(border=True):
            st.markdown("**Step 1. Generate A Starting Scenario**")
            st.caption("Choose the optimization objective, then let the model produce a starting supplier and mitigation structure.")
            recommend_col, spacer_col = st.columns([1, 3])
            with recommend_col:
                if st.button("Recommend Best Scenario", type="secondary", use_container_width=True):
                    recommendation = recommend_best_supplier_scenario(base_analytics, optimization_objective=optimization_objective)
                    st.session_state["pending_scenario_recommendation"] = {
                        "data_source_label": data_source_label,
                        "selected_suppliers": list(recommendation.get("selected_suppliers", ())),
                        "mitigation_assignments": list(recommendation.get("mitigation_assignments", ())),
                        "rationale": recommendation.get("rationale", ""),
                        "tested_scenarios": int(recommendation.get("tested_scenarios", 0)),
                        "score": float(recommendation.get("scorecard", {}).get("score", 0.0)),
                        "optimization_objective": str(recommendation.get("scorecard", {}).get("optimization_objective", optimization_objective)),
                        "covered_spend_share": float(recommendation.get("metrics", {}).get("covered_spend_share", 0.0)),
                        "estimated_savings": float(recommendation.get("metrics", {}).get("estimated_savings", 0.0)),
                        "net_savings": float(recommendation.get("metrics", {}).get("net_savings", 0.0)),
                        "high_risk_components": int(recommendation.get("scorecard", {}).get("high_risk_components", 0)),
                        "uncovered_components": int(recommendation.get("metrics", {}).get("uncovered_components", 0)),
                        "score_breakdown": recommendation.get("scorecard", {}).get("breakdown", pd.DataFrame()),
                    }
                    st.rerun()

        recommendation_state = st.session_state.get("scenario_recommendation", {})
        recommended_assignment_map: Dict[str, List[str]] = {}
        if recommendation_state.get("data_source_label") == data_source_label:
            recommended_supplier_list = recommendation_state.get("selected_suppliers", [])
            recommended_mitigation_list = recommendation_state.get("mitigation_assignments", [])
            for assignment in recommended_mitigation_list:
                if "|||" not in assignment:
                    continue
                component_name, supplier_name = assignment.split("|||", 1)
                recommended_assignment_map.setdefault(component_name, []).append(supplier_name)
            current_draft_suppliers = st.session_state.get("scenario_selected_suppliers", default_selected_suppliers)
            current_selected_tuple = tuple(sorted(current_draft_suppliers))
            recommended_selected_tuple = tuple(sorted(recommended_supplier_list))
            if current_selected_tuple == recommended_selected_tuple:
                st.success("Your current selected supplier set matches the model's recommended scenario.")
            st.info(
                f"{recommendation_state.get('optimization_objective', optimization_objective)} recommendation: score {recommendation_state.get('score', 0.0):,.1f}, {len(recommended_supplier_list)} suppliers, "
                f"{recommendation_state.get('covered_spend_share', 0.0):.0%} spend coverage, "
                f"{recommendation_state.get('high_risk_components', 0)} high-risk components, {recommendation_state.get('uncovered_components', 0)} uncovered components, "
                f"and ${recommendation_state.get('net_savings', recommendation_state.get('estimated_savings', 0.0)):,.0f} net savings."
            )
            st.caption(
                "Recommended suppliers: "
                + (", ".join(recommended_supplier_list) if recommended_supplier_list else "none")
                + ". Mitigation assignments: "
                + (", ".join(recommended_mitigation_list) if recommended_mitigation_list else "none needed")
                + "."
            )
            st.caption(
                f"Scenario engine reviewed {recommendation_state.get('tested_scenarios', 0)} candidate combinations across supplier counts. "
                + str(recommendation_state.get("rationale", ""))
            )
            recommendation_breakdown = recommendation_state.get("score_breakdown", pd.DataFrame())
            if isinstance(recommendation_breakdown, pd.DataFrame) and not recommendation_breakdown.empty:
                with st.expander("Recommended Scenario Score Breakdown"):
                    show_table(recommendation_breakdown)
        if "scenario_selected_suppliers" not in st.session_state:
            st.session_state["scenario_selected_suppliers"] = default_selected_suppliers
        current_draft_supplier_tuple = tuple(
            sorted(
                supplier
                for supplier in st.session_state.get("scenario_selected_suppliers", [])
                if supplier in all_supplier_options
            )
        )
        previous_assignment_signature = st.session_state.get("scenario_assignment_supplier_signature")
        if (
            previous_assignment_signature != current_draft_supplier_tuple
            and len(current_draft_supplier_tuple) == len(all_supplier_options)
        ):
            for key in list(st.session_state.keys()):
                if key.startswith("mitigation_") or key.startswith("uncovered_"):
                    del st.session_state[key]
        st.session_state["scenario_assignment_supplier_signature"] = current_draft_supplier_tuple

        def get_matching_assignment_map(builder_state: Optional[Dict[str, object]]) -> Dict[str, List[str]]:
            if not isinstance(builder_state, dict):
                return {}
            builder_suppliers = tuple(sorted(builder_state.get("selected_suppliers", [])))
            if builder_suppliers != current_draft_supplier_tuple:
                return {}

            assignment_map: Dict[str, List[str]] = {}
            for assignment in builder_state.get("mitigation_assignments", []):
                if "|||" not in assignment:
                    continue
                component_name, supplier_name = assignment.split("|||", 1)
                assignment_map.setdefault(component_name, []).append(supplier_name)
            return assignment_map

        saved_assignment_map = (
            get_matching_assignment_map(evaluated_builder)
            or get_matching_assignment_map(persisted_builder)
            or get_matching_assignment_map(applied_builder)
        )

        mitigation_assignments: List[str] = []
        current_form_widget_values: Dict[str, List[str]] = {}
        with st.container(border=True):
            st.markdown("**Step 2. Adjust The Scenario Draft**")
            st.caption("Select the suppliers you want to keep, then assign explicit mitigation or pickup suppliers where the draft still creates exposure.")
            with st.form("scenario_builder_form"):
                selected_suppliers = st.multiselect(
                    "Choose the suppliers for the general supplier scenario",
                    options=all_supplier_options,
                    max_selections=len(all_supplier_options),
                    format_func=lambda supplier: supplier_label_map.get(supplier, supplier),
                    key="scenario_selected_suppliers",
                )
                st.caption(f"{len(selected_suppliers)} suppliers selected.")
                if selected_suppliers:
                    single_source_components = get_dynamic_single_source_candidates(base_analytics, selected_suppliers)
                    uncovered_components = get_uncovered_candidates(base_analytics, selected_suppliers)
                else:
                    single_source_components = []
                    uncovered_components = []
                    st.caption("No suppliers selected. Evaluating this scenario will show the downside of removing all current supply coverage.")
                if single_source_components:
                    st.caption("Assign mitigation suppliers to the components that are currently single-source under your general consolidation selection.")
                for row in single_source_components:
                    component_name = row["component"]
                    incumbent_supplier = row["dominant_supplier"]
                    mitigation_options = [supplier for supplier in all_supplier_options if supplier != incumbent_supplier]
                    mitigation_key = f"mitigation_{component_name}"
                    recommended_component_suppliers = [supplier for supplier in recommended_assignment_map.get(component_name, []) if supplier in mitigation_options]
                    current_component_suppliers = [supplier for supplier in saved_assignment_map.get(component_name, []) if supplier in mitigation_options]
                    existing_component_suppliers = [supplier for supplier in st.session_state.get(mitigation_key, []) if supplier in mitigation_options]
                    if mitigation_key not in st.session_state or previous_assignment_signature != current_draft_supplier_tuple:
                        if current_component_suppliers:
                            st.session_state[mitigation_key] = current_component_suppliers
                        elif len(current_draft_supplier_tuple) == len(all_supplier_options):
                            st.session_state[mitigation_key] = []
                        else:
                            st.session_state[mitigation_key] = existing_component_suppliers
                    if current_component_suppliers:
                        st.caption(
                            f"Current evaluated mitigation for {component_name}: {', '.join(current_component_suppliers)}"
                        )
                    if recommended_component_suppliers:
                        st.caption(
                            f"App recommendation for {component_name}: {', '.join(recommended_component_suppliers)}"
                        )
                    mitigation_choices = st.multiselect(
                        f"Mitigation suppliers for {component_name} (current supplier: {incumbent_supplier})",
                        options=mitigation_options,
                        key=mitigation_key,
                        format_func=lambda supplier: supplier_label_map.get(supplier, supplier),
                        help=(
                            "App-recommended mitigation: " + ", ".join(recommended_component_suppliers)
                            if recommended_component_suppliers
                            else None
                        ),
                    )
                    current_form_widget_values[mitigation_key] = list(mitigation_choices)
                    for supplier_name in mitigation_choices:
                        mitigation_assignments.append(f"{component_name}|||{supplier_name}")
                if uncovered_components:
                    st.caption("Assign pickup suppliers to components that are currently uncovered under your general consolidation selection.")
                for row in uncovered_components:
                    component_name = row["component"]
                    prior_supplier = row["dominant_supplier"]
                    uncovered_key = f"uncovered_{component_name}"
                    recommended_component_suppliers = [supplier for supplier in recommended_assignment_map.get(component_name, []) if supplier in all_supplier_options]
                    current_component_suppliers = [supplier for supplier in saved_assignment_map.get(component_name, []) if supplier in all_supplier_options]
                    existing_component_suppliers = [supplier for supplier in st.session_state.get(uncovered_key, []) if supplier in all_supplier_options]
                    if uncovered_key not in st.session_state or previous_assignment_signature != current_draft_supplier_tuple:
                        if current_component_suppliers:
                            st.session_state[uncovered_key] = current_component_suppliers
                        elif len(current_draft_supplier_tuple) == len(all_supplier_options):
                            st.session_state[uncovered_key] = []
                        else:
                            st.session_state[uncovered_key] = existing_component_suppliers
                    if current_component_suppliers:
                        st.caption(
                            f"Current evaluated pickup for {component_name}: {', '.join(current_component_suppliers)}"
                        )
                    if recommended_component_suppliers:
                        st.caption(
                            f"App recommendation for {component_name}: {', '.join(recommended_component_suppliers)}"
                        )
                    pickup_choices = st.multiselect(
                        f"Pickup suppliers for uncovered component {component_name} (previous dominant supplier: {prior_supplier})",
                        options=all_supplier_options,
                        key=uncovered_key,
                        format_func=lambda supplier: supplier_label_map.get(supplier, supplier),
                        help=(
                            "App-recommended pickup supplier: " + ", ".join(recommended_component_suppliers)
                            if recommended_component_suppliers
                            else None
                        ),
                    )
                    current_form_widget_values[uncovered_key] = list(pickup_choices)
                    for supplier_name in pickup_choices:
                        mitigation_assignments.append(f"{component_name}|||{supplier_name}")
                st.caption("When the draft looks right, evaluate it to refresh the scenario metrics and impact view below.")
                evaluate_submitted = st.form_submit_button("Evaluate Selected Scenario", type="secondary", use_container_width=True)
        draft_scenario_builder = {
            "selected_suppliers": list(selected_suppliers),
            "mitigation_assignments": sorted(mitigation_assignments),
        }
        has_evaluated_builder = bool(evaluated_builder.get("selected_suppliers") or evaluated_builder.get("mitigation_assignments"))
        draft_matches_evaluated = (
            tuple(sorted(draft_scenario_builder["selected_suppliers"])) == tuple(sorted(evaluated_builder.get("selected_suppliers", [])))
            and tuple(sorted(draft_scenario_builder["mitigation_assignments"])) == tuple(sorted(evaluated_builder.get("mitigation_assignments", [])))
        )
        if has_evaluated_builder and not draft_matches_evaluated:
            st.warning("Draft scenario changes are pending. Click `Evaluate Selected Scenario` to refresh the scenario results below.")
        display_selected_suppliers = (
            evaluated_builder.get("selected_suppliers", draft_scenario_builder["selected_suppliers"])
            if has_evaluated_builder
            else draft_scenario_builder["selected_suppliers"]
        )
        display_mitigation_assignments = (
            evaluated_builder.get("mitigation_assignments", draft_scenario_builder["mitigation_assignments"])
            if has_evaluated_builder
            else draft_scenario_builder["mitigation_assignments"]
        )
        scenario_metrics, scenario_df, scenario_assumptions = build_consolidation_scenario(
            base_analytics, tuple(sorted(display_selected_suppliers)), tuple(sorted(display_mitigation_assignments))
        )
        current_scorecard = score_supplier_scenario(
            base_analytics,
            tuple(sorted(display_selected_suppliers)),
            tuple(sorted(display_mitigation_assignments)),
            scenario_metrics,
            scenario_df,
            optimization_objective=optimization_objective,
        )
        if evaluate_submitted:
            evaluated_scenario_builder = draft_scenario_builder
            for key in list(st.session_state.keys()):
                if (key.startswith("mitigation_") or key.startswith("uncovered_")) and key not in current_form_widget_values:
                    del st.session_state[key]
            for key, values in current_form_widget_values.items():
                st.session_state[key] = list(values)
            st.session_state["scenario_builder"] = evaluated_scenario_builder
            save_persisted_scenario_state(
                {
                    "data_source_label": data_source_label,
                    "scenario_builder": evaluated_scenario_builder,
                    "applied_scenario": None,
                }
            )
            st.rerun()
        with st.container(border=True):
            st.markdown("**Step 3. Finalize The Scenario**")
            st.caption("After evaluation, either reset the exercise or apply the evaluated scenario to the dashboard.")
            action_left, action_right = st.columns(2)
            with action_left:
                if st.button("Revert To Base Scenario", type="secondary", use_container_width=True):
                    st.session_state["pending_revert_to_base_scenario"] = True
                    st.rerun()
            with action_right:
                if st.button("Apply Scenario To Dashboard", type="primary", disabled=not has_evaluated_builder):
                    st.session_state["applied_scenario"] = {
                        "data_source_label": data_source_label,
                        "selected_suppliers": tuple(sorted(display_selected_suppliers)),
                        "mitigation_assignments": tuple(sorted(display_mitigation_assignments)),
                    }
                    save_persisted_scenario_state(
                        {
                            "data_source_label": data_source_label,
                            "scenario_builder": st.session_state.get("scenario_builder", {}),
                            "applied_scenario": None,
                        }
                    )
                    st.rerun()

        base_component_summary = base_analytics["component_summary"].copy()
        base_structural_single_source_names = base_component_summary.loc[
            base_component_summary["single_source_flag"], "component"
        ].tolist()
        base_uncovered_component_names: List[str] = []
        base_high_risk_count = int(base_component_summary["high_risk_flag"].sum()) if not base_component_summary.empty else 0
        empty_or_unevaluated_scenario = scenario_df.empty
        if empty_or_unevaluated_scenario:
            scenario_high_risk_count = base_high_risk_count
            scenario_structural_single_source_names = list(base_structural_single_source_names)
            resolved_single_source_names = []
            remaining_single_source_names = list(base_structural_single_source_names)
            new_single_source_names = []
            uncovered_component_names = []
            mitigated_uncovered_names = []
        else:
            scenario_high_risk_count = int(
                scenario_df["Scenario Risk Level"].isin(["High", "Uncovered"]).sum()
            ) if not scenario_df.empty else 0
            resolved_single_source_names = scenario_df.loc[
                scenario_df["Component"].isin(base_structural_single_source_names)
                & scenario_df["Scenario Status"].ne("Not Covered")
                & scenario_df["Effective Supplier Count"].gt(1),
                "Component",
            ].tolist() if not scenario_df.empty else []
            remaining_single_source_names = scenario_df.loc[
                scenario_df["Component"].isin(base_structural_single_source_names)
                & scenario_df["Scenario Status"].ne("Not Covered")
                & scenario_df["Effective Supplier Count"].eq(1),
                "Component",
            ].tolist() if not scenario_df.empty else []
            scenario_structural_single_source_names = remaining_single_source_names
            new_single_source_names = scenario_df.loc[
                (scenario_df["Effective Supplier Count"].eq(1)) & (~scenario_df["Component"].isin(base_structural_single_source_names)),
                "Component",
            ].tolist() if not scenario_df.empty else []
            uncovered_component_names = scenario_df.loc[
                scenario_df["Scenario Status"].eq("Not Covered"), "Component"
            ].tolist() if not scenario_df.empty else []
            mitigated_uncovered_names = scenario_df.loc[
                scenario_df["Scenario Status"].eq("Mitigated Coverage"), "Component"
            ].tolist() if not scenario_df.empty else []
        fixed_uncovered_count = max(0, len(base_uncovered_component_names) - len(uncovered_component_names))
        new_uncovered_count = max(0, len(uncovered_component_names) - len(base_uncovered_component_names))
        if empty_or_unevaluated_scenario:
            st.info("No suppliers are currently selected for evaluation, so the scenario comparison below shows the base issues as still unresolved.")

        c1, c2, c3, c4 = st.columns(4)
        c1.metric("Existing Single-Source Fixed", f"{len(resolved_single_source_names)}")
        c2.metric("Existing Single-Source Still Present", f"{len(remaining_single_source_names)}")
        c3.metric("New Single-Source Introduced", f"{len(new_single_source_names)}")
        c4.metric("High-Risk Change", f"{scenario_high_risk_count - base_high_risk_count:+d}")

        c5, c6, c7, c8 = st.columns(4)
        c5.metric("Existing Uncovered Fixed", f"{fixed_uncovered_count}")
        c6.metric("New Uncovered Introduced", f"{new_uncovered_count}")
        c7.metric("Covered Spend %", f"{scenario_metrics['covered_spend_share']:.0%}")
        c8.metric("Gross Savings", f"${scenario_metrics['estimated_savings']:,.0f}")

        c9, c10, c11, c12 = st.columns(4)
        c9.metric("Selected Suppliers", f"{scenario_metrics['selected_supplier_count']}")
        c10.metric("Mitigation Suppliers", f"{scenario_metrics['mitigation_supplier_count']}")
        c11.metric("Mitigated Single-Source", f"{scenario_metrics['mitigated_single_source_components']}")
        c12.metric("Mitigated Uncovered", f"{scenario_metrics['mitigated_uncovered_components']}")

        c13, c14, c15, c16 = st.columns(4)
        c13.metric("Mitigation Cost", f"${scenario_metrics['mitigation_cost']:,.0f}")
        c14.metric("Net Savings", f"${scenario_metrics['net_savings']:,.0f}")
        c15.metric("Risk Reduction", f"{scenario_metrics['aggregate_risk_reduction']:.1f}")
        c16.metric(
            "Avg Reduction / Mitigated",
            (
                f"{scenario_metrics['aggregate_risk_reduction'] / max(scenario_metrics['mitigated_single_source_components'], 1):.1f}"
                if scenario_metrics["mitigated_single_source_components"] > 0
                else "0.0"
            ),
        )
        st.caption(
            f"Current scenario score under `{optimization_objective}` is {current_scorecard.get('score', 0.0):,.1f}."
        )
        with st.expander("Current Scenario Score Breakdown"):
            breakdown = current_scorecard.get("breakdown", pd.DataFrame())
            if isinstance(breakdown, pd.DataFrame) and not breakdown.empty:
                show_table(breakdown)

        current_snapshot = build_scenario_compare_snapshot(
            "Current scenario",
            list(display_selected_suppliers),
            list(display_mitigation_assignments),
            scenario_metrics,
            current_scorecard,
        )
        comparison_store = st.session_state.get("scenario_comparisons", {})
        saved_snapshots = (
            list(comparison_store.get(data_source_label, []))
            if isinstance(comparison_store, dict)
            else []
        )
        save_compare_col, clear_compare_col = st.columns(2)
        with save_compare_col:
            if st.button("Save Current Scenario For Comparison", type="secondary", use_container_width=True):
                matching_index = next(
                    (
                        idx for idx, snapshot in enumerate(saved_snapshots)
                        if snapshot.get("signature") == current_snapshot.get("signature")
                    ),
                    None,
                )
                if matching_index is None:
                    scenario_number = len(saved_snapshots) + 1
                    current_snapshot["label"] = f"Scenario {scenario_number}"
                    saved_snapshots.append(current_snapshot)
                else:
                    current_snapshot["label"] = str(saved_snapshots[matching_index].get("label", f"Scenario {matching_index + 1}"))
                    saved_snapshots[matching_index] = current_snapshot
                if len(saved_snapshots) > 4:
                    saved_snapshots = saved_snapshots[-4:]
                st.session_state["scenario_comparisons"] = {
                    **comparison_store,
                    data_source_label: saved_snapshots,
                }
                st.rerun()
        with clear_compare_col:
            if st.button("Clear Saved Comparisons", use_container_width=True):
                st.session_state["scenario_comparisons"] = {
                    **comparison_store,
                    data_source_label: [],
                }
                st.rerun()

        if saved_snapshots:
            st.subheader("Scenario Comparison")
            st.caption("Compare the current evaluated scenario against up to four saved scenarios for this dataset.")
            comparison_table = build_scenario_compare_table(current_snapshot, saved_snapshots)
            show_table(comparison_table)
            comparison_columns = st.columns(min(len(saved_snapshots) + 1, 3))
            scenario_cards = [current_snapshot] + saved_snapshots[:2]
            for idx, snapshot in enumerate(scenario_cards):
                with comparison_columns[idx]:
                    st.markdown(f"### {snapshot['label']}")
                    st.metric("Scenario score", f"{float(snapshot['score']):,.1f}")
                    st.metric("Covered spend %", f"{float(snapshot['covered_spend_share']):.0%}")
                    st.metric("Net savings", f"${float(snapshot['net_savings']):,.0f}")
                    st.metric("Risk reduction", f"{float(snapshot['aggregate_risk_reduction']):.1f}")
                    st.caption(
                        "Suppliers: "
                        + (", ".join(snapshot.get("selected_suppliers", [])) if snapshot.get("selected_suppliers") else "none")
                    )
            if len(saved_snapshots) > 2:
                overflow_labels = ", ".join(str(snapshot.get("label", "Scenario")) for snapshot in saved_snapshots[2:])
                st.caption(f"Additional saved comparisons: {overflow_labels}. They are included in the comparison table above.")

        st.caption(
            "Existing single-source issues fixed: "
            + (", ".join(resolved_single_source_names) if resolved_single_source_names else "none")
        )
        st.caption(
            "Existing single-source issues still present: "
            + (", ".join(remaining_single_source_names) if remaining_single_source_names else "none")
        )
        st.caption(
            "New single-source issues introduced: "
            + (", ".join(new_single_source_names) if new_single_source_names else "none")
        )
        st.caption(
            "Mitigated uncovered components: "
            + (", ".join(mitigated_uncovered_names) if mitigated_uncovered_names else "none")
        )
        st.caption(
            "Uncovered components remaining: "
            + (", ".join(uncovered_component_names) if uncovered_component_names else "none")
        )
        st.caption(
            f"Modeled mitigation cost: ${scenario_metrics['mitigation_cost']:,.0f}. "
            f"Net savings after mitigation: ${scenario_metrics['net_savings']:,.0f}. "
            f"Aggregate scenario risk reduction from mitigation: {scenario_metrics['aggregate_risk_reduction']:.1f} points."
        )

        if not scenario_df.empty:
            st.caption("This visual shows how the evaluated scenario changes component coverage and risk at the same time. Components colored high or uncovered remain the priority constraints, while changes in spend coverage and supplier count show where the scenario improves or worsens resilience.")
            scenario_chart = (
                alt.Chart(scenario_df)
                .mark_bar()
                .encode(
                    x=alt.X("Component:N", axis=alt.Axis(labelAngle=-45, labelOverlap=False, labelLimit=220)),
                    y=alt.Y("Current Spend:Q", title="Current Spend"),
                    color=alt.Color(
                        "Scenario Risk Level:N",
                        scale=alt.Scale(
                            domain=["High", "Medium", "Low", "Uncovered"],
                            range=["#d73027", "#f39c12", "#2e8b57", "#7f8c8d"],
                        ),
                    ),
                    tooltip=[
                        "Component",
                        "Selected Suppliers",
                        "Single-Source Mitigation Suppliers",
                        "Selected Supplier Count",
                        "Effective Supplier Count",
                        "Current Coverage Share",
                        "Scenario Risk Level",
                    ],
                )
                .properties(height=420)
                .configure_axis(labelFontSize=11, titleFontSize=13)
                .configure_legend(labelFontSize=12, titleFontSize=13)
            )
            st.altair_chart(scenario_chart, width="stretch")
            show_table(scenario_df)
        display_assumptions("Scenario assumptions", scenario_assumptions)

    with tabs[8]:
        st.subheader("Supplier Consolidation Plan")
        if not scenario_applied:
            st.caption("Base view: treat these as scenario-testing prompts, not approved supplier moves.")
        st.write(
            (
                "This visual shows how supplier spend lines up with the current scenario-test logic, so teams can see which suppliers may be worth testing for consolidation, closer review, or exit exposure. "
                + (
                    f"The current portfolio suggests {int(supplier_summary['decision'].eq('Keep / Consolidate To').sum())} suppliers to test for consolidation potential, "
                    f"{int(supplier_summary['decision'].eq('Keep and Monitor').sum())} suppliers to test for monitoring need, and "
                    f"{int(supplier_summary['decision'].eq('Eliminate / De-prioritize').sum())} suppliers to test for exit exposure."
                )
                if not scenario_applied
                else
                "This visual shows how supplier spend lines up with the current consolidation logic, so teams can see which suppliers are stronger expansion candidates, which are watchlist suppliers, and which may be exit candidates. "
                + (
                    f"The current portfolio suggests {int(supplier_summary['decision'].eq('Keep / Consolidate To').sum())} consolidation candidates, "
                    f"{int(supplier_summary['decision'].eq('Keep and Monitor').sum())} suppliers to monitor, and "
                    f"{int(supplier_summary['decision'].eq('Eliminate / De-prioritize').sum())} potential exit suppliers."
                )
            )
        )
        st.altair_chart(build_strategic_outcomes_chart(supplier_summary, scenario_applied=scenario_applied), width="stretch")
        show_table(consolidation_plan)
        display_assumptions("Assumptions", consolidation_assumptions)

        st.subheader("Supplier Risk Assessment")
        top_supplier_risk = supplier_summary.sort_values("supplier_risk_score", ascending=False).iloc[0] if not supplier_summary.empty else None
        high_risk_supplier_names = (
            supplier_risk_assessment.loc[supplier_risk_assessment["Risk Tier"].eq("High"), "Supplier"].tolist()
            if not supplier_risk_assessment.empty and "Risk Tier" in supplier_risk_assessment.columns
            else []
        )
        st.write(
            "This visual shows which suppliers combine operational weakness and exposure risk, so teams know where supplier-management effort is most urgent. "
            + (
                f"{top_supplier_risk['supplier']} has the highest supplier risk score at {top_supplier_risk['supplier_risk_score']:.1f}, "
                f"and {'the suppliers currently flagged high risk are ' + ', '.join(high_risk_supplier_names) if high_risk_supplier_names else 'no suppliers are currently flagged high risk'}."
                if top_supplier_risk is not None
                else "Supplier risk is being ranked to highlight where resilience and performance attention should go first."
            )
            + " "
            + build_supplier_risk_methodology_note()
        )
        st.altair_chart(build_supplier_metric_chart(supplier_summary, "supplier_risk_score", "Supplier Risk Score"), width="stretch")
        show_table(supplier_risk_assessment)
        display_assumptions("Assumptions", supplier_risk_assumptions)

        st.subheader("Strategic Sourcing Plan")
        strategic_components = component_summary.loc[component_summary["kraljic_quadrant"].eq("Strategic"), "component"].tolist()
        top_strategic_component = component_summary.sort_values("strategic_priority_score", ascending=False).iloc[0] if not component_summary.empty else None
        st.write(
            "This visual helps show which components need executive attention because they combine supply risk with business impact. "
            + (
                f"The Strategic quadrant currently includes {', '.join(strategic_components[:5]) if strategic_components else 'no components'}, "
                f"and {top_strategic_component['component']}, which has the highest strategic priority score, is the top strategic-priority component."
                if top_strategic_component is not None
                else "Strategic sourcing priorities are being ranked so resilience and commercial effort can be focused where they matter most."
            )
        )
        st.altair_chart(build_kraljic_chart(component_summary), width="stretch")
        show_table(strategic_sourcing_plan)
        display_assumptions("Assumptions", strategic_sourcing_assumptions)

        st.subheader("Step-by-Step Action Plan")
        st.write(
            "This sequenced plan translates the supplier and component findings into a practical order of operations so the analyst can move from diagnosis to action without skipping exposure checks."
        )
        show_table(step_plan)

    with tabs[5]:
        st.subheader("Supplier Action Plans")
        if scenario_state and applied_scenario_metrics is not None:
            st.caption("Base vs. applied scenario supplier action plans")
            applied_supplier_order = supplier_action_plan["Supplier"].tolist()
            base_action_lookup = base_supplier_action_plan.set_index("Supplier")
            aligned_base_supplier_action_plan = (
                base_action_lookup.reindex(applied_supplier_order).reset_index()
                if not base_supplier_action_plan.empty
                else base_supplier_action_plan
            )
            base_col, applied_col = st.columns(2)
            with base_col:
                st.markdown("### Base")
                show_table(aligned_base_supplier_action_plan)
            with applied_col:
                st.markdown("### Applied Scenario")
                show_table(supplier_action_plan)
        else:
            show_table(supplier_action_plan)

    with tabs[6]:
        st.subheader("Overview & Objectives")
        st.write("Use a single decision pipeline to align spend prioritization, supply risk, supplier actions, and executive messaging.")
        bundle = make_download_bundle(
            {
                "normalized_data": normalized_df,
                "input_field_status": input_field_status,
                "supplier_summary": supplier_summary,
                "component_summary": component_summary,
                "component_supplier_detail": component_supplier_detail,
                "executive_actions": executive_actions,
                "supplier_consolidation_plan": consolidation_plan,
                "supplier_risk_assessment": supplier_risk_assessment,
                "strategic_sourcing_plan": strategic_sourcing_plan,
                "supplier_action_plan": supplier_action_plan,
                "step_plan": step_plan,
            }
        )
        powerpoint_bytes = None
        powerpoint_error = None
        try:
            powerpoint_bytes = make_powerpoint_export(summary_text, analytics, scenario_applied=scenario_applied)
        except Exception as exc:
            powerpoint_error = str(exc)
        export_col1, export_col2 = st.columns(2)
        with export_col1:
            st.download_button(
                "Download Analysis Bundle (CSV text pack)",
                data=bundle,
                file_name="supplier_analysis_bundle.csv",
                mime="text/csv",
                use_container_width=True,
            )
        with export_col2:
            if powerpoint_bytes is not None:
                st.download_button(
                    "Download Executive PowerPoint",
                    data=powerpoint_bytes,
                    file_name="supplier_analysis_executive_pack.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                )
            else:
                st.warning(powerpoint_error or "PowerPoint export is currently unavailable in this environment.")
        st.subheader("Key Insights")
        top_risk = component_summary.sort_values("supply_risk_score", ascending=False).head(3)["component"].tolist()
        top_priority = component_summary.sort_values("strategic_priority_score", ascending=False).head(3)["component"].tolist()
        retained = supplier_summary.loc[supplier_summary["decision"] == "Keep / Consolidate To", "supplier"].tolist()
        st.write(
            f"Highest-risk components: {', '.join(top_risk)}. Highest strategic priorities: {', '.join(top_priority)}. "
            f"Retain/consolidate suppliers: {', '.join(retained) if retained else 'none'}."
        )
        for visual in visual_pack:
            st.markdown(f"### {visual['title']}")
            if visual.get("summary"):
                st.write(visual["summary"])
            if visual["title"] == "Spend by Supplier":
                st.altair_chart(build_supplier_metric_chart(supplier_summary, "spend", "Spend"), width="stretch")
            elif visual["title"] == "Spend by Component":
                st.altair_chart(build_component_risk_bar_chart(component_summary, "spend", "Spend", top_n=None), width="stretch")
            elif visual["title"] == "Component Analysis Bubble":
                st.altair_chart(build_component_analysis_bubble_chart(component_summary), width="stretch")
            elif visual["title"] == "Component-Supplier Detail":
                st.altair_chart(build_component_supplier_detail_chart(component_supplier_detail, top_n=None), width="stretch")
            elif visual["title"] == "Spend Pareto (ABC)":
                st.altair_chart(
                    build_pareto_chart(spend_pareto, "component", "spend", "spend_cum_share", "spend_abc", "Component", "Spend"),
                    width="stretch",
                )
            elif visual["title"] == "Supplier Concentration by Component":
                st.altair_chart(build_supplier_concentration_chart(component_summary, top_n=None), width="stretch")
            elif visual["title"] == "Supplier Risk Score":
                st.altair_chart(build_supplier_metric_chart(supplier_summary, "supplier_risk_score", "Supplier Risk Score"), width="stretch")
            elif visual["title"] == "Risk-Adjusted Pareto":
                st.altair_chart(
                    build_pareto_chart(
                        risk_pareto,
                        "component",
                        "risk_adjusted_spend",
                        "risk_cum_share",
                        "risk_abc",
                        "Component",
                        "Risk-Adjusted Spend",
                        color_col="sourcing_risk_level",
                        color_title="Sourcing Risk",
                        color_domain=RISK_LEVEL_DOMAIN,
                        color_range=RISK_LEVEL_RANGE,
                    ),
                    width="stretch",
                )
            elif visual["title"] == "Strategic Priority Pareto":
                st.altair_chart(
                    build_pareto_chart(
                        strategic_pareto,
                        "component",
                        "strategic_priority_score",
                        "strategic_cum_share",
                        "strategic_abc",
                        "Component",
                        "Strategic Priority Score",
                    ),
                    width="stretch",
                )
            elif visual["title"] == "Kraljic Positioning":
                st.altair_chart(build_kraljic_chart(component_summary), width="stretch")
            elif visual["title"] == "Supply Risk Score":
                st.altair_chart(build_component_risk_bar_chart(component_summary, "supply_risk_score", "Supply Risk Score", top_n=None), width="stretch")
            elif visual["title"] == "Supplier Count by Component":
                st.altair_chart(
                    build_component_risk_bar_chart(component_summary, "supplier_count", "Supplier Count", top_n=None, ascending=True),
                    width="stretch",
                )
            elif visual["title"] == "Strategic Sourcing Outcomes":
                st.altair_chart(build_strategic_outcomes_chart(supplier_summary, scenario_applied=scenario_applied), width="stretch")
            elif visual["title"] == "Supplier Spend by Component Mix":
                st.altair_chart(build_supplier_component_mix_chart(component_supplier_detail, component_summary, top_n_suppliers=None), width="stretch")
            elif visual["title"] == "Supplier Spend by Kraljic Mix":
                st.altair_chart(build_supplier_quadrant_mix_chart(component_supplier_detail, component_summary, top_n_suppliers=None), width="stretch")
            else:
                st.altair_chart(build_component_risk_bar_chart(component_summary, "supply_risk_score", "Supply Risk Score"), width="stretch")
            st.caption("Supporting Data")
            show_table(visual["data"].reset_index())
            professor_notes = build_professor_notes(str(visual["title"]), analytics, scenario_applied=scenario_applied)
            activity_note = build_professor_activity_note(str(visual["title"]), analytics, scenario_applied=scenario_applied)
            takeaway_points = (professor_notes or list(visual["talking_points"])) + ([activity_note] if activity_note else [])
            st.caption("Key Takeaways")
            for point in takeaway_points:
                st.write(f"- {point}")

    with st.expander("Normalized Input Preview"):
        show_table(build_normalized_input_preview(normalized_df, input_diagnostics))


if __name__ == "__main__":
    render_app()
